from fastapi import FastAPI, Query, HTTPException, status, Depends, Form
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Optional, Dict
from datetime import datetime
import uvicorn
import os
from pydantic import BaseModel
from sqlalchemy.orm import Session
from sqlalchemy import func
from database import get_db, engine
import sql_models as models
import models as schemas
import models as schemas
from services.openai_service import generate_embedding
from mock_db import (
    PATHWAYS_DB, QUIZZES_DB, QUIZ_SUBMISSIONS_DB, DISCUSSIONS_DB,
    DISCUSSION_COMMENTS_DB, NOTIFICATIONS_DB, CONVERSATIONS_DB,
    CREDENTIALS_DB
)
from services.seeding_service import init_sample_data

# Create tables
models.Base.metadata.create_all(bind=engine)


class GradingQueueItem(BaseModel):
    id: int
    project_id: int
    student_name: str
    project_title: str
    submission_type: str
    course_title: str
    submitted_at: str
    status: str


app = FastAPI(
    title="TSEA-X API",
    description="Backend for TSEA-X Platform with CDIO Framework and NUSA IRIS Cycle",
    version="2.1.0"
)

# DEBUG: Print DB Config
print(f"DEBUG: CWD = {os.getcwd()}")
try:
    from database import DATABASE_URL
    print(f"DEBUG: DATABASE_URL = {DATABASE_URL}")
except ImportError:
    print("DEBUG: Could not import DATABASE_URL")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:3001",
        "https://t6.tsea.asia",
        "http://t6.tsea.asia",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

from routers import naska, ai, courses
from routers.debug import router as debug_router
app.include_router(naska.router)
app.include_router(ai.router)
app.include_router(courses.router)
app.include_router(debug_router)

@app.get("/")
def read_root():
    return {"message": "TSEA-X Backend API is running"}

@app.get("/api/health")
def health_check():
    return {"status": "ok"}

# --- Mock Data ---
# COURSES_DB removed - using database
# Imported from mock_db.py

# ID counters
project_counter = 1
charter_counter = 1
blueprint_counter = 1
implementation_counter = 1
deployment_counter = 1
quiz_counter = 1
quiz_submission_counter = 1
discussion_counter = 1
comment_counter = 1
notification_counter = 1
conversation_counter = 1
message_counter = 1
enrollment_counter = 1
credential_counter = 1
snapshot_counter = 1


# Initialize on startup
@app.on_event("startup")
def startup_event():
    # Initialize sample data on startup
<<<<<<< HEAD
    db = next(get_db())
    try:
        init_sample_data(db)
    finally:
        db.close()
=======
    try:
        db = next(get_db())
        try:
            init_sample_data(db)
        finally:
            db.close()
    except Exception as e:
        print(f"Startup Warning: Could not initialize sample data: {e}")
>>>>>>> 36b5f396c275188512a1d135e85e4f68b59f3fef


# ===== AUTH ENDPOINTS =====

class LoginRequest(BaseModel):
    email: str
    password: str

class SyncUserRequest(BaseModel):
    supabase_id: str
    email: str
    full_name: str
    role: str = "student"


@app.post("/api/v1/auth/login")
def demo_login(request: LoginRequest, db: Session = Depends(get_db)):
    """Demo login for testing - accepts 'test' as password for all seeded users"""
    user = db.query(models.User).filter(models.User.email == request.email).first()
    
    if not user:
        raise HTTPException(status_code=401, detail="Invalid login credentials")
    
    # For demo purposes, accept 'test' as password for all seeded users
    if request.password != "test":
        raise HTTPException(status_code=401, detail="Invalid login credentials")
    
    return {
        "id": user.id,
        "email": user.email,
        "name": user.full_name,
        "role": user.role,
        "institution_id": user.institution_id,
        "instructor_type": user.instructor_type
    }


@app.post("/api/v1/auth/sync-user")
def sync_user(request: SyncUserRequest, db: Session = Depends(get_db)):
    """Sync Supabase user to local database"""
    # Check if user exists by supabase_id
    user = db.query(models.User).filter(models.User.supabase_id == request.supabase_id).first()
    
    if not user:
        # Check by email
        user = db.query(models.User).filter(models.User.email == request.email).first()
        if user:
            # Link existing user to supabase
            user.supabase_id = request.supabase_id
        else:
            # Create new user
            user = models.User(
                email=request.email,
                full_name=request.full_name,
                role=request.role,
                supabase_id=request.supabase_id
            )
            db.add(user)
    
    db.commit()
    db.refresh(user)
    
    return {
        "id": user.id,
        "email": user.email,
        "full_name": user.full_name,
        "role": user.role,
        "institution_id": user.institution_id,
        "instructor_type": user.instructor_type
    }

@app.get("/api/v1/courses", response_model=List[schemas.Course])
def get_courses(category: Optional[str] = None, db: Session = Depends(get_db)):
    query = db.query(models.Course)
    if category:
        query = query.filter(models.Course.category == category)
    return query.all()

@app.get("/api/v1/pathways", response_model=List[schemas.Pathway])
def get_pathways():
    # Pathways are still static for now as they are UI configuration
    return PATHWAYS_DB

# --- Institution Routes ---

@app.get("/api/v1/institutions", response_model=List[schemas.Institution])
def get_institutions(
    type: Optional[str] = None,
    country: Optional[str] = None,
    featured_only: bool = False,
    db: Session = Depends(get_db)
):
    """Get all institutions with optional filters"""
    query = db.query(models.Institution)
    
    if type:
        query = query.filter(models.Institution.type == type)
    
    if country:
        query = query.filter(models.Institution.country.ilike(country))
    
    if featured_only:
        query = query.filter(models.Institution.is_featured == True)
    
    return query.order_by(models.Institution.name).all()

@app.get("/api/v1/institutions/{institution_id}", response_model=schemas.InstitutionDetail)
def get_institution_detail(institution_id: int, db: Session = Depends(get_db)):
    """Get detailed institution profile with courses and programs"""
    institution = db.query(models.Institution).filter(models.Institution.id == institution_id).first()
    if not institution:
        raise HTTPException(status_code=404, detail="Institution not found")
    
    # Get courses for this institution
    institution_courses = db.query(models.Course).filter(models.Course.institution_id == institution_id).all()
    
    # Get pathways (mock for now)
    institution_pathways = []
    
    # Mock featured instructors
    featured_instructors = []
    
    return schemas.InstitutionDetail(
        institution=institution,
        courses=institution_courses,
        featured_instructors=featured_instructors,
        pathways=institution_pathways
    )

@app.get("/api/v1/institutions/{institution_id}/courses", response_model=List[schemas.Course])
def get_institution_courses(institution_id: int, db: Session = Depends(get_db)):
    """Get all courses from a specific institution"""
    institution = db.query(models.Institution).filter(models.Institution.id == institution_id).first()
    if not institution:
        raise HTTPException(status_code=404, detail="Institution not found")
    
    courses = db.query(models.Course).filter(models.Course.institution_id == institution_id).all()
    return courses

@app.get("/api/v1/institutions/{institution_id}/stats", response_model=schemas.InstitutionStats)
def get_institution_stats(institution_id: int, db: Session = Depends(get_db)):
    """Get dashboard statistics for institution (admin only)"""
    institution = db.query(models.Institution).filter(models.Institution.id == institution_id).first()
    if not institution:
        raise HTTPException(status_code=404, detail="Institution not found")
    
    # Calculate stats
    # 1. Total Active Students (unique students enrolled in institution's courses)
    total_active_students = db.query(models.Enrollment.user_id).join(models.Course).filter(
        models.Course.institution_id == institution_id,
        models.Enrollment.status == 'active'
    ).distinct().count()
    
    # 2. Total Completions (credentials issued for institution's courses)
    total_completions = db.query(models.Credential).join(models.Course).filter(
        models.Course.institution_id == institution_id
    ).count()
    
    # 3. Total Certificates (same as completions for now)
    total_certificates_issued = total_completions
    
    # 4. Average Course Rating
    avg_rating = db.query(func.avg(models.Course.rating)).filter(
        models.Course.institution_id == institution_id
    ).scalar() or 0.0
    
    # 5. Total Revenue
    total_revenue = db.query(func.sum(models.RevenueTransaction.amount)).join(models.Course).filter(
        models.Course.institution_id == institution_id
    ).scalar() or 0.0
    
    return schemas.InstitutionStats(
        institution_id=institution_id,
        total_active_students=total_active_students,
        total_completions=total_completions,
        total_certificates_issued=total_certificates_issued,
        average_course_rating=round(avg_rating, 1),
        total_revenue=total_revenue,
        month_over_month_growth=12.5 # Placeholder for now
    )


@app.get("/api/v1/institutions/{institution_id}/dashboard")
def get_institution_dashboard(institution_id: int, db: Session = Depends(get_db)):
    """Get comprehensive dashboard data for institution admin (stats, courses, enrollments)"""
    institution = db.query(models.Institution).filter(models.Institution.id == institution_id).first()
    if not institution:
        raise HTTPException(status_code=404, detail="Institution not found")
    
    # Get all courses for this institution
    courses = db.query(models.Course).filter(models.Course.institution_id == institution_id).all()
    
    # Calculate stats
    total_courses = len(courses)
    
    # Get unique enrolled students
    total_enrollments = db.query(models.Enrollment).join(models.Course).filter(
        models.Course.institution_id == institution_id
    ).count()
    
    active_students = db.query(models.Enrollment.user_id).join(models.Course).filter(
        models.Course.institution_id == institution_id,
        models.Enrollment.status == 'active'
    ).distinct().count()
    
    completed_enrollments = db.query(models.Enrollment).join(models.Course).filter(
        models.Course.institution_id == institution_id,
        models.Enrollment.status == 'completed'
    ).count()
    
    completion_rate = round((completed_enrollments / total_enrollments * 100) if total_enrollments > 0 else 0, 1)
    
    # Average rating
    avg_rating = db.query(func.avg(models.Course.rating)).filter(
        models.Course.institution_id == institution_id
    ).scalar() or 0.0
    
    # Recent enrollments (last 7 days)
    from datetime import timedelta
    week_ago = datetime.utcnow() - timedelta(days=7)
    new_enrollments_this_week = db.query(models.Enrollment).join(models.Course).filter(
        models.Course.institution_id == institution_id,
        models.Enrollment.enrolled_at >= week_ago
    ).count()
    
    # Build courses list with enrollment counts and completion rates
    courses_data = []
    for course in courses:
        course_enrollments = db.query(models.Enrollment).filter(
            models.Enrollment.course_id == course.id
        ).count()
        
        course_completed = db.query(models.Enrollment).filter(
            models.Enrollment.course_id == course.id,
            models.Enrollment.status == 'completed'
        ).count()
        
        course_completion_rate = round((course_completed / course_enrollments * 100) if course_enrollments > 0 else 0, 1)
        
        # Determine status based on approval_status
        status = 'published'
        if hasattr(course, 'approval_status'):
            if course.approval_status == 'pending_approval':
                status = 'pending'
            elif course.approval_status == 'rejected':
                status = 'draft'
        
        courses_data.append({
            "id": course.id,
            "title": course.title,
            "instructor": course.instructor,
            "enrollments": course_enrollments,
            "rating": course.rating or 0.0,
            "status": status,
            "completionRate": course_completion_rate,
            "thumbnail": course.image or f"https://images.unsplash.com/photo-1516321318423-f06f85e504b3?w=400",
            "category": course.category or "Technology"
        })
    
    # Get recent enrollments with student info
    recent_enrollments_query = db.query(
        models.Enrollment, models.User, models.Course
    ).join(
        models.User, models.Enrollment.user_id == models.User.id
    ).join(
        models.Course, models.Enrollment.course_id == models.Course.id
    ).filter(
        models.Course.institution_id == institution_id
    ).order_by(
        models.Enrollment.enrolled_at.desc()
    ).limit(10).all()
    
    recent_enrollments = []
    for enrollment, user, course in recent_enrollments_query:
        # Calculate progress (mock for now - could be based on project phases)
        progress = 0
        if enrollment.status == 'completed':
            progress = 100
        elif enrollment.status == 'active':
            # Check if there's a project and calculate based on phase
            project = db.query(models.CDIOProject).filter(
                models.CDIOProject.course_id == course.id,
                models.CDIOProject.user_id == user.id
            ).first()
            if project:
                progress = project.completion_percentage or 0
        
        recent_enrollments.append({
            "id": enrollment.id,
            "student_name": user.full_name or user.email.split('@')[0],
            "student_email": user.email,
            "course_title": course.title,
            "enrolled_date": enrollment.enrolled_at.isoformat() if enrollment.enrolled_at else None,
            "progress": progress,
            "status": enrollment.status
        })
    
    # Format learners count
    learners_display = f"{active_students:,}" if active_students < 10000 else f"{active_students // 1000}K+"
    
    return {
        "institution": {
            "id": institution.id,
            "name": institution.name,
            "short_name": institution.short_name or institution.name[:3].upper(),
            "logo_url": institution.logo_url or f"https://ui-avatars.com/api/?name={institution.short_name or institution.name}&background=0066cc&color=fff&size=128",
            "type": institution.type
        },
        "stats": {
            "total_courses": total_courses,
            "total_learners": learners_display,
            "total_programs": 0,  # Programs feature not implemented yet
            "total_enrollments": total_enrollments,
            "completion_rate": completion_rate,
            "average_rating": round(avg_rating, 1),
            "revenue_this_month": 0,  # Revenue tracking not implemented
            "new_enrollments_this_week": new_enrollments_this_week
        },
        "courses": courses_data,
        "recent_enrollments": recent_enrollments
    }


# ===== PARTNER INSTRUCTOR MANAGEMENT =====

@app.get("/api/v1/institutions/{institution_id}/instructors")
def get_institution_instructors(institution_id: int, db: Session = Depends(get_db)):
    """Get all instructors belonging to an institution"""
    institution = db.query(models.Institution).filter(models.Institution.id == institution_id).first()
    if not institution:
        raise HTTPException(status_code=404, detail="Institution not found")
    
    instructors = db.query(models.User).filter(
        models.User.institution_id == institution_id,
        models.User.role == "instructor"
    ).all()
    
    result = []
    for inst in instructors:
        # Get course count for this instructor
        course_count = db.query(models.Course).filter(
            models.Course.instructor_id == inst.id
        ).count()
        # Get pending courses count
        pending_count = db.query(models.Course).filter(
            models.Course.instructor_id == inst.id,
            models.Course.approval_status == "pending_approval"
        ).count()
        
        result.append({
            "id": inst.id,
            "name": inst.full_name,
            "email": inst.email,
            "instructor_type": inst.instructor_type,
            "courses_count": course_count,
            "pending_courses": pending_count,
            "created_at": inst.created_at.isoformat() if inst.created_at else None
        })
    
    return result


@app.post("/api/v1/institutions/{institution_id}/instructors")
def add_institution_instructor(
    institution_id: int,
    email: str = Form(...),
    full_name: str = Form(...),
    db: Session = Depends(get_db)
):
    """Add a new instructor to an institution"""
    institution = db.query(models.Institution).filter(models.Institution.id == institution_id).first()
    if not institution:
        raise HTTPException(status_code=404, detail="Institution not found")
    
    # Check if user already exists
    existing = db.query(models.User).filter(models.User.email == email).first()
    if existing:
        # Update existing user to be institutional instructor
        existing.institution_id = institution_id
        existing.instructor_type = "institutional"
        existing.role = "instructor"
        db.commit()
        return {"message": f"User {email} added as instructor", "user_id": existing.id}
    
    # Create new instructor
    new_instructor = models.User(
        email=email,
        full_name=full_name,
        role="instructor",
        institution_id=institution_id,
        instructor_type="institutional"
    )
    db.add(new_instructor)
    db.commit()
    db.refresh(new_instructor)
    
    return {"message": "Instructor created", "user_id": new_instructor.id}


@app.delete("/api/v1/institutions/{institution_id}/instructors/{user_id}")
def remove_institution_instructor(institution_id: int, user_id: int, db: Session = Depends(get_db)):
    """Remove an instructor from an institution (doesn't delete the user)"""
    user = db.query(models.User).filter(
        models.User.id == user_id,
        models.User.institution_id == institution_id
    ).first()
    if not user:
        raise HTTPException(status_code=404, detail="Instructor not found in this institution")
    
    # Remove institution link, convert to individual
    user.institution_id = None
    user.instructor_type = "individual"
    db.commit()
    
    return {"message": "Instructor removed from institution"}


@app.get("/api/v1/institutions/{institution_id}/pending-courses")
def get_pending_courses(institution_id: int, db: Session = Depends(get_db)):
    """Get courses pending approval for an institution"""
    institution = db.query(models.Institution).filter(models.Institution.id == institution_id).first()
    if not institution:
        raise HTTPException(status_code=404, detail="Institution not found")
    
    pending_courses = db.query(models.Course).filter(
        models.Course.institution_id == institution_id,
        models.Course.approval_status == "pending_approval"
    ).all()
    
    result = []
    for course in pending_courses:
        instructor = db.query(models.User).filter(models.User.id == course.instructor_id).first()
        result.append({
            "id": course.id,
            "title": course.title,
            "category": course.category,
            "level": course.level,
            "instructor_id": course.instructor_id,
            "instructor_name": instructor.full_name if instructor else course.instructor,
            "created_at": None  # Add created_at to Course model if needed
        })
    
    return result


# ===== INSTRUCTOR GRADING QUEUE =====

@app.get("/api/v1/instructor/grading-queue")
def get_grading_queue(
    instructor_id: Optional[int] = Query(None),
    db: Session = Depends(get_db)
):
    """
    Get all pending submissions for an instructor's courses.
    Returns charters, blueprints, and implementations that need grading.
    """
    grading_items = []
    
    # Get all projects with submissions
    projects = db.query(models.CDIOProject).all()
    
    for project in projects:
        # Get student info
        student = db.query(models.User).filter(models.User.id == project.user_id).first()
        if not student:
            continue
            
        # Get course info
        course = db.query(models.Course).filter(models.Course.id == project.course_id).first()
        if not course:
            continue
        
        # If instructor_id is provided, filter by instructor's courses
        if instructor_id and course.instructor_id != instructor_id:
            continue
        
        # Check for charter submission
        if project.charter:
            grading_items.append({
                "id": project.charter.id,
                "project_id": project.id,
                "student_name": student.full_name,
                "project_title": project.title,
                "submission_type": "Charter",
                "course_title": course.title,
                "submitted_at": project.charter.created_at.strftime("%Y-%m-%d %H:%M") if project.charter.created_at else "N/A",
                "status": "Pending"
            })
        
        # Check for blueprint submission
        if project.blueprint:
            grading_items.append({
                "id": project.blueprint.id,
                "project_id": project.id,
                "student_name": student.full_name,
                "project_title": project.title,
                "submission_type": "Blueprint",
                "course_title": course.title,
                "submitted_at": project.blueprint.created_at.strftime("%Y-%m-%d %H:%M") if project.blueprint.created_at else "N/A",
                "status": "Pending"
            })
        
        # Check for implementation submission
        if project.implementation:
            # Check if already graded (has AI feedback)
            status = "Graded" if project.implementation.ai_feedback else "Pending"
            grading_items.append({
                "id": project.implementation.id,
                "project_id": project.id,
                "student_name": student.full_name,
                "project_title": project.title,
                "submission_type": "Implementation",
                "course_title": course.title,
                "submitted_at": project.implementation.created_at.strftime("%Y-%m-%d %H:%M") if project.implementation.created_at else "N/A",
                "status": status
            })
    
    # Sort by submitted_at (most recent first)
    grading_items.sort(key=lambda x: x["submitted_at"], reverse=True)
    
    return grading_items


@app.post("/api/v1/courses/{course_id}/approve")
def approve_course(
    course_id: int,
    approver_id: int = Form(...),
    db: Session = Depends(get_db)
):
    """Approve a pending course for publishing"""
    course = db.query(models.Course).filter(models.Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    
    if course.approval_status != "pending_approval":
        raise HTTPException(status_code=400, detail="Course is not pending approval")
    
    course.approval_status = "approved"
    course.approved_by = approver_id
    course.approved_at = datetime.utcnow()
    course.rejection_reason = None
    db.commit()
    
    return {"message": "Course approved", "course_id": course_id}


@app.post("/api/v1/courses/{course_id}/reject")
def reject_course(
    course_id: int,
    approver_id: int = Form(...),
    reason: str = Form(...),
    db: Session = Depends(get_db)
):
    """Reject a pending course with feedback"""
    course = db.query(models.Course).filter(models.Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    
    if course.approval_status != "pending_approval":
        raise HTTPException(status_code=400, detail="Course is not pending approval")
    
    course.approval_status = "rejected"
    course.approved_by = approver_id
    course.approved_at = datetime.utcnow()
    course.rejection_reason = reason
    db.commit()
    
    return {"message": "Course rejected", "course_id": course_id, "reason": reason}



# ===== Course Editor Endpoints =====

<<<<<<< HEAD

=======
class ContentBlock(BaseModel):
    id: str
    type: str
    content: str
    metadata: Optional[Dict] = {}

class CourseModuleSchema(BaseModel):
    title: str
    order: int
    content_blocks: List[ContentBlock]
    status: str

@app.get("/api/v1/courses/{course_id}")
def get_course_details(course_id: int, db: Session = Depends(get_db)):
    course = db.query(models.Course).filter(models.Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    return {
        "id": course.id,
        "title": course.title,
        "description": course.description,
        "instructor": course.instructor
    }

@app.get("/api/v1/courses/{course_id}/modules")
def get_course_modules(course_id: int, db: Session = Depends(get_db)):
    modules = db.query(models.CourseModule).filter(
        models.CourseModule.course_id == course_id
    ).order_by(models.CourseModule.order).all()
    
    return [
        {
            "id": m.id,
            "title": m.title,
            "order": m.order,
            "content_blocks": m.content_blocks,
            "status": m.status
        }
        for m in modules
    ]

@app.post("/api/v1/courses/{course_id}/modules")
def save_course_modules(
    course_id: int, 
    modules: List[CourseModuleSchema], 
    db: Session = Depends(get_db)
):
    # Verify course exists
    course = db.query(models.Course).filter(models.Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
        
    try:
        # Simplistic sync: Delete existing and recreate (or update)
        # For this demo/prototype, let's update by title or recreate.
        # Better strategy: 
        # 1. Get existing modules
        # 2. Update matches, delete removed, add new. 
        # But since we send the WHOLE list, we can just sync.
        
        # Current implementation: Update if title matches (preserving IDs), else create.
        # NOTE: This implies we can't rename modules easily without losing ID.
        # But getting correct IDs from frontend would be better.
        # The frontend sends 'CourseModule' which doesn't seem to track module ID in state explicitely yet?
        # Let's check frontend. Frontend has 'CourseModule' interface but no ID.
        # So we might need to wipe and recreate or match by order/title.
        # Let's match by title for now.
        
        for i, mod_data in enumerate(modules):
            existing = db.query(models.CourseModule).filter(
                models.CourseModule.course_id == course_id,
                models.CourseModule.title == mod_data.title
            ).first()
            
            blocks = [b.dict() for b in mod_data.content_blocks]
            
            if existing:
                existing.order = i
                existing.content_blocks = blocks
                existing.status = mod_data.status
                existing.updated_at = datetime.utcnow()
            else:
                new_mod = models.CourseModule(
                    course_id=course_id,
                    title=mod_data.title,
                    order=i,
                    content_blocks=blocks,
                    status=mod_data.status
                )
                db.add(new_mod)
        
        db.commit()
        return {"message": "Modules saved successfully"}
        
    except Exception as e:
        db.rollback()
        print(f"Error saving modules: {e}")
        raise HTTPException(status_code=500, detail=str(e))
>>>>>>> 36b5f396c275188512a1d135e85e4f68b59f3fef

@app.post("/api/v1/ai/suggest-content")
async def ai_suggest_content(request: Dict[str, str]):
    """AI Helper for Course Editor"""
    context = request.get("context", "")
    course_title = request.get("course_title", "")
    module_title = request.get("module_title", "")
    
    prompt = f"""
    You are an expert teaching assistant.
    Course: {course_title}
    Module: {module_title}
    Current Content context: "{context}"
    
    Provide a specific, constructive suggestion to improve this educational content.
    Keep it short (max 2 sentences).
    """
    
    try:
        from services.openai_service import clients
        if "gemini" in clients:
            response = await clients["gemini"].generate_content_async(prompt)
            return {"suggestion": response.text.strip()}
        return {"suggestion": "Try adding a concrete example here to illustrate the concept."}
    except Exception as e:
        return {"suggestion": "Consider adding an interactive element."}


class CourseModuleCreate(BaseModel):
    title: str
    order: int = 0
    content_blocks: List[Dict] = []
    status: str = "draft"

@app.get("/api/v1/courses/{course_id}/modules", response_model=List)
def get_course_modules(course_id: int, db: Session = Depends(get_db)):
    """Get all content modules for a course"""
    modules = db.query(models.CourseModule).filter(models.CourseModule.course_id == course_id).order_by(models.CourseModule.order).all()
    # If no modules exist, create a default one
    if not modules:
        default_module = models.CourseModule(
            course_id=course_id,
            title="Module 1: Getting Started",
            order=0,
            content_blocks=[],
            status="draft"
        )
        db.add(default_module)
        db.commit()
        db.refresh(default_module)
        return [default_module]
    return modules

@app.post("/api/v1/courses/{course_id}/modules")
def update_course_modules(course_id: int, modules: List[CourseModuleCreate], db: Session = Depends(get_db)):
    """Update all modules (full replace for simplicity in this MVP)"""
    # Verify course exists
    course = db.query(models.Course).filter(models.Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    # In a real app we'd do smart diffing, but for MVP we can delete old and re-create
    # Make sure to keep IDs if possible, but simpler to wipe for now or check by title/order
    
    # 1. Delete existing (brute force strategy for rapid prototyping)
    db.query(models.CourseModule).filter(models.CourseModule.course_id == course_id).delete()
    
    # 2. Create new
    new_modules = []
    for m in modules:
        new_mod = models.CourseModule(
            course_id=course_id,
            title=m.title,
            order=m.order,
            content_blocks=m.content_blocks,
            status=m.status
        )
        db.add(new_mod)
        new_modules.append(new_mod)
    
    db.commit()
    
    return {"message": "Course content saved", "count": len(new_modules)}


@app.get("/api/v1/instructors/{user_id}/type")
def get_instructor_type(user_id: int, db: Session = Depends(get_db)):
    """Get instructor type (institutional or individual) for course creation workflow"""
    user = db.query(models.User).filter(models.User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    return {
        "user_id": user_id,
        "instructor_type": user.instructor_type,
        "institution_id": user.institution_id,
        "requires_approval": user.instructor_type == "institutional"
    }

@app.post("/api/v1/search")
async def vector_search(query: str, db: Session = Depends(get_db)):
    """Semantic search using pgvector"""
    try:
        embedding = await generate_embedding(query)
        
        # Search knowledge nodes using L2 distance
        results = db.query(models.KnowledgeNode).order_by(
            models.KnowledgeNode.embedding.l2_distance(embedding)
        ).limit(5).all()
        
        return {
            "query": query,
            "results": results
        }
    except Exception as e:
        print(f"Vector Search Error: {e}")
        return {"query": query, "results": [], "error": str(e)}

class CreateCourseRequest(BaseModel):
    topic: str
    target_audience: str = "Beginner"

from fastapi import UploadFile, File, Form
import io

@app.post("/api/v1/courses/generate")
async def generate_course(
    topic: str = Form(...),
    target_audience: str = Form("Beginner"),
    files: List[UploadFile] = File(None)
):
    """Generate course structure using AI, optionally with uploaded materials"""
    
    material_content = ""
    if files:
        for file in files:
            try:
                content = await file.read()
                # Simple text decoding, assume utf-8
                text = content.decode("utf-8", errors="ignore")
                material_content += f"\n\n--- File: {file.filename} ---\n{text}"
            except Exception as e:
                print(f"Error reading file {file.filename}: {e}")
    
    return await openai_service.generate_course_structure(topic, target_audience, material_content)


# ===== SMART COURSE CREATION (Alexandria AI Engine) =====

class SmartParseRequest(BaseModel):
    file_base64: str
    file_name: str
    mime_type: str = "application/pdf"

class SmartAgendaRequest(BaseModel):
    parsed_content: Dict
    target_audience: str = "Intermediate"
    duration_weeks: int = 4

class SmartKnowledgeGraphRequest(BaseModel):
    content: str

class InternationalSyllabusRequest(BaseModel):
    course_title: str
    course_description: str
    duration_weeks: int = 4
    level: str = "Intermediate"
    material_content: str = ""

# Import syllabus generator
from services import syllabus_generator

@app.post("/api/v1/courses/smart-create/parse")
async def smart_parse_materials(request: SmartParseRequest):
    """
    Parse uploaded materials using AI (Gemini Vision for multi-modal).
    Extracts structured content from slides, PDFs, and images.
    
    Powered by Alexandria AI Engine for intelligent content extraction.
    """
    result = await openai_service.parse_materials_multimodal(
        file_base64=request.file_base64,
        file_name=request.file_name,
        mime_type=request.mime_type
    )
    return result


@app.post("/api/v1/courses/smart-create/agenda")
async def smart_generate_agenda(request: SmartAgendaRequest):
    """
    Generate Alexandria AI Engine teaching agenda with heterogeneous teaching actions.
    
    Teaching action types: EXPLAIN, DISCUSS, PRACTICE, QUIZ, DEMO, REFLECT, COLLABORATE
    """
    result = await openai_service.generate_teaching_agenda(
        parsed_content=request.parsed_content,
        target_audience=request.target_audience,
        duration_weeks=request.duration_weeks
    )
    return result


@app.post("/api/v1/courses/smart-create/knowledge-graph")
async def smart_extract_knowledge_graph(request: SmartKnowledgeGraphRequest):
    """
    Extract knowledge points for visual knowledge graph.
    Returns nodes and edges for graph visualization.
    """
    result = await openai_service.extract_knowledge_points(content=request.content)
    return result


@app.post("/api/v1/courses/smart-create/international-syllabus")
async def generate_international_syllabus_endpoint(request: InternationalSyllabusRequest):
    """
    Generate an international-standard syllabus aligned with:
    - MIT OpenCourseWare structure
    - CDIO Framework (Conceive-Design-Implement-Operate)
    - IRIS Framework (Immerse-Realize-Iterate-Scale)
    - Bloom's Taxonomy learning objectives
    
    Powered by Alexandria AI Engine.
    """
    # Get the Gemini client for AI generation
    from services.openai_service import clients
    
    ai_client = clients.get("gemini")
    provider = "gemini" if ai_client else "mock"
    
    result = await syllabus_generator.generate_international_syllabus(
        course_title=request.course_title,
        course_description=request.course_description,
        duration_weeks=request.duration_weeks,
        level=request.level,
        material_content=request.material_content,
        ai_client=ai_client,
        provider=provider
    )
    return result


@app.post("/api/v1/courses/smart-create/upload")
async def smart_upload_materials(
    files: List[UploadFile] = File(...)
):
    """
    Upload and parse materials for Smart Course Creation.
    Accepts PDF, PPTX, images, and text files.
    Returns base64 encoded content for further processing.
    """
    parsed_files = []
    
    for file in files:
        try:
            content = await file.read()
            import base64
            file_base64 = base64.b64encode(content).decode('utf-8')
            
            # Determine MIME type
            mime_type = file.content_type or "application/octet-stream"
            
            parsed_files.append({
                "filename": file.filename,
                "mime_type": mime_type,
                "size_bytes": len(content),
                "base64": file_base64
            })
        except Exception as e:
            print(f"Error processing file {file.filename}: {e}")
            parsed_files.append({
                "filename": file.filename,
                "error": str(e)
            })
    
    return {
        "success": True,
        "files": parsed_files,
        "total_files": len(parsed_files)
    }


class AISuggestRequest(BaseModel):
    context: str
    course_title: str = ""
    module_title: str = ""

@app.post("/api/v1/ai/suggest-content")
async def ai_suggest_content(request: AISuggestRequest):
    """
    Get AI-powered content suggestions for the course editor.
    Provides contextual recommendations based on current content.
    """
    from services.openai_service import clients, PRIORITY
    
    prompt = f"""You are Alexandria AI, an expert instructional designer. Provide a brief, actionable suggestion to improve the following course content.

CONTEXT:
Course: {request.course_title or 'Untitled Course'}
Module: {request.module_title or 'Current Module'}
Current Content: {request.context[:500]}

Provide ONE concise suggestion (1-2 sentences) to make this content more engaging, clear, or educationally effective.
Focus on practical improvements like:
- Adding examples or case studies
- Including interactive elements
- Clarifying complex concepts
- Adding visual aids or multimedia
- Incorporating assessment checkpoints

Respond with just the suggestion, no preamble."""

    for provider in PRIORITY:
        if provider == "mock":
            suggestions = [
                "Consider adding a real-world case study to illustrate this concept and help learners connect theory to practice.",
                "Add an interactive quiz checkpoint here to reinforce learning before moving to the next section.",
                "Include a video demonstration or diagram to visualize this process for visual learners.",
                "Break this section into smaller chunks with reflection prompts between each concept.",
                "Add a Socratic AI checkpoint where students can ask questions about this material."
            ]
            import random
            return {"success": True, "suggestion": random.choice(suggestions), "provider": "mock"}
        
        if provider not in clients:
            continue
            
        try:
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                return {"success": True, "suggestion": response.text.strip(), "provider": provider}
            else:
                response = await clients[provider].chat.completions.create(
                    model="gpt-4o" if provider == "openai" else "google/gemini-2.0-flash-exp:free",
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=150
                )
                return {"success": True, "suggestion": response.choices[0].message.content.strip(), "provider": provider}
        except Exception as e:
            print(f"[{provider}] AI Suggest Error: {e}")
            continue
    
    return {"success": False, "suggestion": "Unable to generate suggestion at this time.", "provider": "error"}


class AgentConsultRequest(BaseModel):
    agent_type: str  # 'alexandria', 'sme', 'pedagogy'
    query: str
    course_context: Optional[Dict] = None
    system_prompt: Optional[str] = None

@app.post("/api/v1/ai/agent-consult")
async def ai_agent_consult(request: AgentConsultRequest):
    """
    Consult with MAIC-inspired AI agents for course development.
    
    Agents:
    - alexandria: Curriculum Designer (structure, IRIS/CDIO, learning outcomes)
    - sme: Subject Matter Expert (content depth, examples, exercises)
    - pedagogy: Pedagogy Advisor (engagement, accessibility, differentiation)
    """
    from services.openai_service import clients, PRIORITY
    
    # Build context from course
    context_str = ""
    if request.course_context:
        ctx = request.course_context
        context_str = f"""
COURSE CONTEXT:
- Title: {ctx.get('title', 'Untitled')}
- Description: {ctx.get('description', 'No description')}
- Level: {ctx.get('level', 'Intermediate')}
- Topics: {', '.join(ctx.get('topics', [])[:5]) if ctx.get('topics') else 'Not specified'}
"""
    
    # Agent-specific system prompts
    agent_prompts = {
        'alexandria': """You are Alexandria, an expert curriculum designer AI. Focus on:
- Course structure and logical flow using IRIS (Immerse-Realize-Iterate-Scale) and CDIO frameworks
- Learning outcomes aligned with Bloom's Taxonomy
- Assessment strategies and rubrics
- International standards like MIT OCW structure
Provide specific, actionable suggestions with examples when relevant.""",
        
        'sme': """You are Dr. Insight, a subject matter expert AI. Focus on:
- Deep domain knowledge and accuracy
- Real-world examples and industry case studies
- Practical exercises that mirror workplace scenarios
- Current research and emerging trends in the field
- Knowledge prerequisites and learning pathways
Provide expert-level insights that add depth to course content.""",
        
        'pedagogy': """You are Professor Engage, a pedagogy expert AI. Focus on:
- Active learning and learner engagement strategies
- Differentiated instruction for diverse learners
- Accessibility and inclusive design (Universal Design for Learning)
- Formative assessment integration
- Learner motivation and retention techniques
Provide pedagogically sound recommendations that enhance learning outcomes."""
    }
    
    system_prompt = request.system_prompt or agent_prompts.get(request.agent_type, agent_prompts['alexandria'])
    
    prompt = f"""{system_prompt}

{context_str}

USER QUERY: {request.query}

Provide a helpful, specific response. Use formatting like **bold** for emphasis and numbered lists for multi-step suggestions. Be concise but thorough."""

    for provider in PRIORITY:
        if provider == "mock":
            # Return intelligent mock responses based on agent type
            mock_responses = {
                'alexandria': f"""Based on my analysis, here are my recommendations:

**Structure Suggestions:**
1. Begin with a clear "Immerse" phase focusing on context and stakeholder analysis
2. Progress to "Realize" where learners design their approach with clear milestones
3. Include hands-on "Iterate" activities with peer feedback and revision cycles
4. Conclude with "Scale" focusing on deployment, measurement, and reflection

**Learning Outcomes (Bloom's Aligned):**
- *Remember*: Define key terminology and concepts
- *Understand*: Explain how concepts interconnect
- *Apply*: Demonstrate practical implementation
- *Analyze*: Evaluate different approaches and trade-offs
- *Create*: Design original solutions for real problems

Would you like me to generate specific module structures or assessment rubrics?""",
                
                'sme': f"""Here are my expert recommendations for enriching your content:

**Content Depth:**
1. Include case studies from organizations like Google, Amazon, or local success stories
2. Reference industry-standard frameworks (e.g., Agile, Design Thinking, Lean)
3. Connect to current research and emerging trends from recent publications

**Practical Applications:**
- Design hands-on exercises that simulate real workplace challenges
- Include templates, checklists, or tools learners can adapt and use
- Build progressive project components toward the capstone

**Prerequisites Check:**
- Ensure foundational concepts are introduced or have prerequisite links
- Provide optional "refresher" resources for those who need to fill gaps

What specific topic would you like me to elaborate on?""",
                
                'pedagogy': f"""To maximize engagement and learning outcomes, consider:

**Engagement Strategies:**
1. Add Socratic AI checkpoints every 2-3 sections for reflection
2. Include collaborative discussion prompts that encourage peer learning
3. Vary content formats: combine video (< 10 min), text, and interactive elements
4. Use storytelling to make abstract concepts concrete

**Accessibility (UDL Principles):**
- Provide multiple means of representation (text + video + diagrams)
- Offer multiple means of action (typed responses, voice, project options)
- Allow multiple means of engagement (choice in topics, flexible deadlines)

**Differentiation:**
- Extension activities for advanced learners who finish early
- Scaffolding with hints and guided templates for those who need support
- Self-paced modules with checkpoints for immediate feedback

Which aspect would you like me to develop further?"""
            }
            return {
                "success": True,
                "response": mock_responses.get(request.agent_type, mock_responses['alexandria']),
                "agent": request.agent_type,
                "provider": "mock"
            }
        
        if provider not in clients:
            continue
            
        try:
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                return {
                    "success": True,
                    "response": response.text.strip(),
                    "agent": request.agent_type,
                    "provider": provider
                }
            else:
                response = await clients[provider].chat.completions.create(
                    model="gpt-4o" if provider == "openai" else "google/gemini-2.0-flash-exp:free",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"{context_str}\n\n{request.query}"}
                    ],
                    max_tokens=1000,
                    temperature=0.7
                )
                return {
                    "success": True,
                    "response": response.choices[0].message.content.strip(),
                    "agent": request.agent_type,
                    "provider": provider
                }
        except Exception as e:
            print(f"[{provider}] Agent Consult Error: {e}")
            continue
    
    return {
        "success": False,
        "response": "Unable to get agent response at this time. Please try again.",
        "agent": request.agent_type,
        "provider": "error"
    }


class EnhancedSocraticRequest(BaseModel):
    user_message: str
    phase: str = "immerse"
    context: Optional[Dict] = None
    user_role: str = "student"  # 'student' or 'instructor'
    scaffolding_level: str = "guided"  # 'exploration', 'guided', 'independent'
    conversation_history: Optional[List[Dict]] = None
    question_quality: str = "good"  # 'needs_work', 'good', 'excellent'

@app.post("/api/v1/ai/socratic-enhanced")
async def enhanced_socratic_chat(request: EnhancedSocraticRequest):
    """
    Enhanced Socratic Tutor with scaffolding based on SocraticAI research.
    
    Features:
    - Question quality-aware responses
    - Scaffolding levels (exploration, guided, independent)
    - Phase-aware guidance (IRIS framework)
    - Dual-role support (student/instructor)
    """
    from services.openai_service import clients, PRIORITY
    
    # Socratic questioning strategies
    strategies = {
        'clarifying': "What do you mean by that? Can you be more specific?",
        'probing': "Why do you think that's the case? What evidence supports this?",
        'assumptions': "What are you assuming here? Is that always true?",
        'consequences': "What might be the implications of this approach?",
        'viewpoints': "How might someone else see this differently?"
    }
    
    # Build system prompt based on role and scaffolding level
    if request.user_role == "instructor":
        system_prompt = """You are a Socratic course development assistant helping instructors create better learning experiences.

Your approach:
- Ask probing questions to help instructors think deeper about their course design
- Suggest improvements aligned with IRIS (Immerse-Realize-Iterate-Scale) framework
- Help with learning outcomes, assessments, and engagement strategies
- Never give direct answers for pedagogical choices - guide through questions

Be supportive but challenging. Push for excellence in course design."""
    else:
        scaffolding_prompts = {
            'exploration': "Be very open and encouraging. Ask broad questions. Let the learner explore freely.",
            'guided': "Provide structured guidance through questions. If the learner is stuck, offer hints. Praise good thinking.",
            'independent': "Challenge the learner. Ask tough questions. Expect well-formulated answers. Push for deeper thinking."
        }
        
        system_prompt = f"""You are a Socratic Tutor for the {request.phase} phase of learning.

CORE PRINCIPLE: Never give direct answers. Guide through questions.

{scaffolding_prompts.get(request.scaffolding_level, scaffolding_prompts['guided'])}

Socratic Techniques to Use:
1. Clarifying Questions: "What do you mean by...?"
2. Probing Questions: "Why do you think...?"
3. Assumption Testing: "What if the opposite were true?"
4. Consequence Exploration: "What would happen if...?"
5. Viewpoint Shifting: "How might others see this?"

Question Quality Response (current: {request.question_quality}):
- needs_work: Ask for clarification and context before engaging
- good: Engage with Socratic questions
- excellent: Challenge with deeper probing questions

End each response with a thought-provoking question."""

    prompt = f"""{system_prompt}

PHASE: {request.phase}
USER MESSAGE: {request.user_message}

Respond as a Socratic tutor. Remember: guide, don't tell."""

    for provider in PRIORITY:
        if provider == "mock":
            # Generate intelligent mock responses based on scaffolding level
            if request.question_quality == "needs_work":
                response = """I'd like to help you explore this! Before I can guide you effectively, could you share:

1. **What specific aspect** are you trying to understand?
2. **What do you already know** about this topic?
3. **Where exactly** are you feeling stuck?

The more context you give me, the better questions I can ask to help you discover the answers yourself. 🤔"""
            elif request.user_role == "instructor":
                response = f"""That's an interesting approach for your {request.phase} phase content! Let me ask you some questions to deepen your thinking:

**About Learning Outcomes:**
• What specific skills should learners demonstrate by the end of this section?
• How will you know they've achieved mastery?

**About Engagement:**
• How can you make this more interactive or hands-on?
• What real-world examples could bring this concept to life?

*What aspect would you like to explore further?*"""
            else:
                strategy = list(strategies.values())[hash(request.user_message) % len(strategies)]
                response = f"""🤔 That's a thoughtful question! Let me help you think through this...

{strategy}

Also consider: How does this connect to what you learned in earlier sections?

*Take a moment to reflect before responding. The goal isn't speed—it's deep understanding.*"""
            
            return {
                "success": True,
                "response": response,
                "scaffolding_level": request.scaffolding_level,
                "provider": "mock"
            }
        
        if provider not in clients:
            continue
            
        try:
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                return {
                    "success": True,
                    "response": response.text.strip(),
                    "scaffolding_level": request.scaffolding_level,
                    "provider": provider
                }
            else:
                response = await clients[provider].chat.completions.create(
                    model="gpt-4o" if provider == "openai" else "google/gemini-2.0-flash-exp:free",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": request.user_message}
                    ],
                    max_tokens=500,
                    temperature=0.7
                )
                return {
                    "success": True,
                    "response": response.choices[0].message.content.strip(),
                    "scaffolding_level": request.scaffolding_level,
                    "provider": provider
                }
        except Exception as e:
            print(f"[{provider}] Enhanced Socratic Error: {e}")
            continue
    
    return {
        "success": False,
        "response": "I'm having trouble connecting. Please try rephrasing your question.",
        "scaffolding_level": request.scaffolding_level,
        "provider": "error"
    }


@app.get("/api/v1/courses", response_model=List[schemas.Course])
def get_courses(db: Session = Depends(get_db)):
    """Get all courses"""
    return db.query(models.Course).all()

@app.post("/api/v1/courses", response_model=schemas.Course, status_code=status.HTTP_201_CREATED)
def create_course(course_data: schemas.CourseCreate, db: Session = Depends(get_db)):
    """Create a new course with validation"""
    
    new_course = models.Course(
        title=course_data.title,
        instructor=course_data.instructor,
        instructor_id=course_data.instructor_id,  # Link to instructor
        org=course_data.org,
        image=course_data.image,
        tag=course_data.tag,
        level=course_data.level,
        category=course_data.category,
        description=course_data.description,
        duration=course_data.duration,
        rating=0.0,
        students_count="0",
        institution_id=course_data.institution_id
    )
    
    db.add(new_course)
    db.commit()
    db.refresh(new_course)
    
    return new_course


# ===== INSTRUCTOR DASHBOARD ENDPOINTS =====

@app.get("/api/v1/instructor/courses")
def get_instructor_courses(user_email: str = None, db: Session = Depends(get_db)):
    """Get courses created by instructor (based on email lookup)"""
    if not user_email:
        # Return all courses if no email specified (for testing)
        courses = db.query(models.Course).all()
    else:
        # Find instructor by email
        user = db.query(models.User).filter(models.User.email == user_email).first()
        if not user:
            return []
        
        # First try by instructor_id FK
        courses = db.query(models.Course).filter(models.Course.instructor_id == user.id).all()
        
        # Fallback: match by instructor name (for seed data compatibility)
        if not courses and user.full_name:
            courses = db.query(models.Course).filter(
                models.Course.instructor == user.full_name
            ).all()
    
    return [
        {
            "id": c.id,
            "title": c.title,
            "students_count": c.students_count or "0",
            "rating": c.rating or 0.0,
            "image": c.image or "https://images.unsplash.com/photo-1516321318423-f06f85e504b3?w=400",
            "level": c.level or "Beginner",
            "category": c.category or "Technology",
            "description": c.description,
            "duration": c.duration
        }
        for c in courses
    ]


@app.get("/api/v1/instructor/grading-queue")
def get_instructor_grading_queue(user_email: str = None, db: Session = Depends(get_db)):
    """Get pending submissions for instructor's courses"""
    # Get instructor's course IDs
    instructor_course_ids = []
    if user_email:
        user = db.query(models.User).filter(models.User.email == user_email).first()
        if user:
            # Get courses by instructor_id OR instructor name
            courses = db.query(models.Course).filter(
                (models.Course.instructor_id == user.id) | (models.Course.instructor == user.full_name)
            ).all()
            instructor_course_ids = [c.id for c in courses]
    
    # Get projects from instructor's courses (or all if no filter)
    if instructor_course_ids:
        projects = db.query(models.CDIOProject).filter(
            models.CDIOProject.course_id.in_(instructor_course_ids),
            models.CDIOProject.overall_status.in_(["in_progress", "submitted"])
        ).all()
    else:
        projects = db.query(models.CDIOProject).filter(
            models.CDIOProject.overall_status.in_(["in_progress", "submitted"])
        ).all()
    
    queue = []
    for project in projects:
        # Get course title
        course = db.query(models.Course).filter(models.Course.id == project.course_id).first()
        course_title = course.title if course else "Unknown Course"
        
        # Get student name
        student = db.query(models.User).filter(models.User.id == project.user_id).first()
        student_name = student.full_name if student else "Unknown Student"
        
        # Check each phase for pending submissions
        if project.charter and not getattr(project.charter, 'grade', None):
            queue.append({
                "id": len(queue) + 1,
                "project_id": project.id,
                "student_name": student_name,
                "project_title": project.title,
                "submission_type": "charter",
                "course_title": course_title,
                "submitted_at": project.updated_at.isoformat() if project.updated_at else datetime.utcnow().isoformat(),
                "status": "pending"
            })
        
        if project.blueprint and not getattr(project.blueprint, 'grade', None):
            queue.append({
                "id": len(queue) + 1,
                "project_id": project.id,
                "student_name": student_name,
                "project_title": project.title,
                "submission_type": "design",
                "course_title": course_title,
                "submitted_at": project.updated_at.isoformat() if project.updated_at else datetime.utcnow().isoformat(),
                "status": "pending"
            })
        
        if project.implementation and not getattr(project.implementation, 'grade', None):
            queue.append({
                "id": len(queue) + 1,
                "project_id": project.id,
                "student_name": student_name,
                "project_title": project.title,
                "submission_type": "implementation",
                "course_title": course_title,
                "submitted_at": project.updated_at.isoformat() if project.updated_at else datetime.utcnow().isoformat(),
                "status": "pending"
            })
    
    return queue


@app.get("/api/v1/instructor/students")
def get_instructor_students(user_email: str = None, db: Session = Depends(get_db)):
    """Get students enrolled in instructor's courses"""
    # Get instructor's course IDs
    instructor_course_ids = []
    if user_email:
        user = db.query(models.User).filter(models.User.email == user_email).first()
        if user:
            courses = db.query(models.Course).filter(
                (models.Course.instructor_id == user.id) | (models.Course.instructor == user.full_name)
            ).all()
            instructor_course_ids = [c.id for c in courses]
    
    # Get enrollments (filtered if instructor specified)
    if instructor_course_ids:
        enrollments = db.query(models.Enrollment).filter(
            models.Enrollment.course_id.in_(instructor_course_ids)
        ).all()
    else:
        enrollments = db.query(models.Enrollment).all()
    
    students = []
    for enrollment in enrollments:
        user = db.query(models.User).filter(models.User.id == enrollment.user_id).first()
        course = db.query(models.Course).filter(models.Course.id == enrollment.course_id).first()
        
        if user and course:
            # Calculate progress from CDIO project if exists
            project = db.query(models.CDIOProject).filter(
                models.CDIOProject.user_id == user.id,
                models.CDIOProject.course_id == course.id
            ).first()
            progress = project.completion_percentage if project else 0
            
            students.append({
                "id": user.id,
                "name": user.full_name or user.email.split('@')[0],
                "email": user.email,
                "course": course.title,
                "status": enrollment.status or "active",
                "enrolled_at": enrollment.enrolled_at.strftime("%Y-%m-%d") if enrollment.enrolled_at else "N/A",
                "progress": progress,
                "grade": "N/A"
            })
    
    return students


@app.get("/")
def root():
    return {"message": "TSEA-X Mission Control Online - CDIO Framework Enabled"}


# --- CDIO Project Management Routes ---

@app.post("/api/v1/projects", response_model=schemas.CDIOProject, status_code=status.HTTP_201_CREATED)
def create_project(project_data: schemas.CDIOProjectCreate, db: Session = Depends(get_db)):
    """Create a new CDIO project for a course"""
    
    # Validate course exists
    course = db.query(models.Course).filter(models.Course.id == project_data.course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    
    # Create project
    now = datetime.utcnow()
    project = models.CDIOProject(
        course_id=project_data.course_id,
        user_id=project_data.user_id,
        title=project_data.project_title or f"{course.title} - My Project",
        current_phase="conceive",
        overall_status="in_progress",
        completion_percentage=0,
        created_at=now,
        updated_at=now
    )
    
    db.add(project)
    db.commit()
    db.refresh(project)
    
    return project


@app.get("/api/v1/projects/{project_id}", response_model=schemas.CDIOProjectDetail)
def get_project(project_id: int, db: Session = Depends(get_db)):
    """Get detailed project information with all phase data"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Fetch related artifacts (SQLAlchemy relationships handle this, but we need to ensure they are loaded)
    # Since we defined relationships in sql_models.py, accessing project.charter etc. should work if lazy loading is on,
    # or we can join them. For simplicity, let's rely on relationships.
    
    return schemas.CDIOProjectDetail(
        project=project,
        charter=project.charter,
        blueprint=project.blueprint,
        implementation=project.implementation,
        deployment=project.deployment
    )


@app.get("/api/v1/courses/{course_id}/projects", response_model=List[schemas.CDIOProject])
def get_course_projects(course_id: int, user_id: Optional[int] = None, db: Session = Depends(get_db)):
    """Get all projects for a course, optionally filtered by user"""
    query = db.query(models.CDIOProject).filter(models.CDIOProject.course_id == course_id)
    if user_id:
        query = query.filter(models.CDIOProject.user_id == user_id)
    return query.all()


# --- Conceive Phase Routes ---

@app.post("/api/v1/projects/{project_id}/charter", response_model=schemas.ProjectCharter, status_code=status.HTTP_201_CREATED)
async def create_charter(project_id: int, charter_data: schemas.ProjectCharterCreate, db: Session = Depends(get_db)):
    """Create or update project charter (Conceive phase)"""
    
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Generate AI suggestions
    suggestions = await openai_service.generate_charter_suggestions(
        problem_statement=charter_data.problem_statement,
        success_metrics=charter_data.success_metrics,
        course_id=project.course_id,
        target_outcome=charter_data.target_outcome
    )
    
    now = datetime.utcnow()
    
    # Check if charter exists
    charter = db.query(models.ProjectCharter).filter(models.ProjectCharter.project_id == project_id).first()
    
    if charter:
        # Update existing
        charter.problem_statement = charter_data.problem_statement
        charter.success_metrics = charter_data.success_metrics
        charter.target_outcome = charter_data.target_outcome
        charter.constraints = charter_data.constraints
        charter.stakeholders = charter_data.stakeholders
        # Update AI fields if not set or if we want to overwrite? Let's overwrite for now as it's a "re-conception"
        charter.suggested_tools = suggestions.get("suggested_tools", [])
        charter.reasoning = suggestions.get("reasoning")
        charter.estimated_duration = suggestions.get("estimated_duration")
        charter.difficulty_level = suggestions.get("difficulty_level")
    else:
        # Create new
        charter = models.ProjectCharter(
            project_id=project_id,
            problem_statement=charter_data.problem_statement,
            success_metrics=charter_data.success_metrics,
            target_outcome=charter_data.target_outcome,
            constraints=charter_data.constraints,
            stakeholders=charter_data.stakeholders,
            suggested_tools=suggestions.get("suggested_tools", []),
            reasoning=suggestions.get("reasoning"),
            estimated_duration=suggestions.get("estimated_duration"),
            difficulty_level=suggestions.get("difficulty_level")
        )
        db.add(charter)
        db.flush() # Get ID

    
    # Update project status

    project.current_phase = schemas.CDIOPhase.DESIGN
    project.completion_percentage = 25
    project.updated_at = now
    
    # --- PKC Auto-Capture Trigger ---
    try:
        # Create Knowledge Node for this Charter
        node_title = f"Charter: {project.title}"
        node_content = f"Problem: {charter.problem_statement}\nMetrics: {charter.success_metrics}\nOutcome: {charter.target_outcome}"
        
        # Generate embedding
        embedding = await generate_embedding(f"{node_title}\n{node_content}")
        
        # Check if node already exists for this source
        source_id = f"project_{project_id}_charter"
        existing_node = db.query(models.KnowledgeNode).filter(
            models.KnowledgeNode.user_id == project.user_id,
            models.KnowledgeNode.source_id == source_id
        ).first()
        
        if existing_node:
            existing_node.title = node_title
            existing_node.content = node_content
            existing_node.embedding = embedding
            existing_node.updated_at = now
        else:
            new_node = models.KnowledgeNode(
                user_id=project.user_id,
                title=node_title,
                content=node_content,
                node_type="project_artifact",
                source_id=source_id,
                embedding=embedding,
                created_at=now,
                updated_at=now
            )
            db.add(new_node)
            
    except Exception as e:
        print(f"PKC Auto-Capture Error: {e}")
        # Don't fail the request if PKC fails
    
    db.commit()
    db.refresh(charter)
    return charter


@app.get("/api/v1/projects/{project_id}/charter", response_model=schemas.ProjectCharter)
def get_charter(project_id: int, db: Session = Depends(get_db)):
    """Get project charter"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if not project.charter:
        raise HTTPException(status_code=404, detail="Charter not found for this project")
    
    return project.charter


# --- Design Phase Routes (Stubs for now) ---

@app.post("/api/v1/projects/{project_id}/blueprint", response_model=schemas.DesignBlueprint)
async def create_blueprint(project_id: int, blueprint_data: schemas.DesignBlueprintCreate, db: Session = Depends(get_db)):
    """Create or update design blueprint"""
    
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    now = datetime.utcnow()
    
    # Check if blueprint exists
    blueprint = db.query(models.DesignBlueprint).filter(models.DesignBlueprint.project_id == project_id).first()
    
    if blueprint:
        # Update existing
        blueprint.architecture_diagram = blueprint_data.architecture_diagram
        blueprint.logic_flow = blueprint_data.logic_flow
        blueprint.component_list = blueprint_data.component_list or []
        # blueprint.data_flow = blueprint_data.data_flow # Not in SQL model yet? Check sql_models.py
    else:
        # Create new
        blueprint = models.DesignBlueprint(
            project_id=project_id,
            architecture_diagram=blueprint_data.architecture_diagram,
            logic_flow=blueprint_data.logic_flow,
            component_list=blueprint_data.component_list or []
        )
        db.add(blueprint)
        db.flush()

    
    # Update project

    project.current_phase = schemas.CDIOPhase.IMPLEMENT
    project.completion_percentage = 50
    project.updated_at = now
    
    # --- PKC Auto-Capture Trigger ---
    try:
        # Create Knowledge Node for this Blueprint
        node_title = f"Blueprint: {project.title}"
        node_content = f"Logic Flow: {blueprint.logic_flow}\nComponents: {', '.join(blueprint.component_list or [])}"
        
        # Generate embedding
        embedding = await generate_embedding(f"{node_title}\n{node_content}")
        
        # Check if node already exists
        source_id = f"project_{project_id}_blueprint"
        existing_node = db.query(models.KnowledgeNode).filter(
            models.KnowledgeNode.user_id == project.user_id,
            models.KnowledgeNode.source_id == source_id
        ).first()
        
        if existing_node:
            existing_node.title = node_title
            existing_node.content = node_content
            existing_node.embedding = embedding
            existing_node.updated_at = now
        else:
            new_node = models.KnowledgeNode(
                user_id=project.user_id,
                title=node_title,
                content=node_content,
                node_type="project_artifact",
                source_id=source_id,
                embedding=embedding,
                created_at=now,
                updated_at=now
            )
            db.add(new_node)
            
    except Exception as e:
        print(f"PKC Auto-Capture Error: {e}")
    
    db.commit()
    db.refresh(blueprint)
    return blueprint


# --- Implementation Phase Routes ---

@app.post("/api/v1/projects/{project_id}/implementation", response_model=schemas.Implementation)
def create_implementation(project_id: int, impl_data: schemas.ImplementationCreate, db: Session = Depends(get_db)):
    """Submit implementation (code)"""
    
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    now = datetime.utcnow()
    
    # Check if implementation exists
    implementation = db.query(models.Implementation).filter(models.Implementation.project_id == project_id).first()
    
    if implementation:
        # Update existing
        implementation.code_snapshot = impl_data.code_snapshot
        implementation.notes = impl_data.notes
        implementation.updated_at = now
    else:
        # Create new
        implementation = models.Implementation(
            project_id=project_id,
            code_snapshot=impl_data.code_snapshot,
            notes=impl_data.notes,
            created_at=now,
            updated_at=now
        )
        db.add(implementation)
        db.flush()

    
    # Mock validation
    validation = implementation_service.validate_submission(impl_data.code_snapshot or "")
    implementation.security_check_passed = validation["passed"]
    
    # Update project status

    project.current_phase = schemas.CDIOPhase.OPERATE
    project.completion_percentage = 75
    project.updated_at = now
    
    db.commit()
    db.refresh(implementation)
    return implementation

@app.get("/api/v1/projects/{project_id}/implementation", response_model=schemas.Implementation)
def get_implementation(project_id: int, db: Session = Depends(get_db)):
    """Get implementation details"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
        
    if not project.implementation:
        raise HTTPException(status_code=404, detail="Implementation not found")
        
    return project.implementation


# --- Operate Phase Routes ---

# --- AI-Powered Features ---

from services import openai_service, implementation_service, blockchain_service
from pydantic import BaseModel

# --- Operate Phase Routes ---



# --- Code Persistence Routes ---

@app.post("/api/v1/projects/{project_id}/save-code", response_model=schemas.CodeSnapshot)
def save_code(project_id: int, snapshot_data: schemas.CodeSnapshotCreate, db: Session = Depends(get_db)):
    """Save a code snapshot"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
        
    now = datetime.utcnow()
    
    snapshot = models.CodeSnapshot(
        project_id=project_id,
        code=snapshot_data.code,
        language=snapshot_data.language,
        timestamp=now,
        auto_saved=snapshot_data.auto_saved,
        file_name=f"main.{'py' if snapshot_data.language == 'python' else 'js' if snapshot_data.language == 'javascript' else 'java'}"
    )
    
    db.add(snapshot)
    
    # Also update the latest implementation
    implementation = db.query(models.Implementation).filter(models.Implementation.project_id == project_id).first()
    if implementation:
        implementation.code_snapshot = snapshot_data.code
        implementation.updated_at = now
    else:
        implementation = models.Implementation(
            project_id=project_id,
            code_snapshot=snapshot_data.code,
            created_at=now,
            updated_at=now
        )
        db.add(implementation)
        project.implementation_id = implementation.id
        
    db.commit()
    db.refresh(snapshot)
    
    return snapshot

@app.get("/api/v1/projects/{project_id}/snapshots", response_model=List[schemas.CodeSnapshot])
def get_snapshots(project_id: int, db: Session = Depends(get_db)):
    """Get code version history"""
    snapshots = db.query(models.CodeSnapshot).filter(
        models.CodeSnapshot.project_id == project_id
    ).order_by(models.CodeSnapshot.timestamp.desc()).limit(50).all()
    return snapshots

@app.post("/api/v1/projects/{project_id}/deployment", response_model=schemas.Deployment)
async def create_deployment(project_id: int, deploy_data: schemas.DeploymentCreate, db: Session = Depends(get_db)):
    """Submit deployment details (Operate phase)"""
    
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    now = datetime.utcnow()
    
    # Check if deployment exists
    deployment = db.query(models.Deployment).filter(models.Deployment.project_id == project_id).first()
    
    if deployment:
        # Update existing
        deployment.deployment_url = deploy_data.deployment_url
        deployment.deployment_platform = deploy_data.deployment_platform
        deployment.readme = deploy_data.readme
        deployment.updated_at = now
    else:
        # Create new
        deployment = models.Deployment(
            project_id=project_id,
            deployment_url=deploy_data.deployment_url,
            deployment_platform=deploy_data.deployment_platform,
            verification_status="submitted"
        )
        db.add(deployment)
        db.flush()
        project.deployment_id = deployment.id
    
    # Update project status
    project.operate_completed = True
    project.current_phase = schemas.CDIOPhase.OPERATE 
    project.overall_status = schemas.ProjectStatus.COMPLETED
    project.completion_percentage = 100
    project.updated_at = now
    # project.completed_at = now # Not in SQL model

    # --- Blockchain Integration (SBT Minting) ---
    # In a real app, this might happen asynchronously via a queue
    try:
        mint_result = await blockchain_service.mint_credential(
            user_id=project.user_id,
            course_id=project.course_id,
            project_id=project.id
        )
        
        deployment.sbt_minted = True
        deployment.sbt_token_id = mint_result["token_id"]
        deployment.transaction_hash = mint_result["transaction_hash"]
        deployment.explorer_url = mint_result["explorer_url"]
        
    except Exception as e:
        print(f"Minting failed: {e}")
        # Don't fail the request, just log it. 
        # In production, we'd have a retry mechanism.
    
    db.commit()
    db.refresh(deployment)
    return deployment

class RunCodeRequest(BaseModel):
    code: str
    language: str = "python"

@app.post("/api/v1/implementation/run")
async def run_code(request: RunCodeRequest):
    """Mock execution of student code"""
    return await implementation_service.run_code_mock(request.code, request.language)


class SaveCodeRequest(BaseModel):
    project_id: int
    code: str
    language: str
    auto_saved: bool = False

@app.post("/api/v1/projects/{project_id}/save-code")
def save_code(project_id: int, request: SaveCodeRequest, db: Session = Depends(get_db)):
    """Save code snapshot and update implementation"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Create snapshot
    snapshot = models.CodeSnapshot(
        project_id=project_id,
        code=request.code,
        language=request.language,
        auto_saved=request.auto_saved,
        timestamp=datetime.utcnow()
    )
    db.add(snapshot)
    
    # Update or Create Implementation
    impl = db.query(models.Implementation).filter(models.Implementation.project_id == project_id).first()
    if impl:
        impl.code_snapshot = request.code
        impl.updated_at = datetime.utcnow()
    else:
        impl = models.Implementation(
            project_id=project_id,
            code_snapshot=request.code,
            created_at=datetime.utcnow(),
            updated_at=datetime.utcnow()
        )
        db.add(impl)
        db.flush() # get ID
        project.implementation = impl # Link if relationship not auto-handled
    
    db.commit()
    
    return {"success": True, "snapshot_id": snapshot.id}


class CharterSuggestionsRequest(BaseModel):
    problem_statement: str
    success_metrics: str
    course_id: int
    target_outcome: Optional[str] = None

class SocraticChatRequest(BaseModel):
    project_id: int
    user_message: str
    conversation_history: List[dict]
    design_context: dict

@app.post("/api/v1/ai/charter-suggestions")
async def get_charter_suggestions(request: CharterSuggestionsRequest):
    """
    Generate AI-powered charter suggestions using GPT-4
    """
    try:
        suggestions = await openai_service.generate_charter_suggestions(
            problem_statement=request.problem_statement,
            success_metrics=request.success_metrics,
            course_id=request.course_id,
            target_outcome=request.target_outcome
        )
        return suggestions
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"AI service error: {str(e)}"
        )

@app.post("/api/v1/ai/socratic-chat")
async def socratic_chat(request: SocraticChatRequest):
    """
    Get Socratic tutor response for design phase
    """
    try:
        response = await openai_service.socratic_response(
            user_message=request.user_message,
            design_context=request.design_context,
            conversation_history=request.conversation_history,
            project_id=request.project_id
        )
        return response
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"AI service error: {str(e)}"
        )

class IterationSubmission(schemas.IterationArtifactCreate):
    user_id: int

@app.post("/api/v1/projects/{course_id}/iteration")
def submit_iteration(course_id: int, submission: IterationSubmission, db: Session = Depends(get_db)):
    """Submit BML iteration"""
    # Find project for user + course
    project = db.query(models.CDIOProject).filter(
        models.CDIOProject.course_id == course_id,
        models.CDIOProject.user_id == submission.user_id
    ).first()
    
    if not project:
         # Auto-create project for demo flow if missing (SKIP logic for brevity, assume exists or throw 404)
         raise HTTPException(status_code=404, detail="Project not found")

    # Update Implementation/Iteration Artifact
    impl = project.implementation
    if not impl:
        impl = models.Implementation(project_id=project.id)
        db.add(impl)
    
    impl.iteration_number = submission.iteration_number
    impl.hypothesis = submission.hypothesis
    impl.code_repository_url = submission.prototype_url
    impl.code_snapshot = submission.code_snapshot
    impl.learnings = submission.learnings
    impl.next_hypothesis = submission.next_hypothesis
    # impl.measurements = submission.measurements 
    impl.updated_at = datetime.utcnow()
    
    db.commit()
    db.refresh(impl)
    return impl

@app.post("/api/v1/ai/grade-iteration/{iteration_id}")
async def grade_iteration(iteration_id: int, db: Session = Depends(get_db)):
    """Grade an iteration artifact using AI"""
    iteration = db.query(models.IterationArtifact).filter(models.IterationArtifact.id == iteration_id).first()
    if not iteration:
        raise HTTPException(status_code=404, detail="Iteration artifact not found")
        
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == iteration.project_id).first()
    if not project:
         raise HTTPException(status_code=404, detail="Project not found")
         
    user = db.query(models.User).filter(models.User.id == project.user_id).first()
    student_name = user.full_name if user else "Student"

    submission_content = {
        "hypothesis": iteration.hypothesis,
        "learnings": iteration.learnings,
        "code_snapshot": iteration.code_snapshot,
        "next_hypothesis": iteration.next_hypothesis,
        "measurements": iteration.measurements
    }

    try:
        from services import openai_service
        feedback = await openai_service.generate_grading_feedback(
            submission_type="Iteration (Build-Measure-Learn)",
            project_title=project.title,
            student_name=student_name,
            submission_content=submission_content
        )
        
        iteration.ai_feedback = feedback
        db.commit()
        
        return feedback
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI Grading failed: {str(e)}")

@app.get("/api/v1/instructor/grading-queue", response_model=List[GradingQueueItem])
def get_grading_queue(instructor_id: Optional[int] = None, db: Session = Depends(get_db)):
    """Get all submissions requiring grading or recently graded"""
    log_path = r"C:\Users\PT\Desktop\TSEA-X\api_debug.log"
    try:
        with open(log_path, "a") as f:
            f.write(f"\n[{datetime.now()}] get_grading_queue called\n")
    except Exception as e:
        print(f"LOGGING ERROR: {e}")
        
    queue = []
    
    # Get all projects (in real app, filter by instructor's courses)
    projects = db.query(models.CDIOProject).all()
    
    try:
        with open(log_path, "a") as f:
            f.write(f"[{datetime.now()}] Found {len(projects)} projects in DB\n")
    except: pass
    
    for project in projects:
        try:
            with open(log_path, "a") as f:
                 f.write(f"[{datetime.now()}] Proc Project {project.id}: Charter={project.charter is not None}, Blueprint={project.blueprint is not None}\n")
        except: pass

        course = db.query(models.Course).filter(models.Course.id == project.course_id).first()
        course_title = course.title if course else "Unknown Course"
        student_name = f"Student {project.user_id}" # Mock name lookup
        
        # Check Charter
        print(f"DEBUG: Check Charter for {project.id}: {project.charter}")
        if project.charter:
            status = "Pending"
            # Logic to determine if graded: e.g. if feedback exists. 
            # For now, simplistic: if overall status is NOT under_review, it's graded.
            if project.overall_status != schemas.ProjectStatus.UNDER_REVIEW:
                status = "Graded"

            print(f"DEBUG: Adding Charter item for {project.id}")
            queue.append(GradingQueueItem(
                id=project.charter.id,
                project_id=project.id,
                student_name=student_name,
                project_title=project.title,
                submission_type="charter",
                course_title=course_title,
                submitted_at=project.charter.created_at.strftime("%Y-%m-%d %H:%M"),
                status=status
            ))

        # Check Blueprint (Frontend expects 'design')
        print(f"DEBUG: Check Blueprint for {project.id}: {project.blueprint}")
        if project.blueprint:
            status = "Pending"
            if project.design_completed:
                status = "Graded"
                
            print(f"DEBUG: Adding Design item for {project.id}")
            queue.append(GradingQueueItem(
                id=project.blueprint.id,
                project_id=project.id,
                student_name=student_name,
                project_title=project.title,
                submission_type="design",
                course_title=course_title,
                submitted_at=project.blueprint.created_at.strftime("%Y-%m-%d %H:%M"),
                status=status
            ))

        # Check Implementation
        if project.implementation:
            status = "Pending"
            if project.implement_completed:
                status = "Graded"
                
            queue.append(GradingQueueItem(
                id=project.implementation.id,
                project_id=project.id,
                student_name=student_name,
                project_title=project.title,
                submission_type="implementation",
                course_title=course_title,
                submitted_at=project.implementation.created_at.strftime("%Y-%m-%d %H:%M"),
                status=status
            ))

    print(f"DEBUG: Returning queue with {len(queue)} items")
    # Sort by status (Pending first) then date
    # queue.sort(key=lambda x: x.submitted_at, reverse=True)
    queue.sort(key=lambda x: x.status == "Pending", reverse=True)
    
    return queue


@app.post("/api/v1/projects/{course_id}/scale")
def submit_scale(course_id: int, submission: schemas.ScaleArtifactCreate, user_id: int = Query(...), db: Session = Depends(get_db)):
    """Submit Scale phase artifact (Institutional Deployment)"""
    project = db.query(models.CDIOProject).filter(
        models.CDIOProject.course_id == course_id,
        models.CDIOProject.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    # Update or Create Deployment
    deployment = project.deployment
    if not deployment:
        deployment = models.Deployment(project_id=project.id)
        db.add(deployment)
    
    if submission.deployment_url:
        deployment.deployment_url = submission.deployment_url
    if submission.deployment_platform:
        deployment.deployment_platform = submission.deployment_platform
        
    deployment.institutional_handoff = submission.institutional_handoff
    deployment.stakeholder_training = submission.stakeholder_training
    deployment.impact_metrics = submission.impact_metrics
    deployment.sfia_evidence = submission.sfia_evidence
    
    deployment.verification_status = "submitted"
    deployment.updated_at = datetime.utcnow()
    
    # Auto-advance project status
    project.current_phase = "scale"
    project.overall_status = "completed"
    project.completion_percentage = 100
    
    db.commit()
    db.refresh(deployment)
    return deployment


class CompanionChatRequest(BaseModel):
    user_message: str
    current_page: str
    user_context: Dict = {}
    conversation_history: List[Dict] = []

@app.post("/api/v1/ai/companion-chat")
async def companion_chat(request: CompanionChatRequest):
    """
    Get response from Global AI Companion
    """
    try:
        response = await openai_service.companion_response(
            user_message=request.user_message,
            current_page=request.current_page,
            user_context=request.user_context,
            conversation_history=request.conversation_history
        )
        return response
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"AI service error: {str(e)}"
        )


# --- Health Check ---

@app.get("/health")
def health_check(db: Session = Depends(get_db)):
    return {
        "status": "healthy",
        "cdio_enabled": True,
        "ai_enabled": True,
        "projects_count": db.query(models.CDIOProject).count(),
        "charters_count": db.query(models.ProjectCharter).count()
    }


# --- Quiz Routes ---

from models import Quiz, QuizSubmission, QuizSubmissionCreate

@app.get("/api/v1/courses/{course_id}/quizzes", response_model=List[Quiz])
def get_course_quizzes(course_id: int):
    """Get all quizzes for a course"""
    quizzes = [q for q in QUIZZES_DB.values() if q.course_id == course_id]
    return quizzes

@app.get("/api/v1/quizzes/{quiz_id}", response_model=Quiz)
def get_quiz(quiz_id: int):
    """Get quiz details with questions"""
    quiz = QUIZZES_DB.get(quiz_id)
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz not found")
    return quiz

@app.post("/api/v1/quizzes/{quiz_id}/submit", response_model=schemas.QuizSubmission)
def submit_quiz(quiz_id: int, submission: schemas.QuizSubmissionCreate):
    """Submit quiz answers and get results"""
    global quiz_submission_counter
    
    quiz = QUIZZES_DB.get(quiz_id)
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz not found")
    
    # Calculate score
    correct_count = 0
    for question in quiz.questions:
        user_answer = submission.answers.get(question.id)
        if user_answer == question.correct_answer:
            correct_count += 1
    
    total_questions = len(quiz.questions)
    passing_score = int(total_questions * 0.7)  # 70% to pass
    
    result = schemas.QuizSubmission(
        id=quiz_submission_counter,
        quiz_id=quiz_id,
        user_id=submission.user_id,
        answers=submission.answers,
        score=correct_count,
        total_questions=total_questions,
        passed=correct_count >= passing_score,
        submitted_at=datetime.now()
    )
    
    QUIZ_SUBMISSIONS_DB[quiz_submission_counter] = result
    quiz_submission_counter += 1
    
    return result


# --- Discussion Routes ---

# from models import DiscussionThread, DiscussionThreadCreate, DiscussionComment, DiscussionCommentCreate

@app.get("/api/v1/courses/{course_id}/discussions", response_model=List[schemas.DiscussionThread])
def get_course_discussions(course_id: int):
    """Get all discussion threads for a course"""
    threads = [t for t in DISCUSSIONS_DB.values() if t.course_id == course_id]
    # Sort by most recent
    threads.sort(key=lambda x: x.created_at, reverse=True)
    return threads

@app.post("/api/v1/courses/{course_id}/discussions", response_model=schemas.DiscussionThread)
def create_discussion(course_id: int, thread_data: schemas.DiscussionThreadCreate):
    """Create a new discussion thread"""
    global discussion_counter
    
    thread = schemas.DiscussionThread(
        id=discussion_counter,
        course_id=course_id,
        user_id=thread_data.user_id,
        author_name=thread_data.author_name,
        title=thread_data.title,
        content=thread_data.content,
        created_at=datetime.now(),
        comments_count=0,
        likes=0,
        comments=[]
    )
    
    DISCUSSIONS_DB[discussion_counter] = thread
    discussion_counter += 1
    
    return thread

@app.post("/api/v1/discussions/{thread_id}/comments", response_model=schemas.DiscussionComment)
def add_comment(thread_id: int, comment_data: schemas.DiscussionCommentCreate):
    """Add a comment to a discussion thread"""
    global comment_counter
    
    thread = DISCUSSIONS_DB.get(thread_id)
    if not thread:
        raise HTTPException(status_code=404, detail="Thread not found")
    
    comment = schemas.DiscussionComment(
        id=comment_counter,
        thread_id=thread_id,
        user_id=comment_data.user_id,
        author_name=comment_data.author_name,
        content=comment_data.content,
        created_at=datetime.now(),
        likes=0
    )
    
    DISCUSSION_COMMENTS_DB[comment_counter] = comment
    thread.comments.append(comment)
    thread.comments_count += 1
    comment_counter += 1
    
    return comment


# --- Notification Routes ---

@app.get("/api/v1/notifications", response_model=List[schemas.Notification])
def get_notifications(user_id: int = 1): # Default to user 1 for now
    """Get all notifications for a user"""
    notifs = [n for n in NOTIFICATIONS_DB.values() if n.user_id == user_id]
    notifs.sort(key=lambda x: x.created_at, reverse=True)
    return notifs

@app.post("/api/v1/notifications/{notification_id}/read")
def mark_notification_read(notification_id: int):
    """Mark a notification as read"""
    notif = NOTIFICATIONS_DB.get(notification_id)
    if not notif:
        raise HTTPException(status_code=404, detail="Notification not found")
    
    notif.is_read = True
    return {"status": "success"}


# --- Inbox Routes ---

@app.get("/api/v1/inbox/conversations", response_model=List[schemas.Conversation])
def get_conversations(user_id: int = 1):
    """Get all conversations for a user"""
    convs = [c for c in CONVERSATIONS_DB.values() if user_id in c.participants]
    convs.sort(key=lambda x: x.updated_at, reverse=True)
    return convs

@app.get("/api/v1/inbox/conversations/{conversation_id}/messages", response_model=List[schemas.Message])
def get_messages(conversation_id: int, user_id: int = 1):
    """Get messages for a conversation"""
    conv = CONVERSATIONS_DB.get(conversation_id)
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    
    if user_id not in conv.participants:
        raise HTTPException(status_code=403, detail="Not a participant")
        
    return conv.messages

class MessageCreate(BaseModel):
    content: str

@app.post("/api/v1/inbox/conversations/{conversation_id}/messages", response_model=schemas.Message)
def send_message(conversation_id: int, message_data: schemas.MessageCreate, user_id: int = 1):
    """Send a message in a conversation"""
    global message_counter
    
    conv = CONVERSATIONS_DB.get(conversation_id)
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    
    if user_id not in conv.participants:
        raise HTTPException(status_code=403, detail="Not a participant")
    
    msg = schemas.Message(
        id=message_counter,
        sender_id=user_id,
        sender_name="You", # Simplified
        content=message_data.content,
        created_at=datetime.now(),
        is_read=False
    )
    
    conv.messages.append(msg)
    conv.last_message = msg
    conv.updated_at = datetime.now()
    
    message_counter += 1
    
    return msg



@app.get("/api/v1/test-ai")
async def test_ai_endpoint():
    """Test the AI connection"""
    return await test_ai_connection()


# ===== Student Dashboard Endpoints =====

@app.get("/api/v1/student/dashboard", response_model=schemas.StudentDashboardData)
def get_student_dashboard(user_id: int = 1, db: Session = Depends(get_db)):
    """Get aggregated dashboard data for a student"""
    
    # Get user to retrieve career_pathway_id
    user = db.query(models.User).filter(models.User.id == user_id).first()
    career_pathway_id = user.career_pathway_id if user else None
    
    # 1. Get Enrolled Courses (Active Projects)
    enrolled_courses = []
    total_progress = 0
    course_count = 0
    
    # Find projects for this user
    user_projects = db.query(models.CDIOProject).filter(models.CDIOProject.user_id == user_id).all()
    
    for project in user_projects:
        course = db.query(models.Course).filter(models.Course.id == project.course_id).first()
        if course:
            # Calculate progress
            progress = project.completion_percentage
            
            enrolled_courses.append(schemas.DashboardCourseProgress(
                course_id=course.id,
                project_id=project.id,
                title=course.title,
                image=course.image,
                category=course.tag or "General",
                progress=progress,
                current_phase=project.current_phase.title() if hasattr(project.current_phase, 'title') else str(project.current_phase).title(),
                remaining_time="4h remaining" # Mock
            ))
            total_progress += progress
            course_count += 1
            
    avg_progress = int(total_progress / course_count) if course_count > 0 else 0
    
    # 2. Get Credentials
    credentials = []
    user_creds = db.query(models.Credential).filter(models.Credential.user_id == user_id).all()
    
    for cred in user_creds:
        credentials.append(schemas.DashboardCredential(
            id=cred.id,
            title=cred.title,
            org=cred.issuer_name,
            date=cred.issued_at.strftime("%b %Y") if cred.issued_at else "Pending",
            hash=cred.transaction_hash or "Pending"
        ))
        
    # 3. Recommended Courses (Simple logic: courses not enrolled in)
    enrolled_ids = [p.course_id for p in user_projects]
    recommended_courses = db.query(models.Course).filter(models.Course.id.notin_(enrolled_ids)).limit(3).all()
    
    # 4. Upcoming Deadlines (Mock)
    upcoming_deadlines = [
        {"id": 1, "title": "Charter Submission", "course": "Sustainable Development", "date": "Tomorrow, 11:59 PM", "urgent": True},
        {"id": 2, "title": "Peer Review", "course": "AI Governance", "date": "Dec 15, 2024", "urgent": False}
    ]
    
    return schemas.StudentDashboardData(
        enrolled_courses=enrolled_courses,
        credentials=credentials,
        total_learning_hours=42, # Mock
        average_progress=avg_progress,
        recommended_courses=recommended_courses,
        upcoming_deadlines=upcoming_deadlines,
        career_pathway_id=career_pathway_id
    )


class CareerPathwayUpdate(BaseModel):
    career_pathway_id: Optional[str] = None

@app.put("/api/v1/users/{user_id}/career-pathway")
def update_career_pathway(user_id: int, data: CareerPathwayUpdate, db: Session = Depends(get_db)):
    """Update user's selected career pathway"""
    user = db.query(models.User).filter(models.User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    user.career_pathway_id = data.career_pathway_id
    db.commit()
    db.refresh(user)
    
    return {
        "success": True,
        "message": f"Career pathway updated to: {data.career_pathway_id or 'None'}",
        "user_id": user.id,
        "career_pathway_id": user.career_pathway_id
    }


# ===== Enrollment Endpoints =====

@app.post("/api/v1/enrollments", response_model=schemas.Enrollment, status_code=status.HTTP_201_CREATED)
def create_enrollment(enrollment_data: schemas.EnrollmentCreate, db: Session = Depends(get_db)):
    """Enroll a student in a course"""
    
    # Check if course exists
    course = db.query(models.Course).filter(models.Course.id == enrollment_data.course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    
    # Check if already enrolled
    existing_enrollment = db.query(models.Enrollment).filter(
        models.Enrollment.user_id == enrollment_data.user_id,
        models.Enrollment.course_id == enrollment_data.course_id
    ).first()
    
    if existing_enrollment:
        raise HTTPException(status_code=400, detail="User already enrolled in this course")
    
    try:
        # Create enrollment
        new_enrollment = models.Enrollment(
            user_id=enrollment_data.user_id,
            course_id=enrollment_data.course_id,
            status="active",
            enrolled_at=datetime.utcnow()
        )
        db.add(new_enrollment)
        db.commit()
        db.refresh(new_enrollment)
        
        # Auto-create Project for this course
        existing_project = db.query(models.CDIOProject).filter(
            models.CDIOProject.user_id == enrollment_data.user_id,
            models.CDIOProject.course_id == enrollment_data.course_id
        ).first()
        
        if not existing_project:
            new_project = models.CDIOProject(
                course_id=enrollment_data.course_id,
                user_id=enrollment_data.user_id,
                title=f"{course.title} - My Project",
                current_phase="conceive",
                overall_status="in_progress",
                completion_percentage=0,
                created_at=datetime.utcnow(),
                updated_at=datetime.utcnow(),
                started_at=datetime.utcnow(),
                last_activity_at=datetime.utcnow()
            )
            db.add(new_project)
            db.commit()
            print(f"DEBUG: Created new project {new_project.id} for user {new_project.user_id}")
        else:
            print(f"DEBUG: Found existing project {existing_project.id} for user {existing_project.user_id}")
        
        return new_enrollment
    except Exception as e:
        print(f"ERROR in enroll_course: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/api/v1/enrollments/{enrollment_id}", status_code=status.HTTP_204_NO_CONTENT)
@app.delete("/api/v1/enrollments/{enrollment_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_enrollment(enrollment_id: int, db: Session = Depends(get_db)):
    """Unenroll a student from a course"""
    enrollment = db.query(models.Enrollment).filter(models.Enrollment.id == enrollment_id).first()
    if not enrollment:
        raise HTTPException(status_code=404, detail="Enrollment not found")
    
    # Instead of deleting, mark as dropped
    enrollment.status = "dropped"
    db.commit()
    return None


@app.get("/api/v1/users/{user_id}/enrollments", response_model=List[schemas.Enrollment])
@app.get("/api/v1/users/{user_id}/enrollments", response_model=List[schemas.Enrollment])
def get_user_enrollments(user_id: int, status_filter: Optional[str] = None, db: Session = Depends(get_db)):
    """Get all enrollments for a user"""
    query = db.query(models.Enrollment).filter(models.Enrollment.user_id == user_id)
    
    if status_filter:
        query = query.filter(models.Enrollment.status == status_filter)
    
    return query.all()


@app.get("/api/v1/users/{user_id}/projects", response_model=List[schemas.CDIOProject])
def get_user_projects(user_id: int, db: Session = Depends(get_db)):
    """Get all projects for a user"""
    projects = db.query(models.CDIOProject).filter(models.CDIOProject.user_id == user_id).all()
    return projects


@app.get("/api/v1/courses/{course_id}/enrollments")

def get_course_enrollments(course_id: int, db: Session = Depends(get_db)):
    """Get enrollment count and status for a course"""
    enrollments = db.query(models.Enrollment).filter(
        models.Enrollment.course_id == course_id,
        models.Enrollment.status == "active"
    ).all()
    
    return {
        "course_id": course_id,
        "total_enrolled": len(enrollments),
        "enrollments": enrollments
    }


@app.get("/api/v1/enrollments/check")
def check_enrollment(user_id: int, course_id: int, db: Session = Depends(get_db)):
    """Check if a user is enrolled in a course"""
    enrollment = db.query(models.Enrollment).filter(
        models.Enrollment.user_id == user_id,
        models.Enrollment.course_id == course_id,
        models.Enrollment.status == "active"
    ).first()
    
    return {
        "enrolled": enrollment is not None,
        "enrollment": enrollment
    }


# ===== Blockchain Credential Endpoints =====

@app.post("/api/v1/credentials", response_model=schemas.Credential, status_code=status.HTTP_201_CREATED)

def create_credential(credential_data: schemas.CredentialCreate, db: Session = Depends(get_db)):
    """Create a new blockchain credential (initiates minting process)"""
    
    # Validate course/project exists if provided
    if credential_data.course_id:
        course = db.query(models.Course).filter(models.Course.id == credential_data.course_id).first()
        if not course:
            raise HTTPException(status_code=404, detail="Course not found")
    
    if credential_data.project_id:
        project = db.query(models.CDIOProject).filter(models.CDIOProject.id == credential_data.project_id).first()
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")
    
    # Create credential record
    now = datetime.utcnow()
    
    new_credential = models.Credential(
        user_id=credential_data.user_id,
        course_id=credential_data.course_id,
        # project_id=credential_data.project_id, # Note: Credential model in sql_models.py doesn't have project_id yet, assuming generic credential
        title=credential_data.title,
        description=credential_data.description,
        issuer_name=credential_data.issuer_name,
        credential_type=credential_data.credential_type,
        wallet_address=credential_data.wallet_address,
        status="minting",
        blockchain_network="Polygon Mumbai Testnet",
        issued_at=now
    )
    
    db.add(new_credential)
    db.commit()
    db.refresh(new_credential)
    
    return new_credential


@app.post("/api/v1/credentials/{credential_id}/mint")
def mint_credential(credential_id: int, db: Session = Depends(get_db)):
    """Mint the SBT on blockchain (simulated)"""
    credential = db.query(models.Credential).filter(models.Credential.id == credential_id).first()
    if not credential:
        raise HTTPException(status_code=404, detail="Credential not found")
    
    if credential.status == "issued":
        raise HTTPException(status_code=400, detail="Credential already minted")
    
    # Simulate blockchain minting
    import hashlib
    import secrets
    
    # Generate mock transaction hash
    tx_data = f"{credential.id}{credential.user_id}{credential.title}{secrets.token_hex(16)}"
    tx_hash = hashlib.sha256(tx_data.encode()).hexdigest()
    
    # Generate mock token ID
    token_id = f"T6C-{credential.id:06d}"
    
    # Update credential
    credential.status = "issued"
    credential.token_id = token_id
    credential.transaction_hash = f"0x{tx_hash}"
    credential.issued_at = datetime.utcnow()
    credential.metadata_uri = f"ipfs://QmMockHash{credential.id}"
    
    db.commit()
    db.refresh(credential)
    
    return {
        "success": True,
        "message": "Credential successfully minted on blockchain",
        "credential": credential,
        "explorer_url": f"https://mumbai.polygonscan.com/tx/{credential.transaction_hash}"
    }


@app.get("/api/v1/users/{user_id}/credentials", response_model=List[schemas.Credential])
def get_user_credentials(
    user_id: int,
    credential_type: Optional[str] = None,
    status: Optional[str] = None
):
    """Get all credentials for a user (their digital wallet)"""
    credentials = [c for c in CREDENTIALS_DB.values() if c.user_id == user_id]
    
    # Filter by type if provided
    if credential_type:
        try:
            type_enum = schemas.CredentialType(credential_type)
            credentials = [c for c in credentials if c.credential_type == type_enum]
        except ValueError:
            pass
    
    # Filter by status if provided
    if status:
        try:
            status_enum = schemas.CredentialStatus(status)
            credentials = [c for c in credentials if c.status == status_enum]
        except ValueError:
            pass
    
    # Sort by issue date (newest first)
    credentials.sort(key=lambda x: x.issued_at or x.created_at, reverse=True)
    
    return credentials


@app.get("/api/v1/credentials/{credential_id}", response_model=schemas.Credential)
def get_credential(credential_id: int):
    """Get a specific credential by ID"""
    credential = CREDENTIALS_DB.get(credential_id)
    if not credential:
        return schemas.CredentialVerification(
            is_valid=False,
            credential=None,
            verification_message="Credential not found. This token ID does not exist in our system.",
            verified_at=datetime.now()
        )
    
    if credential.status == schemas.CredentialStatus.REVOKED:
        return schemas.CredentialVerification(
            is_valid=False,
            credential=credential,
            verification_message="This credential has been revoked and is no longer valid.",
            verified_at=datetime.now()
        )
    
    if credential.expires_at and credential.expires_at < datetime.now():
        return schemas.CredentialVerification(
            is_valid=False,
            credential=credential,
            verification_message=f"This credential expired on {credential.expires_at.strftime('%Y-%m-%d')}. The skill may need to be refreshed.",
            verified_at=datetime.now()
        )
    
    if credential.status != schemas.CredentialStatus.ISSUED:
        return schemas.CredentialVerification(
            is_valid=False,
            credential=credential,
            verification_message="This credential is still being processed and has not been issued yet.",
            verified_at=datetime.now()
        )
    
    return schemas.CredentialVerification(
        is_valid=True,
        credential=credential,
        verification_message="✓ This credential is valid and verified on the blockchain.",
        verified_at=datetime.now()
    )


@app.post("/api/v1/credentials/{credential_id}/revoke")
def revoke_credential(credential_id: int, reason: str = ""):
    """Revoke a credential (admin/issuer only)"""
    credential = CREDENTIALS_DB.get(credential_id)
    if not credential:
        raise HTTPException(status_code=404, detail="Credential not found")
    
    if credential.status == schemas.CredentialStatus.REVOKED:
        raise HTTPException(status_code=400, detail="Credential already revoked")
    
    credential.status = schemas.CredentialStatus.REVOKED
    credential.updated_at = datetime.now()
    
    CREDENTIALS_DB[credential_id] = credential
    
    return {
        "success": True,
        "message": f"Credential revoked. Reason: {reason or 'Not specified'}",
        "credential": credential
    }


@app.post("/api/v1/projects/{project_id}/issue-credential")
def issue_project_credential(project_id: int, db: Session = Depends(get_db)):
    """Automatically issue credential when project is completed (Operate phase)"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if not project.operate_completed:
        # In stateless mode, frontend might have just updated it, but let's trust the DB
        raise HTTPException(
            status_code=400, 
            detail="Project must complete the Operate phase before credential can be issued"
        )
    
    # Check if credential already exists for this project
    existing_cred = db.query(models.Credential).filter(
        models.Credential.user_id == project.user_id,
        models.Credential.course_id == project.course_id
    ).first()

    if existing_cred:
        return {
            "success": True,
            "message": "Credential already issued",
            "credential": existing_cred
        }
    
    # Issue new credential
    now = datetime.now()
    
    # Get course details for title
    course = db.query(models.Course).filter(models.Course.id == project.course_id).first()
    course_title = course.title if course else "Unknown Course"

    new_credential = models.Credential(
        user_id=project.user_id,
        course_id=project.course_id,
        title=f"Certified: {course_title}",
        description=f"Successfully completed CDIO project: {project.project_title}",
        issuer_name="TSEA-X Institute",
        credential_type="Soulbound Token (SBT)",
        status="issued", # In real app, might be 'minting'
        issued_at=now
    )
    
    db.add(new_credential)
    db.commit()
    db.refresh(new_credential)
    
    # Auto-mint the credential (Mock for now, or call blockchain service)
    # mint_result = mint_credential(new_credential.id)
    
    return {
        "success": True,
        "message": "Project credential issued and minted successfully",
        "credential": new_credential
    }



# ===== CDIO Phase Endpoints (SQLAlchemy) =====

@app.post("/api/v1/charters", response_model=schemas.ProjectCharter, status_code=status.HTTP_201_CREATED)
def create_charter(charter_data: schemas.ProjectCharterCreate, project_id: int = Query(..., alias="project_id"), db: Session = Depends(get_db)):
    """Submit a project charter (Conceive phase)"""
    # Check if project exists
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Check if charter already exists
    if project.charter:
        raise HTTPException(status_code=400, detail="Charter already exists for this project")
    
    # Create charter
    new_charter = models.ProjectCharter(
        project_id=project_id,
        problem_statement=charter_data.problem_statement,
        success_metrics=charter_data.success_metrics,
        target_outcome=charter_data.target_outcome,
        constraints=charter_data.constraints,
        stakeholders=charter_data.stakeholders,
        suggested_tools=[], # AI can populate this later
        created_at=datetime.utcnow(),
        updated_at=datetime.utcnow()
    )
    
    db.add(new_charter)
    db.commit()
    db.refresh(new_charter)
    
    # Update project status
    project.current_phase = "design"
    # project.conceive_completed = True # Computed property
    project.overall_status = "in_progress"
    project.last_activity_at = datetime.utcnow()
    db.commit()
    
    # Auto-capture to PKC (Knowledge Node)
    # For MVP, we just create the node
    try:
        node_content = f"Problem: {new_charter.problem_statement}\nMetrics: {new_charter.success_metrics}"
        pkc_node = models.KnowledgeNode(
            user_id=project.user_id,
            title=f"Charter: {project.title}",
            content=node_content,
            node_type="project_artifact",
            source_id=str(new_charter.id),
            created_at=datetime.utcnow(),
            updated_at=datetime.utcnow()
        )
        db.add(pkc_node)
        db.commit()
    except Exception as e:
        print(f"Failed to auto-capture to PKC: {e}")

    return new_charter

@app.get("/api/v1/projects/{project_id}/charter", response_model=schemas.ProjectCharter)
def get_charter(project_id: int, db: Session = Depends(get_db)):
    """Get project charter"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
        
    if not project.charter:
        raise HTTPException(status_code=404, detail="Charter not found")
        
    return project.charter

@app.post("/api/v1/blueprints", response_model=schemas.DesignBlueprint, status_code=status.HTTP_201_CREATED)
def create_blueprint(blueprint_data: schemas.DesignBlueprintCreate, project_id: int = Query(..., alias="project_id"), db: Session = Depends(get_db)):
    """Submit a design blueprint (Design phase)"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Update existing blueprint or create new one
    if project.blueprint:
        # Update existing
        project.blueprint.architecture_diagram = blueprint_data.architecture_diagram
        project.blueprint.logic_flow = blueprint_data.logic_flow
        project.blueprint.component_list = blueprint_data.component_list
        project.blueprint.data_flow = blueprint_data.data_flow
        project.blueprint.updated_at = datetime.utcnow()
        db.commit()
        db.refresh(project.blueprint)
        blueprint = project.blueprint
    else:
        # Create new
        new_blueprint = models.DesignBlueprint(
            project_id=project_id,
            architecture_diagram=blueprint_data.architecture_diagram,
            logic_flow=blueprint_data.logic_flow,
            component_list=blueprint_data.component_list,
            data_flow=blueprint_data.data_flow,
            created_at=datetime.utcnow(),
            updated_at=datetime.utcnow()
        )
        db.add(new_blueprint)
        db.commit()
        db.refresh(new_blueprint)
        blueprint = new_blueprint
    
    # Only advance phase if it's the first save
    if project.current_phase == "design":
        project.current_phase = "implement"
    project.last_activity_at = datetime.utcnow()
    db.commit()
    
    return blueprint

@app.get("/api/v1/projects/{project_id}/blueprint", response_model=schemas.DesignBlueprint)
def get_blueprint(project_id: int, db: Session = Depends(get_db)):
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project or not project.blueprint:
        raise HTTPException(status_code=404, detail="Blueprint not found")
    return project.blueprint

@app.post("/api/v1/implementations", response_model=schemas.Implementation, status_code=status.HTTP_201_CREATED)
def create_implementation(impl_data: schemas.ImplementationCreate, project_id: int = Query(..., alias="project_id"), db: Session = Depends(get_db)):
    """Submit implementation (Implement phase)"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
        
    # Check if exists (optional, maybe allow updates? For now assume new)
    if project.implementation:
        # Update existing
        impl = project.implementation
        impl.code_repository_url = impl_data.code_repository_url
        impl.code_snapshot = impl_data.code_snapshot
        impl.notes = impl_data.notes
        impl.updated_at = datetime.utcnow()
    else:
        impl = models.Implementation(
            project_id=project_id,
            code_repository_url=impl_data.code_repository_url,
            code_snapshot=impl_data.code_snapshot,
            notes=impl_data.notes,
            created_at=datetime.utcnow(),
            updated_at=datetime.utcnow()
        )
        db.add(impl)
    
    db.commit()
    db.refresh(impl)
    
    project.current_phase = "operate"
    # project.implement_completed = True # Computed property
    project.last_activity_at = datetime.utcnow()
    db.commit()
    
    return impl

@app.get("/api/v1/projects/{project_id}/implementation", response_model=schemas.Implementation)
def get_implementation(project_id: int, db: Session = Depends(get_db)):
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project or not project.implementation:
        raise HTTPException(status_code=404, detail="Implementation not found")
    return project.implementation

@app.post("/api/v1/deployments", response_model=schemas.Deployment, status_code=status.HTTP_201_CREATED)
def create_deployment(deploy_data: schemas.DeploymentCreate, project_id: int = Query(..., alias="project_id"), db: Session = Depends(get_db)):
    """Submit deployment (Operate phase)"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
        
    if project.deployment:
        deploy = project.deployment
        deploy.deployment_url = deploy_data.deployment_url
        deploy.deployment_platform = deploy_data.deployment_platform
        deploy.readme = deploy_data.readme
        deploy.updated_at = datetime.utcnow()
    else:
        deploy = models.Deployment(
            project_id=project_id,
            deployment_url=deploy_data.deployment_url,
            deployment_platform=deploy_data.deployment_platform,
            readme=deploy_data.readme,
            created_at=datetime.utcnow(),
            updated_at=datetime.utcnow()
        )
        db.add(deploy)
        
    db.commit()
    db.refresh(deploy)
    
    # project.operate_completed = True # Computed property
    project.overall_status = "completed"
    project.completion_percentage = 100
    project.last_activity_at = datetime.utcnow()
    db.commit()
    
    return deploy

@app.get("/api/v1/projects/{project_id}/deployment", response_model=schemas.Deployment)
def get_deployment(project_id: int, db: Session = Depends(get_db)):
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project or not project.deployment:
        raise HTTPException(status_code=404, detail="Deployment not found")
    return project.deployment


# ===== Grading Agent Endpoints =====

from services.grading_service import grading_service

class CharterGradingRequest(BaseModel):
    project_id: int
    problem_statement: Optional[str] = None
    success_metrics: Optional[str] = None
    target_outcome: Optional[str] = None
    constraints: Optional[str] = None
    stakeholders: Optional[str] = None

class ImplementationGradingRequest(BaseModel):
    project_id: int
    code_snapshot: Optional[str] = None

@app.post("/api/v1/grading/charter")
async def grade_charter(request: CharterGradingRequest, db: Session = Depends(get_db)):
    """Grade a project charter using AI analysis"""
    
    # Check if data is provided in request (Stateless mode)
    if request.problem_statement and request.success_metrics:
        problem_statement = request.problem_statement
        success_metrics = request.success_metrics
        target_outcome = request.target_outcome
        constraints = request.constraints
        stakeholders = request.stakeholders
    else:
        # Fetch from DB
        project = db.query(models.CDIOProject).filter(models.CDIOProject.id == request.project_id).first()
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")
        
        if not project.charter:
            raise HTTPException(status_code=404, detail="Charter not found")
        
        charter = project.charter
            
        problem_statement = charter.problem_statement
        success_metrics = charter.success_metrics
        target_outcome = charter.target_outcome
        constraints = charter.constraints
        stakeholders = charter.stakeholders
    
    # Grade the charter
    grading_result = await grading_service.grade_charter(
        problem_statement=problem_statement,
        success_metrics=success_metrics,
        target_outcome=target_outcome,
        constraints=constraints,
        stakeholders=stakeholders
    )
    
    return grading_result


class BlueprintGradingRequest(BaseModel):
    project_id: int
    logic_flow: Optional[str] = None
    component_list: Optional[List[str]] = None
    data_flow: Optional[str] = None
    architecture_diagram: Optional[str] = None

@app.post("/api/v1/grading/blueprint")
async def grade_blueprint(request: BlueprintGradingRequest, db: Session = Depends(get_db)):
    """Grade a design blueprint using AI analysis"""
    
    # Check if data is provided in request (Stateless mode)
    if request.logic_flow or request.component_list:
        logic_flow = request.logic_flow
        component_list = request.component_list
        data_flow = request.data_flow
        architecture_diagram = request.architecture_diagram
    else:
        # Fetch from DB
        project = db.query(models.CDIOProject).filter(models.CDIOProject.id == request.project_id).first()
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")
        
        if not project.blueprint:
            raise HTTPException(status_code=404, detail="Blueprint not found")
        
        blueprint = project.blueprint
            
        logic_flow = blueprint.logic_flow
        component_list = blueprint.component_list
        data_flow = blueprint.data_flow
        architecture_diagram = blueprint.architecture_diagram
    
    # Grade the blueprint
    grading_result = await grading_service.grade_blueprint(
        logic_flow=logic_flow,
        component_list=component_list,
        data_flow=data_flow,
        architecture_diagram=architecture_diagram
    )
    
    return grading_result


@app.post("/api/v1/grading/implementation")
async def grade_implementation(request: ImplementationGradingRequest, db: Session = Depends(get_db)):
    """Grade an implementation submission using AI analysis"""
    
    code_snapshot = request.code_snapshot
    repository_url = None
    framework_used = "Unknown"
    dependencies = []

    if not code_snapshot:
        # Fetch from DB
        project = db.query(models.CDIOProject).filter(models.CDIOProject.id == request.project_id).first()
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")
        
        if not project.implementation:
            raise HTTPException(status_code=404, detail="Implementation not found")
        
        implementation = project.implementation
        
        code_snapshot = implementation.code_snapshot
        repository_url = implementation.code_repository_url
        framework_used = implementation.framework_used
        dependencies = implementation.dependencies
    
    # Grade the implementation
    grading_result = await grading_service.grade_implementation(
        code_snapshot=code_snapshot,
        repository_url=repository_url,
        framework_used=framework_used,
        dependencies=dependencies
    )
    
    return grading_result


# ===== Socratic Tutor Endpoint =====

class SocraticRequest(BaseModel):
    user_message: str
    project_id: int
    context: Optional[Dict] = {}
    history: Optional[List[Dict]] = []

@app.post("/api/v1/ai/socratic")
async def socratic_tutor(request: SocraticRequest):
    """Get a Socratic response from the AI tutor"""
    
    # Fetch project context if not provided
    project_context = request.context
    if not project_context and request.project_id:
        project = PROJECTS_DB.get(request.project_id)
        if project:
            charter = CHARTERS_DB.get(project.charter_id) if project.charter_id else None
            project_context = {
                "project_title": project.project_title,
                "phase": project.current_phase,
                "problem_statement": charter.problem_statement if charter else "Not defined",
                "success_metrics": charter.success_metrics if charter else "Not defined"
            }
    
    return await openai_service.socratic_response(
        request.user_message,
        project_context,
        request.history,
        request.project_id
    )

# ===== Code Execution Endpoints =====

@app.post("/api/v1/code/execute")
async def execute_code(request: schemas.CodeExecutionRequest):
    """Execute code safely and return output"""
    from services.code_execution_service import code_execution_service
    result = await code_execution_service.execute_code(request.code, request.language)
    return result


# ===== Admin Dashboard Endpoints =====

@app.get("/api/v1/admin/stats")
def get_admin_stats(db: Session = Depends(get_db)):
    """Get platform-wide statistics for admin dashboard"""
    # Calculate total users
    total_users = db.query(models.User).count()
    
    # Calculate enrollments
    total_enrollments = db.query(models.Enrollment).count()
    
    # Calculate credentials
    total_credentials = db.query(models.Credential).count()
   
    # Calculate active projects
    active_projects = db.query(models.CDIOProject).filter(models.CDIOProject.overall_status == "in_progress").count()
    
    # Calculate completion rate
    completed_projects = db.query(models.CDIOProject).filter(models.CDIOProject.overall_status == "completed").count()
    total_projects = db.query(models.CDIOProject).count()
    total_projects = total_projects if total_projects > 0 else 1
    completion_rate = int((completed_projects / total_projects) * 100)
    
    # Calculate total courses
    total_courses = db.query(models.Course).count()

    # Calculate revenue (mock for now or sum revenue transactions)
    revenue_this_month = db.query(func.sum(models.RevenueTransaction.amount)).scalar() or 0.0

    return {
        "totalUsers": total_users,
        "totalCourses": total_courses,
        "totalEnrollments": total_enrollments,
        "totalCredentials": total_credentials,
        "activeProjects": active_projects,
        "completionRate": completion_rate,
        "revenueThisMonth": revenue_this_month, 
        "systemHealth": "healthy"
    }


@app.get("/api/v1/admin/recent-activity")
def get_recent_activity(db: Session = Depends(get_db)):
    """Get recent platform activity for admin dashboard"""
    activities = []
    
    # Add enrollment activities
    recent_enrollments = db.query(models.Enrollment).order_by(models.Enrollment.enrolled_at.desc()).limit(5).all()
    
    for enrollment in recent_enrollments:
        course = db.query(models.Course).filter(models.Course.id == enrollment.course_id).first()
        if course:
            activities.append({
                "id": enrollment.id,
                "type": "user",
                "description": f"User {enrollment.user_id} enrolled in {course.title}",
                "timestamp": enrollment.enrolled_at.strftime("%Y-%m-%d %H:%M"),
                "status": "success"
            })
    
    # Add credential activities
    recent_credentials = db.query(models.Credential).order_by(models.Credential.issued_at.desc()).limit(3).all()
    
    for cred in recent_credentials:
        activities.append({
            "id": cred.id + 1000,
            "type": "credential",
            "description": f"Credential '{cred.title}' issued to user {cred.user_id}",
            "timestamp": cred.issued_at.strftime("%Y-%m-%d %H:%M") if cred.issued_at else "",
            "status": "success" if cred.status == "issued" else "pending"
        })
    
    # Sort by timestamp and limit to 10
    activities.sort(key=lambda x: x["timestamp"], reverse=True)
    return activities[:10]


@app.get("/api/v1/admin/user-analytics")
def get_user_analytics():
    """Get user growth analytics"""
    # Mock data - in real app, query database
    return {
        "userGrowth": [
            {"month": "Jul", "users": 850},
            {"month": "Aug", "users": 920},
            {"month": "Sep", "users": 1050},
            {"month": "Oct", "users": 1150},
            {"month": "Nov", "users": 1247}
        ],
        "usersByRole": {
            "students": 1180,
            "instructors": 45,
            "institutions": 22
        },
        "activeUsersToday": 342
    }


@app.get("/api/v1/admin/course-analytics")
def get_course_analytics():
    """Get course performance analytics"""
    # Calculate course stats
    course_stats = []
    for course in COURSES_DB:
        enrollments = [e for e in ENROLLMENTS_DB.values() if e.course_id == course["id"]]
        course_stats.append({
            "id": course["id"],
            "title": course["title"],
            "enrollments": len(enrollments),
            "rating": course["rating"],
            "category": course.get("category", "General")
        })
    
    return {
        "topCourses": sorted(course_stats, key=lambda x: x["enrollments"], reverse=True)[:5],
        "coursesByCategory": {
            "Sustainable Development": 12,
            "Public Policy": 8,
            "Technology": 15,
            "Business": 10
        }
    }



# ===== Code Execution Endpoints =====

class CodeExecutionRequest(BaseModel):
    code: str
    language: str

@app.post("/api/v1/code/execute")
async def execute_code(request: schemas.CodeExecutionRequest):
    """Execute code in a secure environment"""
    from services.code_execution_service import code_execution_service
    result = await code_execution_service.execute_code(request.code, request.language)
    return result


# ===== AI Socratic Tutor Endpoints =====

class SocraticChatRequest(BaseModel):
    """Request model for Socratic chat"""
    project_id: int
    user_message: str
    conversation_history: List[Dict] = []
    design_context: Dict = {}
    phase: str = "conceive"  # conceive, design, implement, operate

# ===== Instructor Dashboard Endpoints =====

@app.get("/api/v1/instructor/grading-queue")
def get_grading_queue(db: Session = Depends(get_db)):
    """Get all pending submissions for instructor review"""
    queue = []
    
    # 1. Check Charters (Conceive Phase)
    # Find projects in conceive phase that have a charter
    projects_conceive = db.query(models.CDIOProject).join(models.ProjectCharter).filter(
        models.CDIOProject.current_phase == "conceive"
    ).all()
    
    for p in projects_conceive:
        student = db.query(models.User).filter(models.User.id == p.user_id).first()
        course = db.query(models.Course).filter(models.Course.id == p.course_id).first()
        
        queue.append({
            "id": p.id,
            "project_id": p.id,
            "student_name": student.full_name if student else f"Student {p.user_id}",
            "project_title": p.title,
            "submission_type": "Charter",
            "course_title": course.title if course else "Unknown Course",
            "submitted_at": p.updated_at.strftime("%Y-%m-%d"),
            "status": "Pending"
        })

    # 2. Check Blueprints (Design Phase)
    projects_design = db.query(models.CDIOProject).join(models.DesignBlueprint).filter(
        models.CDIOProject.current_phase == "design"
    ).all()
    
    for p in projects_design:
        student = db.query(models.User).filter(models.User.id == p.user_id).first()
        course = db.query(models.Course).filter(models.Course.id == p.course_id).first()
        
        queue.append({
            "id": p.id + 1000,
            "project_id": p.id,
            "student_name": student.full_name if student else f"Student {p.user_id}",
            "project_title": p.title,
            "submission_type": "Blueprint",
            "course_title": course.title if course else "Unknown Course",
            "submitted_at": p.updated_at.strftime("%Y-%m-%d"),
            "status": "Pending"
        })

    # 3. Check Deployments (Operate Phase - Verification)
    projects_operate = db.query(models.CDIOProject).join(models.Deployment).filter(
        models.CDIOProject.current_phase == "operate",
        models.Deployment.verification_status == "submitted"
    ).all()
    
    for p in projects_operate:
        student = db.query(models.User).filter(models.User.id == p.user_id).first()
        course = db.query(models.Course).filter(models.Course.id == p.course_id).first()
        
        queue.append({
            "id": p.id + 3000,
            "project_id": p.id,
            "student_name": student.full_name if student else f"Student {p.user_id}",
            "project_title": p.title,
            "submission_type": "Deployment Verification",
            "course_title": course.title if course else "Unknown Course",
            "submitted_at": p.updated_at.strftime("%Y-%m-%d"),
            "status": "Pending"
        })

    return queue

@app.get("/api/v1/instructor/courses")
def get_instructor_courses(user_email: str = Query(...), db: Session = Depends(get_db)):
    """Get courses taught by the instructor"""
    user = db.query(models.User).filter(func.lower(models.User.email) == user_email.lower()).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Match by full_name as per current schema design
    courses = db.query(models.Course).filter(models.Course.instructor == user.full_name).all()
    return courses

@app.get("/api/v1/instructor/students")
def get_instructor_students(user_email: str = Query(...), db: Session = Depends(get_db)):
    """Get students enrolled in instructor's courses"""
    user = db.query(models.User).filter(func.lower(models.User.email) == user_email.lower()).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Get courses taught by instructor
    courses = db.query(models.Course).filter(models.Course.instructor == user.full_name).all()
    course_ids = [c.id for c in courses]
    
    if not course_ids:
        return []
    
    # Get enrollments for these courses
    enrollments = db.query(models.Enrollment).filter(models.Enrollment.course_id.in_(course_ids)).all()
    
    students = []
    for enrollment in enrollments:
        student = db.query(models.User).filter(models.User.id == enrollment.user_id).first()
        course = next((c for c in courses if c.id == enrollment.course_id), None)
        
        # Calculate progress (mock or fetch from project)
        # Try to find a project for this student in this course
        project = db.query(models.CDIOProject).filter(
            models.CDIOProject.user_id == enrollment.user_id,
            models.CDIOProject.course_id == enrollment.course_id
        ).first()
        
        progress = project.completion_percentage if project else 0
        grade = "N/A" # Could be calculated
        
        if student and course:
            students.append({
                "id": student.id,
                "name": student.full_name,
                "email": student.email,
                "course": course.title,
                "status": enrollment.status,
                "enrolled_at": enrollment.enrolled_at.strftime("%Y-%m-%d"),
                "progress": progress,
                "grade": grade
            })
            
    return students

@app.post("/api/v1/instructor/courses", response_model=schemas.Course)
def create_course(course: schemas.CourseCreate, db: Session = Depends(get_db)):
    """Create a new course"""
    db_course = models.Course(
        title=course.title,
        instructor=course.instructor,
        org=course.org,
        image=course.image,
        tag=course.tag,
        level=course.level,
        category=course.category,
        description=course.description,
        duration=course.duration,
        institution_id=course.institution_id,
        rating=0.0,
        students_count="0"
    )
    db.add(db_course)
    db.commit()
    db.refresh(db_course)
    return db_course

@app.put("/api/v1/instructor/courses/{course_id}", response_model=schemas.Course)
def update_course(course_id: int, course_update: schemas.CourseUpdate, db: Session = Depends(get_db)):
    """Update an existing course"""
    db_course = db.query(models.Course).filter(models.Course.id == course_id).first()
    if not db_course:
        raise HTTPException(status_code=404, detail="Course not found")
    
    update_data = course_update.model_dump(exclude_unset=True)
    for key, value in update_data.items():
        setattr(db_course, key, value)
    
    db.commit()
    db.refresh(db_course)
    return db_course

@app.get("/api/v1/instructor/{instructor_id}/analytics")
def get_instructor_analytics(instructor_id: int):
    """Get analytics for instructor dashboard"""
    return {
        "engagement_trend": [45, 52, 49, 60, 55, 68, 72],
        "avg_completion_rate": 78,
        "total_students": 1245,
        "active_courses": 3,
        "student_performance": {
            "excellent": 15,
            "good": 45,
            "average": 30,
            "needs_help": 10
        }
    }

class FeedbackRequest(BaseModel):
    grade: str # approve, reject
    feedback: str

@app.post("/api/v1/projects/{project_id}/feedback")
def submit_feedback(project_id: int, feedback: FeedbackRequest, db: Session = Depends(get_db)):
    """Submit feedback and grade for a project submission"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # In a real app, we would update the specific submission status (charter, blueprint, etc.)
    # For this mock, we'll just return success
    
    return {
        "status": "success",
        "message": f"Feedback submitted for project {project_id}",
        "grade": feedback.grade,
        "timestamp": datetime.now()
    }


@app.post("/api/v1/projects/{project_id}/grade")
async def grade_project(project_id: int, db: Session = Depends(get_db)):
    """Auto-grade a project submission using AI"""
    from services import openai_service
    
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
        
    submission_content = {}
    submission_type = "Unknown"
    
    submission_content = {}
    submission_type = "Unknown"
    
    phase_to_grade = project.current_phase
    
    # Fallback logic: If current phase has no content, try previous phase
    # This handles the case where submission auto-advanced the phase
    if phase_to_grade == schemas.CDIOPhase.DESIGN:
        blueprint = db.query(models.DesignBlueprint).filter(models.DesignBlueprint.project_id == project_id).first()
        if not blueprint:
            phase_to_grade = schemas.CDIOPhase.CONCEIVE
            
    elif phase_to_grade == schemas.CDIOPhase.IMPLEMENT:
        # Check for snapshot or implementation record
        snapshot = db.query(models.CodeSnapshot).filter(models.CodeSnapshot.project_id == project_id).first()
        if not snapshot:
            phase_to_grade = schemas.CDIOPhase.DESIGN
            
    elif phase_to_grade == schemas.CDIOPhase.OPERATE:
        deployment = db.query(models.Deployment).filter(models.Deployment.project_id == project_id).first()
        if not deployment:
            phase_to_grade = schemas.CDIOPhase.IMPLEMENT

    if phase_to_grade == schemas.CDIOPhase.CONCEIVE:
        submission_type = "Project Charter"
        charter = db.query(models.ProjectCharter).filter(models.ProjectCharter.project_id == project_id).first()
        if charter:
            submission_content = {
                "problem_statement": charter.problem_statement,
                "success_metrics": charter.success_metrics,
                "constraints": charter.constraints
            }
            
    elif phase_to_grade == schemas.CDIOPhase.DESIGN:
        submission_type = "Design Blueprint"
        blueprint = db.query(models.DesignBlueprint).filter(models.DesignBlueprint.project_id == project_id).first()
        if blueprint:
            submission_content = {
                "architecture": blueprint.architecture_diagram,
                "logic_flow": blueprint.logic_flow,
                "components": blueprint.component_list
            }
            
    elif phase_to_grade == schemas.CDIOPhase.IMPLEMENT:
        submission_type = "Implementation Code"
        # Get latest snapshot
        snapshot = db.query(models.CodeSnapshot).filter(
            models.CodeSnapshot.project_id == project_id
        ).order_by(models.CodeSnapshot.timestamp.desc()).first()
        
        if snapshot:
            submission_content = {
                "code": snapshot.code,
                "language": snapshot.language
            }
        # Also check Implementation record if snapshot missing
        elif project.implementation:
             submission_content = {
                "code": project.implementation.code_snapshot,
                "language": "python" # Default
            }
            
    elif phase_to_grade == schemas.CDIOPhase.OPERATE:
        submission_type = "Deployment"
        deployment = db.query(models.Deployment).filter(models.Deployment.project_id == project_id).first()
        if deployment:
            submission_content = {
                "url": deployment.deployment_url,
                "platform": deployment.deployment_platform,
                "readme": deployment.readme
            }
            
    if not submission_content:
        # Fallback for demo if DB is empty but we want to show the feature
        submission_content = {"note": "No submission data found in database. Grading based on metadata."}
    
    student_name = f"Student {project.user_id}" # Resolve real name if possible
    
    feedback = await openai_service.generate_grading_feedback(
        submission_type=submission_type,
        project_title=project.title,
        student_name=student_name,
        submission_content=submission_content
    )
    
    return feedback


# ===== AI Socratic Tutor Endpoints =====
async def socratic_chat(request: SocraticChatRequest):
    """AI Socratic Tutor endpoint with phase-specific guidance"""
    from services import openai_service
    
    try:
        response = await openai_service.socratic_response(
            user_message=request.user_message,
            design_context=request.design_context,
            conversation_history=request.conversation_history,
            project_id=request.project_id,
            phase=request.phase
        )
        return response
    except Exception as e:
        print(f"Socratic chat error: {e}")
        return {
            "ai_response": "I'm having trouble connecting right now. Can you rephrase your question?",
            "followup_questions": []
        }


# ===== PKC Endpoints =====

@app.post("/api/v1/pkc/nodes", response_model=schemas.KnowledgeNode)
async def create_knowledge_node(node_data: schemas.KnowledgeNodeCreate, db: Session = Depends(get_db)):
    """Create a new knowledge node with embedding"""
    # For MVP, assume user_id=1
    user_id = 1
    
    # Generate embedding
    text_to_embed = f"{node_data.title}\n{node_data.content or ''}"
    embedding = await generate_embedding(text_to_embed)
    
    new_node = models.KnowledgeNode(
        user_id=user_id,
        title=node_data.title,
        content=node_data.content,
        node_type=node_data.node_type,
        source_id=node_data.source_id,
        embedding=embedding,
        created_at=datetime.utcnow(),
        updated_at=datetime.utcnow()
    )
    
    db.add(new_node)
    db.commit()
    db.refresh(new_node)
    return new_node

@app.get("/api/v1/pkc/nodes", response_model=List[schemas.KnowledgeNode])
def get_knowledge_nodes(user_id: int = 1, db: Session = Depends(get_db)):
    """Get all knowledge nodes for a user"""
    return db.query(models.KnowledgeNode).filter(models.KnowledgeNode.user_id == user_id).all()

@app.get("/api/v1/pkc/nodes/{node_id}", response_model=schemas.KnowledgeNode)
def get_knowledge_node(node_id: int, db: Session = Depends(get_db)):
    """Get a specific knowledge node"""
    node = db.query(models.KnowledgeNode).filter(models.KnowledgeNode.id == node_id).first()
    if not node:
        raise HTTPException(status_code=404, detail="Node not found")
    return node

@app.post("/api/v1/pkc/journal", response_model=schemas.LearningJournal)
def create_journal_entry(entry_data: schemas.LearningJournalCreate, db: Session = Depends(get_db)):
    """Create a new learning journal entry"""
    # For MVP, assume user_id=1
    user_id = 1
    
    new_entry = models.LearningJournal(
        user_id=user_id,
        entry_type=entry_data.entry_type,
        content=entry_data.content,
        tags=entry_data.tags,
        created_at=datetime.utcnow()
    )
    
    db.add(new_entry)
    db.commit()
    db.refresh(new_entry)
    return new_entry

@app.get("/api/v1/pkc/journal", response_model=List[schemas.LearningJournal])
def get_journal_entries(user_id: int = 1, db: Session = Depends(get_db)):
    """Get all journal entries for a user"""
    return db.query(models.LearningJournal).filter(models.LearningJournal.user_id == user_id).order_by(models.LearningJournal.created_at.desc()).all()


if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)

# Note: Auth routes (/api/v1/auth/login and /api/v1/auth/sync-user) are defined earlier in the file

@app.get("/api/v1/users/email/{email}")
def get_user_by_email(email: str, db: Session = Depends(get_db)):
    user = db.query(models.User).filter(models.User.email == email).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    return user

# --- Admin Endpoints ---

@app.get("/api/v1/admin/stats")
def get_admin_stats(db: Session = Depends(get_db)):
    # Calculate real stats
    total_users = db.query(models.User).count()
    total_courses = db.query(models.Course).count()
    total_enrollments = db.query(models.Enrollment).count()
    total_credentials = db.query(models.Credential).count()
    active_projects = db.query(models.CDIOProject).filter(models.CDIOProject.overall_status == "in_progress").count()
    
    # Calculate completion rate (completed projects / total projects)
    total_projects = db.query(models.CDIOProject).count()
    completion_rate = int((db.query(models.CDIOProject).filter(models.CDIOProject.overall_status == "completed").count() / total_projects * 100)) if total_projects > 0 else 0

    return {
        "totalUsers": total_users,
        "totalCourses": total_courses,
        "totalEnrollments": total_enrollments,
        "totalCredentials": total_credentials,
        "activeProjects": active_projects,
        "completionRate": completion_rate,
        "revenueThisMonth": 45000, # Mock revenue
        "systemHealth": "healthy"
    }

@app.get("/api/v1/admin/recent-activity")
def get_admin_activity():
    # Mock activity log
    return [
        {"id": 1, "type": "user", "description": "New user registration: Sarah Jones", "timestamp": "2 mins ago", "status": "success"},
        {"id": 2, "type": "course", "description": "Course 'AI Ethics' published", "timestamp": "1 hour ago", "status": "success"},
        {"id": 3, "type": "credential", "description": "Credential issued to Alex Chen", "timestamp": "3 hours ago", "status": "success"},
        {"id": 4, "type": "system", "description": "Database backup completed", "timestamp": "5 hours ago", "status": "success"},
        {"id": 5, "type": "user", "description": "Failed login attempt (admin)", "timestamp": "1 day ago", "status": "warning"}
    ]

@app.get("/api/v1/admin/users")
def get_admin_users(db: Session = Depends(get_db)):
    """Get all users for admin dashboard"""
    users = db.query(models.User).order_by(models.User.id.desc()).all()
    return [
        {
            "id": u.id,
            "name": u.full_name or u.email.split('@')[0] if u.email else 'Unknown',
            "email": u.email,
            "role": u.role or 'student',
            "status": getattr(u, 'status', 'active') or 'active',
            "joinedDate": u.created_at.strftime('%Y-%m-%d') if u.created_at else 'N/A',
            "lastActive": "Recently",
            "coursesEnrolled": db.query(models.Enrollment).filter(models.Enrollment.user_id == u.id).count()
        }
        for u in users
    ]

class ImportUserRequest(BaseModel):
    name: str
    email: str
    role: str = "student"

@app.post("/api/v1/admin/import-user")
def import_user(user_data: ImportUserRequest, db: Session = Depends(get_db)):
    """Import a single user from CSV data"""
    # Check if user already exists
    existing = db.query(models.User).filter(models.User.email == user_data.email).first()
    if existing:
        return {"success": False, "message": "User already exists", "id": existing.id}
    
    new_user = models.User(
        email=user_data.email,
        full_name=user_data.name,
        role=user_data.role,
        created_at=datetime.now()
    )
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    
    return {"success": True, "message": "User imported", "id": new_user.id}

@app.get("/api/v1/admin/courses")
def get_admin_courses(db: Session = Depends(get_db)):
    """Get all courses for admin dashboard"""
    def parse_students_count(val):
        if not val:
            return 0
        if isinstance(val, int):
            return val
        try:
            s = str(val).replace(',', '').replace('+', '').replace('K', '000').replace('k', '000')
            return int(s) if s.isdigit() else 0
        except:
            return 0
    
    courses = db.query(models.Course).order_by(models.Course.id.asc()).all()
    return [
        {
            "id": c.id,
            "title": c.title,
            "instructor": c.instructor or 'Unknown',
            "institution": c.org or 'Unknown',
            "category": c.category or 'General',
            "status": getattr(c, 'status', 'published') or 'published',
            "enrollments": parse_students_count(c.students_count),
            "rating": c.rating or 0,
            "lastUpdated": getattr(c, 'updated_at', None).strftime('%Y-%m-%d') if getattr(c, 'updated_at', None) else 'N/A'
        }
        for c in courses
    ]

@app.get("/api/v1/admin/credentials")
def get_admin_credentials(db: Session = Depends(get_db)):
    """Get all credentials for admin dashboard"""
    credentials = db.query(models.Credential).order_by(models.Credential.id.desc()).all()
    return [
        {
            "id": f"CRT-{c.id}",
            "user": db.query(models.User).filter(models.User.id == c.user_id).first().full_name if db.query(models.User).filter(models.User.id == c.user_id).first() else 'Unknown',
            "title": c.title or 'Certificate',
            "type": c.credential_type or 'Professional Certificate',
            "issuedDate": c.issued_at.strftime('%Y-%m-%d') if c.issued_at else 'N/A',
            "status": c.status or 'issued',
            "blockchainHash": c.transaction_hash or '-'
        }
        for c in credentials
    ]

# ===== SYLLABUS MANAGEMENT ENDPOINTS =====

class SyllabusSectionCreate(BaseModel):
    order: int
    title: str
    cdio_phase: str
    week_number: Optional[int] = None
    topics: Optional[List[str]] = None
    activities: Optional[List[str]] = None
    assessment: Optional[str] = None
    duration_hours: float = 3.0

class SyllabusCreate(BaseModel):
    title: str
    overview: Optional[str] = None
    learning_outcomes: Optional[List[str]] = None
    assessment_strategy: Optional[dict] = None
    resources: Optional[List[str]] = None
    hexahelix_sectors: Optional[List[str]] = None
    duration_weeks: int = 4
    sections: Optional[List[SyllabusSectionCreate]] = None

@app.get("/api/v1/courses/{course_id}/syllabus")
def get_course_syllabus(course_id: int, db: Session = Depends(get_db)):
    """Get all syllabi for a course"""
    syllabi = db.query(models.Syllabus).filter(models.Syllabus.course_id == course_id).all()
    result = []
    for s in syllabi:
        sections = db.query(models.SyllabusSection).filter(
            models.SyllabusSection.syllabus_id == s.id
        ).order_by(models.SyllabusSection.order).all()
        
        result.append({
            "id": s.id,
            "course_id": s.course_id,
            "version": s.version,
            "title": s.title,
            "overview": s.overview,
            "learning_outcomes": s.learning_outcomes or [],
            "assessment_strategy": s.assessment_strategy or {},
            "resources": s.resources or [],
            "hexahelix_sectors": s.hexahelix_sectors or [],
            "duration_weeks": s.duration_weeks,
            "status": s.status,
            "created_at": s.created_at.isoformat() if s.created_at else None,
            "sections": [
                {
                    "id": sec.id,
                    "order": sec.order,
                    "title": sec.title,
                    "cdio_phase": sec.cdio_phase,
                    "week_number": sec.week_number,
                    "topics": sec.topics or [],
                    "activities": sec.activities or [],
                    "assessment": sec.assessment,
                    "duration_hours": sec.duration_hours
                }
                for sec in sections
            ]
        })
    return result

@app.post("/api/v1/courses/{course_id}/syllabus")
def create_course_syllabus(course_id: int, syllabus_data: SyllabusCreate, db: Session = Depends(get_db)):
    """Create or update syllabus for a course"""
    # Create syllabus
    new_syllabus = models.Syllabus(
        course_id=course_id,
        title=syllabus_data.title,
        overview=syllabus_data.overview,
        learning_outcomes=syllabus_data.learning_outcomes,
        assessment_strategy=syllabus_data.assessment_strategy,
        resources=syllabus_data.resources,
        hexahelix_sectors=syllabus_data.hexahelix_sectors,
        duration_weeks=syllabus_data.duration_weeks,
        status="draft"
    )
    db.add(new_syllabus)
    db.commit()
    db.refresh(new_syllabus)
    
    # Add sections if provided
    if syllabus_data.sections:
        for sec in syllabus_data.sections:
            new_section = models.SyllabusSection(
                syllabus_id=new_syllabus.id,
                order=sec.order,
                title=sec.title,
                cdio_phase=sec.cdio_phase,
                week_number=sec.week_number,
                topics=sec.topics,
                activities=sec.activities,
                assessment=sec.assessment,
                duration_hours=sec.duration_hours
            )
            db.add(new_section)
        db.commit()
    
    return {"id": new_syllabus.id, "message": "Syllabus created successfully"}

@app.post("/api/v1/syllabus/generate")
async def generate_syllabus(
    course_id: int = Form(...),
    topic: str = Form(...),
    description: str = Form(""),
    duration_weeks: int = Form(4),
    level: str = Form("Intermediate"),
    hexahelix_sectors: str = Form(""),  # Comma-separated
    files: List[UploadFile] = File(None)
):
    """Generate T6 syllabus with AI, optionally from uploaded materials"""
    
    # Process uploaded files
    material_content = ""
    if files:
        for file in files:
            try:
                content = await file.read()
                filename = file.filename.lower()
                
                # Handle different file types
                if filename.endswith('.txt') or filename.endswith('.md'):
                    material_content += f"\n\n--- {file.filename} ---\n{content.decode('utf-8', errors='ignore')}"
                elif filename.endswith('.pdf'):
                    # Basic PDF text extraction
                    try:
                        import io
                        from pypdf import PdfReader
                        reader = PdfReader(io.BytesIO(content))
                        pdf_text = ""
                        for page in reader.pages[:20]:  # First 20 pages
                            pdf_text += page.extract_text() or ""
                        material_content += f"\n\n--- {file.filename} ---\n{pdf_text}"
                    except Exception as e:
                        print(f"PDF extraction error: {e}")
                elif filename.endswith('.pptx'):
                    # PowerPoint extraction
                    try:
                        import io
                        from pptx import Presentation
                        prs = Presentation(io.BytesIO(content))
                        pptx_text = ""
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    pptx_text += shape.text + "\n"
                        material_content += f"\n\n--- {file.filename} ---\n{pptx_text}"
                    except Exception as e:
                        print(f"PPTX extraction error: {e}")
                elif filename.endswith(('.png', '.jpg', '.jpeg')):
                    # OCR for images
                    try:
                        import io
                        from PIL import Image
                        import pytesseract
                        img = Image.open(io.BytesIO(content))
                        ocr_text = pytesseract.image_to_string(img)
                        material_content += f"\n\n--- {file.filename} (OCR) ---\n{ocr_text}"
                    except Exception as e:
                        print(f"OCR error: {e}")
                else:
                    # Try basic text decode
                    material_content += f"\n\n--- {file.filename} ---\n{content.decode('utf-8', errors='ignore')}"
            except Exception as e:
                print(f"Error processing {file.filename}: {e}")
    
    # Parse hexahelix sectors
    sectors = [s.strip() for s in hexahelix_sectors.split(",") if s.strip()] if hexahelix_sectors else None
    
    # Generate syllabus using AI
    result = await openai_service.generate_t6_syllabus(
        course_title=topic,
        course_description=description,
        duration_weeks=duration_weeks,
        level=level,
        material_content=material_content,
        hexahelix_sectors=sectors
    )
    
    return result

@app.put("/api/v1/syllabus/{syllabus_id}/publish")
def publish_syllabus(syllabus_id: int, db: Session = Depends(get_db)):
    """Publish a syllabus (change status from draft to published)"""
    syllabus = db.query(models.Syllabus).filter(models.Syllabus.id == syllabus_id).first()
    if not syllabus:
        raise HTTPException(status_code=404, detail="Syllabus not found")
    
    syllabus.status = "published"
    syllabus.updated_at = datetime.now()
    db.commit()
    
    return {"message": "Syllabus published successfully"}


# ===== IRIS Cycle Endpoints (NUSA Framework) =====
# Maps to: Immersion → Reflection → Iteration → Scale

class ImmersionDataCreate(BaseModel):
    """Immersion phase - observing authentic problem context"""
    observation_notes: str
    stakeholder_interviews: Optional[List[str]] = []
    field_photos: Optional[List[str]] = []
    problem_context: str
    target_sfia_level: int = 3

class ReflectionDataCreate(BaseModel):
    """Reflection phase - Q/P analysis and SFIA gap mapping"""
    q_what_i_know: List[str]
    p_what_i_need: List[str]
    sfia_current_level: int
    sfia_target_level: int
    skill_gaps: List[str]
    learning_resources: Optional[List[str]] = []

class IterationCycleCreate(BaseModel):
    """Single Build-Measure-Learn iteration"""
    cycle_number: int
    hypothesis: str
    build_description: str
    measurement_data: Optional[str] = None
    learnings: Optional[str] = None
    pivot_or_persevere: Optional[str] = None
    next_steps: Optional[str] = None

class ScaleDataCreate(BaseModel):
    """Scale phase - institutional handoff and credential issuance"""
    deployment_url: Optional[str] = None
    handoff_documentation: str
    institution_partner: Optional[str] = None
    adoption_metrics: Optional[str] = None
    presentation_url: Optional[str] = None
    sfia_achieved_level: int
    ready_for_credential: bool = False


@app.post("/api/v1/iris/immersion")
async def create_immersion(
    data: ImmersionDataCreate,
    project_id: int = Query(...),
    db: Session = Depends(get_db)
):
    """Save Immersion phase data (NUSA Framework - Phase 1)"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    now = datetime.utcnow()
    
    # Store in ProjectCharter with adapted fields for backward compatibility
    charter = db.query(models.ProjectCharter).filter(models.ProjectCharter.project_id == project_id).first()
    
    if charter:
        charter.problem_statement = data.problem_context
        charter.constraints = data.observation_notes
        charter.stakeholders = data.stakeholder_interviews
        charter.updated_at = now
    else:
        charter = models.ProjectCharter(
            project_id=project_id,
            problem_statement=data.problem_context,
            constraints=data.observation_notes,
            stakeholders=data.stakeholder_interviews,
            success_metrics=f"SFIA Level {data.target_sfia_level}",
            target_outcome="Authentic problem immersion complete",
            created_at=now,
            updated_at=now
        )
        db.add(charter)
    
    # Update project phase (using IRIS naming)
    project.current_phase = "reflection"  # Move to next phase
    project.conceive_completed = True
    project.overall_status = "in_progress"
    project.completion_percentage = 25
    project.last_activity_at = now
    
    db.commit()
    
    return {
        "status": "success",
        "phase": "immersion",
        "next_phase": "reflection",
        "project_id": project_id,
        "completion_percentage": 25
    }


@app.post("/api/v1/iris/reflection")
async def create_reflection(
    data: ReflectionDataCreate,
    project_id: int = Query(...),
    db: Session = Depends(get_db)
):
    """Save Reflection phase data - Q/P analysis (NUSA Framework - Phase 2)"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    now = datetime.utcnow()
    
    # Store in DesignBlueprint with adapted fields
    blueprint = db.query(models.DesignBlueprint).filter(models.DesignBlueprint.project_id == project_id).first()
    
    logic_flow = f"""
Q (What I Know): {', '.join(data.q_what_i_know)}
P (What I Need): {', '.join(data.p_what_i_need)}
SFIA Gap: Level {data.sfia_current_level} -> {data.sfia_target_level}
Skill Gaps: {', '.join(data.skill_gaps)}
"""
    
    if blueprint:
        blueprint.logic_flow = logic_flow
        blueprint.component_list = data.learning_resources
        blueprint.updated_at = now
    else:
        blueprint = models.DesignBlueprint(
            project_id=project_id,
            logic_flow=logic_flow,
            component_list=data.learning_resources,
            architecture_diagram="Q/P Gap Analysis",
            created_at=now,
            updated_at=now
        )
        db.add(blueprint)
    
    project.current_phase = "iteration"
    project.design_completed = True
    project.completion_percentage = 50
    project.last_activity_at = now
    
    db.commit()
    
    return {
        "status": "success",
        "phase": "reflection",
        "next_phase": "iteration",
        "project_id": project_id,
        "sfia_gap": data.sfia_target_level - data.sfia_current_level,
        "completion_percentage": 50
    }


@app.post("/api/v1/iris/iteration")
async def create_iteration(
    data: IterationCycleCreate,
    project_id: int = Query(...),
    db: Session = Depends(get_db)
):
    """Save Build-Measure-Learn iteration (NUSA Framework - Phase 3)"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    now = datetime.utcnow()
    
    # Update implementation
    impl = db.query(models.Implementation).filter(models.Implementation.project_id == project_id).first()
    
    iteration_notes = f"""
Cycle #{data.cycle_number}
Hypothesis: {data.hypothesis}
Build: {data.build_description}
Measurement: {data.measurement_data or 'Pending'}
Learnings: {data.learnings or 'Pending'}
Decision: {data.pivot_or_persevere or 'Pending'}
Next Steps: {data.next_steps or 'TBD'}
"""
    
    if impl:
        impl.notes = (impl.notes or "") + "\n\n" + iteration_notes
        impl.updated_at = now
    else:
        impl = models.Implementation(
            project_id=project_id,
            notes=iteration_notes,
            created_at=now,
            updated_at=now
        )
        db.add(impl)
    
    # Mark implement completed after 3+ cycles
    if data.cycle_number >= 3:
        project.implement_completed = True
        project.current_phase = "scale"
        project.completion_percentage = 75
    else:
        project.current_phase = "iteration"
        project.completion_percentage = 50 + (data.cycle_number * 8)
    
    project.last_activity_at = now
    db.commit()
    
    return {
        "status": "success",
        "phase": "iteration",
        "cycle_number": data.cycle_number,
        "next_phase": "scale" if data.cycle_number >= 3 else "iteration",
        "project_id": project_id,
        "completion_percentage": project.completion_percentage
    }


@app.post("/api/v1/iris/scale")
async def create_scale(
    data: ScaleDataCreate,
    project_id: int = Query(...),
    db: Session = Depends(get_db)
):
    """Save Scale phase and issue credential (NUSA Framework - Phase 4)"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    now = datetime.utcnow()
    
    # Create deployment record
    deployment = db.query(models.Deployment).filter(models.Deployment.project_id == project_id).first()
    
    if deployment:
        deployment.deployment_url = data.deployment_url
        deployment.readme = data.handoff_documentation
        deployment.verification_status = "submitted" if data.ready_for_credential else "pending"
        deployment.updated_at = now
    else:
        deployment = models.Deployment(
            project_id=project_id,
            deployment_url=data.deployment_url,
            deployment_platform=data.institution_partner or "T6 Platform",
            readme=data.handoff_documentation,
            verification_status="submitted" if data.ready_for_credential else "pending",
            created_at=now,
            updated_at=now
        )
        db.add(deployment)
        db.flush()
    
    credential_id = None
    
    # Issue credential if ready
    if data.ready_for_credential:
        try:
            # Mint blockchain credential
            mint_result = await blockchain_service.mint_credential(
                user_id=project.user_id,
                course_id=project.course_id,
                project_id=project.id
            )
            
            deployment.sbt_minted = True
            deployment.sbt_token_id = mint_result.get("token_id")
            deployment.transaction_hash = mint_result.get("transaction_hash")
            
            credential_id = mint_result.get("credential_id")
        except Exception as e:
            print(f"Credential minting failed: {e}")
    
    project.operate_completed = True
    project.current_phase = "completed"
    project.overall_status = "completed"
    project.completion_percentage = 100
    project.last_activity_at = now
    
    db.commit()
    
    return {
        "status": "success",
        "phase": "scale",
        "project_id": project_id,
        "completion_percentage": 100,
        "credential_issued": data.ready_for_credential,
        "credential_id": credential_id,
        "sfia_achieved_level": data.sfia_achieved_level
    }


@app.get("/api/v1/iris/project/{project_id}/status")
def get_iris_status(project_id: int, db: Session = Depends(get_db)):
    """Get current IRIS phase status for a project"""
    project = db.query(models.CDIOProject).filter(models.CDIOProject.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Map CDIO phases to IRIS for backward compatibility
    phase_mapping = {
        'conceive': 'immerse',
        'design': 'realize',
        'implement': 'iterate',
        'operate': 'scale',
        'completed': 'completed',
        # IRIS phases pass through
        'immerse': 'immerse',
        'realize': 'realize',
        'iterate': 'iterate',
        'scale': 'scale',
        # Legacy IRIS names map to new
        'immersion': 'immerse',
        'reflection': 'realize',
        'iteration': 'iterate'
    }
    
    return {
        "project_id": project_id,
        "project_title": project.title,
        "current_phase": phase_mapping.get(project.current_phase, project.current_phase),
        "phases_completed": {
            "immerse": project.conceive_completed or False,
            "realize": project.design_completed or False,
            "iterate": project.implement_completed or False,
            "scale": project.operate_completed or False
        },
        "completion_percentage": project.completion_percentage,
        "overall_status": project.overall_status,
        "last_activity": project.last_activity_at.isoformat() if project.last_activity_at else None
    }


# ============================================
# GAMIFICATION ENDPOINTS
# ============================================

@app.get("/api/v1/users/{user_id}/gamification")
def get_user_gamification(user_id: int, db: Session = Depends(get_db)):
    """Get gamification data (streaks, XP, badges) for a specific user"""
    streak = db.query(models.LearningStreak).filter(models.LearningStreak.user_id == user_id).first()
    
    if not streak:
        # Return default values if no streak record exists
        return {
            "user_id": user_id,
            "current_streak": 0,
            "longest_streak": 0,
            "total_xp": 0,
            "week_activity": [False, False, False, False, False, False, False],
            "badges": [],
            "level": 1,
            "last_activity_date": None
        }
    
    return {
        "user_id": user_id,
        "current_streak": streak.current_streak,
        "longest_streak": streak.longest_streak,
        "total_xp": streak.total_xp,
        "week_activity": streak.week_activity or [False] * 7,
        "badges": streak.badges or [],
        "level": streak.level,
        "last_activity_date": streak.last_activity_date.isoformat() if streak.last_activity_date else None
    }


@app.get("/api/v1/admin/gamification-stats")
def get_gamification_stats(db: Session = Depends(get_db)):
    """Get platform-wide gamification statistics for admin dashboard"""
    
    # Get all streak records
    streaks = db.query(models.LearningStreak).all()
    
    if not streaks:
        return {
            "total_xp_distributed": 0,
            "active_streakers": 0,
            "average_streak": 0,
            "seven_day_achievers": 0,
            "thirty_day_achievers": 0,
            "total_badges_earned": 0,
            "top_learners": []
        }
    
    total_xp = sum(s.total_xp for s in streaks)
    active_streakers = sum(1 for s in streaks if s.current_streak > 0)
    avg_streak = sum(s.current_streak for s in streaks) / len(streaks) if streaks else 0
    seven_day_achievers = sum(1 for s in streaks if s.longest_streak >= 7)
    thirty_day_achievers = sum(1 for s in streaks if s.longest_streak >= 30)
    total_badges = sum(len(s.badges or []) for s in streaks)
    
    # Get top 5 learners by XP
    top_streaks = sorted(streaks, key=lambda x: x.total_xp, reverse=True)[:5]
    top_learners = []
    for s in top_streaks:
        user = db.query(models.User).filter(models.User.id == s.user_id).first()
        if user:
            top_learners.append({
                "user_id": s.user_id,
                "name": user.full_name,
                "xp": s.total_xp,
                "streak": s.current_streak,
                "level": s.level
            })
    
    return {
        "total_xp_distributed": total_xp,
        "active_streakers": active_streakers,
        "average_streak": round(avg_streak, 1),
        "seven_day_achievers": seven_day_achievers,
        "thirty_day_achievers": thirty_day_achievers,
        "total_badges_earned": total_badges,
        "top_learners": top_learners
    }


@app.get("/api/v1/admin/stats")
def get_admin_stats(db: Session = Depends(get_db)):
    """Get platform-wide stats for admin dashboard"""
    total_users = db.query(models.User).count()
    total_courses = db.query(models.Course).count()
    total_enrollments = db.query(models.Enrollment).count()
    total_credentials = db.query(models.Credential).count()
    active_projects = db.query(models.CDIOProject).filter(
        models.CDIOProject.overall_status == "in_progress"
    ).count()
    
    # Calculate completion rate
    completed_enrollments = db.query(models.Enrollment).filter(
        models.Enrollment.status == "completed"
    ).count()
    completion_rate = round((completed_enrollments / total_enrollments * 100) if total_enrollments > 0 else 0, 1)
    
    # Mock revenue for now
    revenue_this_month = 24500
    
    return {
        "totalUsers": total_users,
        "totalCourses": total_courses,
        "totalEnrollments": total_enrollments,
        "totalCredentials": total_credentials,
        "activeProjects": active_projects,
        "completionRate": completion_rate,
        "revenueThisMonth": revenue_this_month,
        "systemHealth": "healthy"
    }


@app.get("/api/v1/admin/users")
def get_admin_users(db: Session = Depends(get_db)):
    """Get all users for admin dashboard"""
    users = db.query(models.User).all()
    result = []
    for user in users:
        enrollments_count = db.query(models.Enrollment).filter(
            models.Enrollment.user_id == user.id
        ).count()
        result.append({
            "id": user.id,
            "name": user.full_name,
            "email": user.email,
            "role": user.role,
            "status": "active",
            "joinedDate": user.created_at.strftime("%Y-%m-%d") if user.created_at else "N/A",
            "lastActive": "Today",
            "coursesEnrolled": enrollments_count
        })
    return result


@app.get("/api/v1/admin/courses")
def get_admin_courses(db: Session = Depends(get_db)):
    """Get all courses for admin dashboard"""
    courses = db.query(models.Course).all()
    result = []
    for course in courses:
        inst = db.query(models.Institution).filter(models.Institution.id == course.institution_id).first()
        enrollments_count = db.query(models.Enrollment).filter(
            models.Enrollment.course_id == course.id
        ).count()
        result.append({
            "id": course.id,
            "title": course.title,
            "instructor": course.instructor,
            "institution": inst.name if inst else course.org,
            "category": course.category,
            "status": "published",
            "enrollments": enrollments_count,
            "rating": course.rating,
            "lastUpdated": "2025-12-01"
        })
    return result


@app.get("/api/v1/admin/credentials")
def get_admin_credentials(db: Session = Depends(get_db)):
    """Get all credentials for admin dashboard"""
    credentials = db.query(models.Credential).all()
    result = []
    for cred in credentials:
        user = db.query(models.User).filter(models.User.id == cred.user_id).first()
        result.append({
            "id": f"CRT-{cred.id}",
            "user": user.full_name if user else "Unknown",
            "title": cred.title,
            "type": cred.credential_type,
            "issuedDate": cred.issued_at.strftime("%Y-%m-%d") if cred.issued_at else "N/A",
            "status": cred.status,
            "blockchainHash": cred.transaction_hash or "Pending"
        })
    return result


@app.get("/api/v1/admin/recent-activity")
def get_recent_activity(db: Session = Depends(get_db)):
    """Get recent platform activity for admin dashboard"""
    activities = []
    
    # Get recent enrollments
    recent_enrollments = db.query(models.Enrollment).order_by(
        models.Enrollment.enrolled_at.desc()
    ).limit(3).all()
    
    for e in recent_enrollments:
        user = db.query(models.User).filter(models.User.id == e.user_id).first()
        course = db.query(models.Course).filter(models.Course.id == e.course_id).first()
        if user and course:
            activities.append({
                "id": e.id,
                "type": "enrollment",
                "description": f"{user.full_name} enrolled in {course.title}",
                "timestamp": e.enrolled_at.strftime("%Y-%m-%d %H:%M") if e.enrolled_at else "Recently",
                "status": "success"
            })
    
    # Get recent credentials
    recent_creds = db.query(models.Credential).order_by(
        models.Credential.created_at.desc()
    ).limit(2).all()
    
    for c in recent_creds:
        user = db.query(models.User).filter(models.User.id == c.user_id).first()
        if user:
            activities.append({
                "id": 100 + c.id,
                "type": "credential",
                "description": f"Credential issued to {user.full_name}: {c.title}",
                "timestamp": c.created_at.strftime("%Y-%m-%d %H:%M") if c.created_at else "Recently",
                "status": "success"
            })
    
    return activities[:5] if activities else [
        {"id": 1, "type": "system", "description": "System initialized", "timestamp": "Just now", "status": "success"}
    ]


@app.get("/api/v1/student/dashboard")
def get_student_dashboard(user_id: int, db: Session = Depends(get_db)):
    """Get student dashboard data including gamification"""
    
    user = db.query(models.User).filter(models.User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Get enrollments with course details
    enrollments = db.query(models.Enrollment).filter(
        models.Enrollment.user_id == user_id
    ).all()
    
    enrolled_courses = []
    for e in enrollments:
        course = db.query(models.Course).filter(models.Course.id == e.course_id).first()
        if course:
            # Calculate progress based on project completion
            project = db.query(models.CDIOProject).filter(
                models.CDIOProject.course_id == e.course_id,
                models.CDIOProject.user_id == user_id
            ).first()
            progress = project.completion_percentage if project else 0
            
            enrolled_courses.append({
                "course_id": course.id,
                "title": course.title,
                "instructor": course.instructor,
                "progress": progress,
                "status": e.status,
                "image": course.image,
                "category": course.category
            })
    
    # Get credentials
    credentials = db.query(models.Credential).filter(
        models.Credential.user_id == user_id
    ).all()
    creds_list = [{
        "id": c.id,
        "title": c.title,
        "type": c.credential_type,
        "status": c.status,
        "issued_at": c.issued_at.isoformat() if c.issued_at else None
    } for c in credentials]
    
    # Get gamification data
    streak = db.query(models.LearningStreak).filter(
        models.LearningStreak.user_id == user_id
    ).first()
    
    learning_streak = None
    if streak:
        learning_streak = {
            "current_streak": streak.current_streak,
            "longest_streak": streak.longest_streak,
            "this_week": streak.week_activity or [False] * 7,
            "total_xp": streak.total_xp
        }
    
    # Get IRIS projects
    projects = db.query(models.CDIOProject).filter(
        models.CDIOProject.user_id == user_id
    ).all()
    
    # Map CDIO to IRIS phases
    phase_mapping = {
        'conceive': 'immerse', 'design': 'realize', 'implement': 'iterate', 
        'operate': 'scale', 'immerse': 'immerse', 'realize': 'realize',
        'iterate': 'iterate', 'scale': 'scale', 'completed': 'completed',
        'immersion': 'immerse', 'reflection': 'realize', 'iteration': 'iterate'
    }
    
    iris_projects = []
    for p in projects:
        course = db.query(models.Course).filter(models.Course.id == p.course_id).first()
        iris_projects.append({
            "project_id": p.id,
            "course_id": p.course_id,
            "project_title": p.title or (f"{course.title} Project" if course else "Untitled"),
            "current_phase": phase_mapping.get(p.current_phase, p.current_phase),
            "completion_percentage": p.completion_percentage,
            "sfia_target_level": 3  # Default SFIA level
        })
    
    # Calculate stats
    total_hours = len(enrolled_courses) * 8  # Estimate 8 hours per course
    avg_progress = sum(c['progress'] for c in enrolled_courses) / len(enrolled_courses) if enrolled_courses else 0
    
    return {
        "enrolled_courses": enrolled_courses,
        "credentials": creds_list,
        "total_learning_hours": total_hours,
        "average_progress": round(avg_progress, 1),
        "recommended_courses": [],  # Can be populated with ML recommendations
        "upcoming_deadlines": [],
        "career_pathway_id": user.career_pathway_id,
        "iris_projects": iris_projects,
        "learning_streak": learning_streak
    }
