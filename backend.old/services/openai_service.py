"""
AI Service Layer for TSEA-X CDIO Framework
Supports: OpenAI (GPT-4), Google (Gemini), and Mock (Simulation) modes.
"""

import os
import json
import asyncio
import random
import base64
import logging
from typing import Any, Dict, List, Optional, Tuple
import io
import pypdf
from pptx import Presentation
from PIL import Image as PILImage
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

logger = logging.getLogger("openai_service")

clients: Dict[str, object] = {}
PRIORITY: List[str] = ["openai", "openrouter", "gemini", "mock"]
MODELS: Dict[str, str] = {
    "openai": os.getenv("OPENAI_MODEL", "gpt-3.5-turbo"),
    "openrouter": os.getenv("OPENROUTER_MODEL", "google/gemma-4-31b-it:free"),  # User's default from dashboard
    "gemini": "gemini-2.0-flash",
}
CLIENT_TYPES: Dict[str, str] = {}


def _init_openai_compatible(api_key: str, base_url: Optional[str] = None) -> object:
    from openai import AsyncOpenAI

    kwargs: Dict[str, Any] = {"api_key": api_key}
    if base_url:
        kwargs["base_url"] = base_url
    return AsyncOpenAI(**kwargs)


def _init_gemini(api_key: str):
    import google.generativeai as genai

    genai.configure(api_key=api_key)
    return genai


def _resolve_key(name: str, overrides: Dict[str, Optional[str]], db_keys: Optional[Dict[str, str]] = None) -> Optional[str]:
    if name in overrides:
        return overrides[name] or None
    if db_keys and name in db_keys:
        return db_keys[name]
    return os.getenv(name) or None


async def _create_chat_completion(client: Any, messages: list, model: str = "", temperature: float = 0.7, max_tokens: Optional[int] = None, **kwargs) -> Any:
    """Create chat completion, omitting model if empty (for OpenRouter dashboard default)."""
    params = {
        "messages": messages,
        "temperature": temperature,
        **kwargs
    }
    if model:
        params["model"] = model
    if max_tokens:
        params["max_tokens"] = max_tokens
    return await client.chat.completions.create(**params)


def _build_clients_from_env(overrides: Dict[str, Optional[str]]) -> Dict[str, object]:
    new: Dict[str, object] = {}

    openai_key = _resolve_key("OPENAI_API_KEY", overrides)
    openrouter_key = _resolve_key("OPENROUTER_API_KEY", overrides)
    gemini_key = _resolve_key("GEMINI_API_KEY", overrides)

    if openai_key:
        new["openai"] = _init_openai_compatible(openai_key)
    if openrouter_key:
        new["openrouter"] = _init_openai_compatible(openrouter_key, base_url="https://openrouter.ai/api/v1")
    if gemini_key:
        genai = _init_gemini(gemini_key)
        new["gemini"] = genai.GenerativeModel("gemini-2.0-flash")
        new["gemini_vision"] = genai.GenerativeModel("gemini-2.0-flash")

    return new


def _build_clients_from_db(
    db_session: Any,
    overrides: Dict[str, Optional[str]],
) -> Tuple[Dict[str, object], Dict[str, str], Dict[str, str], List[str]]:
    from routers.admin_settings import AdminSecret, AiProvider, _decrypt

    secret_rows = db_session.query(AdminSecret).all()
    db_keys: Dict[str, str] = {}
    for row in secret_rows:
        try:
            db_keys[row.key_name] = _decrypt(row.encrypted_value)
        except Exception:
            continue

    providers = (
        db_session.query(AiProvider)
        .filter(AiProvider.is_active == "true")
        .order_by(AiProvider.priority)
        .all()
    )

    new: Dict[str, object] = {}
    client_types: Dict[str, str] = {}
    models: Dict[str, str] = {}
    priority: List[str] = []

    for provider in providers:
        api_key = _resolve_key(provider.api_key_name, overrides, db_keys)
        if not api_key:
            continue

        if provider.provider_type == "gemini":
            genai = _init_gemini(api_key)
            new[provider.name] = genai.GenerativeModel(provider.model)
            client_types[provider.name] = "gemini"
            if provider.vision_capable == "true":
                new[f"{provider.name}_vision"] = genai.GenerativeModel(provider.model)
                client_types[f"{provider.name}_vision"] = "gemini"
        else:
            new[provider.name] = _init_openai_compatible(api_key, base_url=provider.base_url)
            client_types[provider.name] = "openai_compatible"

        models[provider.name] = provider.model
        priority.append(provider.name)

    return new, client_types, models, priority


def reload_clients(
    overrides: Optional[Dict[str, Optional[str]]] = None,
    db_session: Any = None,
) -> List[str]:
    overrides = overrides or {}

    new_clients: Dict[str, object] = {}
    new_client_types: Dict[str, str] = {}
    new_models: Dict[str, str] = {}
    new_priority: List[str] = []

    if db_session is not None:
        try:
            db_clients, db_types, db_models, db_priority = _build_clients_from_db(db_session, overrides)
            if db_clients:
                new_clients = db_clients
                new_client_types = db_types
                new_models = db_models
                new_priority = db_priority
        except Exception as e:
            logger.warning("Failed dynamic provider load; fallback to env: %s", e)

    if not new_clients:
        new_clients = _build_clients_from_env(overrides)
        if "openrouter" in new_clients:
            new_client_types["openrouter"] = "openai_compatible"
            new_models["openrouter"] = os.getenv("OPENROUTER_MODEL", "google/gemma-4-31b-it:free")  # User's dashboard default
            new_priority.append("openrouter")
        if "openai" in new_clients:
            new_client_types["openai"] = "openai_compatible"
            new_models["openai"] = os.getenv("OPENAI_MODEL", "gpt-3.5-turbo")
            new_priority.append("openai")
        if "gemini" in new_clients:
            new_client_types["gemini"] = "gemini"
            new_client_types["gemini_vision"] = "gemini"
            new_models["gemini"] = "gemini-2.0-flash"
            new_priority.append("gemini")

    if "mock" not in new_priority:
        new_priority.append("mock")

    clients.clear()
    clients.update(new_clients)

    CLIENT_TYPES.clear()
    CLIENT_TYPES.update(new_client_types)

    MODELS.clear()
    MODELS.update(new_models)

    PRIORITY[:] = new_priority
    logger.info("Reloaded providers: %s", list(clients.keys()))
    return list(clients.keys())

# Fallback models for OpenRouter if primary fails
OPENROUTER_FALLBACKS = [
    "google/gemini-2.0-flash-001",
    "meta-llama/llama-3.3-70b-instruct",
    "deepseek/deepseek-chat",
    "mistralai/mistral-small-24b-instruct-2501"
]


reload_clients()


print(f"[AI Service] Initialized providers: {list(clients.keys())}")

# Course metadata for context
COURSE_METADATA = {
    9: {
        "name": "AI Agents for Small Businesses",
        "level": "Intermediate",
        "tech_stack": ["Python", "OpenAI API", "LangChain", "FastAPI"],
        "domain": "AI/ML"
    },
    11: {
        "name": "Cybersecurity Fundamentals",
        "level": "Intermediate",
        "tech_stack": ["Security Tools", "Penetration Testing", "Encryption"],
        "domain": "Security"
    },
    1: {
        "name": "Sustainability and Green Technology",
        "level": "Beginner",
        "tech_stack": ["Frameworks", "Carbon Accounting", "Lifecycle Analysis"],
        "domain": "Sustainability"
    }
}

async def generate_charter_suggestions(
    problem_statement: str,
    success_metrics: str,
    course_id: int,
    target_outcome: Optional[str] = None
) -> Dict:
    """Generate intelligent charter suggestions"""
    
    course_info = COURSE_METADATA.get(course_id, {"name": "General", "level": "Intermediate", "domain": "General"})
    
    prompt = f"""You are an expert technical advisor for the course: "{course_info['name']}".
    
STUDENT PROJECT:
Problem: {problem_statement}
Metrics: {success_metrics}
Target: {target_outcome or 'Not specified'}

TASK:
Provide a JSON object with:
1. "suggested_tools": List of 3-5 tools suitable for {course_info['level']} level in {course_info['domain']}.
2. "estimated_duration": e.g. "2-3 weeks".
3. "difficulty_level": "Beginner", "Intermediate", or "Advanced".
4. "reasoning": A brief explanation.

Respond ONLY with valid JSON."""

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(1.5)
            return {
                "suggested_tools": ["Mock Tool A", "Mock Tool B", "Python"],
                "reasoning": f"Simulated response. Errors: {errors}",
                "estimated_duration": "2 weeks",
                "difficulty_level": "Intermediate"
            }

        if provider not in clients:
            continue

        try:
            content = ""
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                content = response.text
            else: # openai, openrouter
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[{"role": "system", "content": "You are a JSON generator."}, {"role": "user", "content": prompt}],
                    temperature=0.7
                )
                content = response.choices[0].message.content

            # Parse JSON
            try:
                content = content.replace("```json", "").replace("```", "").strip()
                return json.loads(content)
            except:
                import re
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    return json.loads(match.group())
                # If parsing fails, try next provider? Or just fail? 
                # Usually parsing fail is model issue, so try next provider is good.
                raise ValueError("Failed to parse JSON")

        except Exception as e:
            print(f"[{provider}] Charter AI Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {
        "suggested_tools": ["Error-Fallback Tool"],
        "reasoning": f"All AI Services Unavailable: {errors}",
        "estimated_duration": "Unknown",
        "difficulty_level": "Unknown"
    }

def raise_error():
    raise ValueError("Failed to parse JSON")

async def socratic_response(
    user_message: str,
    design_context: Dict,
    conversation_history: List[Dict],
    project_id: int,
    phase: str = "conceive"
) -> Dict:
    """Generate Socratic tutor response with phase-specific guidance"""
    
    # Phase-specific system prompts
    PHASE_PROMPTS = {
        "conceive": """You are a Socratic Tutor for the CONCEIVE phase. 
        Focus on helping students:
        - Clarify the problem statement and scope
        - Identify stakeholders and their needs
        - Define success metrics and constraints
        - Understand the business context
        
        NEVER give direct answers. ONLY ask guiding questions.
        Keep responses short (under 50 words). Be encouraging.
        Help them think critically about WHAT they're solving and WHY.""",
        
        "design": """You are a Socratic Tutor for the DESIGN phase.
        Focus on helping students:
        - Think through architecture decisions
        - Consider scalability and maintainability
        - Evaluate different design patterns
        - Plan data flow and component interactions
        
        NEVER give direct answers. ONLY ask guiding questions.
        Keep responses short (under 50 words). Be encouraging.
        Help them think critically about HOW to structure their solution.""",
        
        "implement": """You are a Socratic Tutor for the IMPLEMENT phase.
        Focus on helping students:
        - Debug code issues through reasoning
        - Apply best practices and clean code principles
        - Think about edge cases and error handling
        - Consider code quality and testing
        
        NEVER give direct answers. ONLY ask guiding questions.
        Keep responses short (under 50 words). Be encouraging.
        Help them think critically about code quality and implementation details.""",
        
        "operate": """You are a Socratic Tutor for the OPERATE phase.
        Focus on helping students:
        - Plan deployment strategies
        - Think about monitoring and logging
        - Consider user feedback and iteration
        - Understand maintenance and scalability
        
        NEVER give direct answers. ONLY ask guiding questions.
        Keep responses short (under 50 words). Be encouraging.
        Help them think critically about deployment, operations, and continuous improvement."""
    }
    
    # Get phase-specific prompt, default to conceive
    system_prompt = PHASE_PROMPTS.get(phase.lower(), PHASE_PROMPTS["conceive"])
    
    context_str = f"Design Context: {json.dumps(design_context)}"
    full_prompt = f"{system_prompt}\n\n{context_str}\n\nStudent: {user_message}"

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(1)
            phase_examples = {
                "conceive": "That's an interesting problem. Have you considered who the primary stakeholders are and what their specific needs might be?",
                "design": "Good thinking! How would this architecture handle increased load? What are the potential bottlenecks?",
                "implement": "I see you're working on that feature. What edge cases should you consider? How will you handle errors?",
                "operate": "Interesting deployment approach. How will you monitor the system's health? What metrics matter most?"
            }
            return {
                "ai_response": phase_examples.get(phase.lower(), phase_examples["conceive"]) + " (Simulated AI)",
                "followup_questions": []
            }

        if provider not in clients:
            continue

        try:
            if provider == "gemini":
                chat = clients[provider].start_chat(history=[])
                response = await chat.send_message_async(full_prompt)
                return {"ai_response": response.text, "followup_questions": []}
            else:
                messages = [{"role": "system", "content": system_prompt}]
                messages.append({"role": "user", "content": f"Context: {context_str}\nStudent: {user_message}"})
                
                response = await clients[provider].chat.completions.create(model=MODELS[provider], messages=messages, max_tokens=150)
                return {"ai_response": response.choices[0].message.content, "followup_questions": []}

        except Exception as e:
            print(f"[{provider}] Socratic AI Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {"ai_response": f"I'm pondering that... (All services failed: {errors})", "followup_questions": []}

async def companion_response(
    user_message: str,
    current_page: str,
    user_context: Dict,
    conversation_history: List[Dict]
) -> Dict:
    """Global AI Companion"""
    
    persona = "Guide"
    if "conceive" in current_page: persona = "Business Analyst"
    elif "design" in current_page: persona = "Architect"
    elif "implement" in current_page: persona = "Pair Programmer"
    elif "operate" in current_page: persona = "DevOps"
    
    prompt = f"""Role: {persona}. Context: User is on {current_page}.
    User says: "{user_message}".
    Reply helpfully and concisely (max 2 sentences)."""

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(0.5)
            return {
                "ai_response": f"[{persona} Mode] I see you're working on {current_page}. That's great! (Simulated Response)",
                "role_played": persona
            }

        if provider not in clients:
            continue

        try:
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                return {"ai_response": response.text, "role_played": persona}
            else:
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[{"role": "system", "content": f"You are a {persona}."}, {"role": "user", "content": user_message}],
                    max_tokens=100
                )
                return {"ai_response": response.choices[0].message.content, "role_played": persona}
                
        except Exception as e:
            print(f"[{provider}] AI Companion Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {"ai_response": f"I'm here to help, but my brain is offline momentarily. (Errors: {errors})", "role_played": "System"}

async def validate_design(blueprint: Dict) -> Dict:
    return {"validation_score": 0, "suggestions": []}

async def generate_course_structure(topic: str, target_audience: str = "Beginner", material_content: str = "") -> Dict:
    """
    Generate a CDIO/IRIS-based course structure with detailed module content.
    Optionally uses provided instructor materials to personalize content.
    """
    
    context_prompt = ""
    if material_content:
        context_prompt = f"""
        CONTEXT FROM INSTRUCTOR MATERIALS:
        The instructor has provided the following content. You MUST incorporate the key concepts, terminology, and themes from this material into the course structure.
        
        --- MATERIAL START ---
        {material_content[:12000]} 
        --- MATERIAL END ---
        """

    prompt = f"""You are Alexandria AI, an expert curriculum designer. Create a comprehensive course structure for "{topic}" targeting {target_audience} learners.
    
The course MUST follow both CDIO (Conceive, Design, Implement, Operate) and IRIS (Immerse, Realize, Iterate, Scale) frameworks.
{context_prompt}

OUTPUT FORMAT (valid JSON only):
{{
    "title": "Compelling course title that captures the essence of {topic}",
    "description": "2-3 sentence course description highlighting key learning outcomes and value proposition",
    "level": "{target_audience}",
    "duration": "4 weeks",
    "category": "Category based on topic",
    "image": "https://images.unsplash.com/photo-RELEVANT-IMAGE?w=800",
    "modules": [
        {{
            "phase": "Conceive",
            "iris_phase": "Immerse", 
            "title": "Module 1: Understanding & Context",
            "content": "Rich markdown text introduction to the module. Use bolding, lists, and clear structure.",
            "topics": ["Topic 1: Specific topic name", "Topic 2: Another topic", "Topic 3: Third topic"],
            "learning_goals": ["Students will be able to...", "Students will understand..."],
            "activities": [
                {{"type": "lecture", "title": "Introduction to {topic}", "duration_minutes": 30, "content": "Detailed markdown explanation of the topic. Include definitions, examples, and key takeaways."}},
                {{"type": "discussion", "title": "Real-world connections", "duration_minutes": 20, "prompt": "Thought-provoking question for students to discuss..."}},
                {{"type": "quiz", "title": "Key Concepts Check", "duration_minutes": 15, "questions": [{{"question": "Question text?", "options": ["A", "B", "C", "D"], "correctAnswer": "A", "explanation": "Why correct"}}]}},
                {{"type": "exercise", "title": "Hands-on Practice", "duration_minutes": 40, "instructions": "Step-by-step instructions for the activity..."}}
            ]
        }},
        ... (repeat for other phases)
    ],
    "capstone_project": "..."
}}

IMPORTANT: 
- Provide RICH, DETAILED content for every activity.
- Use Markdown formatting in descriptions and content fields.
- Include at least one QUIZ with 3 questions per module.
- Include specific discussion prompts.
- Make the 'content' field in activities substantial (at least 2-3 paragraphs).

Respond ONLY with valid JSON."""

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(2)
            # Generate intelligent mock based on topic
            words = topic.split()
            key_word = words[0] if words else "Course"
            
            return {
                "title": f"Mastering {topic}: From Concept to Scale",
                "description": f"A comprehensive {target_audience.lower()}-level course on {topic}. Learn to understand, design, implement, and deploy real-world solutions through hands-on projects and expert guidance.",
                "level": target_audience,
                "duration": "4 weeks",
                "category": topic.split()[0] if topic.split() else "General",
                "image": "https://images.unsplash.com/photo-1516321318423-f06f85e504b3?w=800",
                "modules": [
                    {
                        "phase": "Conceive", 
                        "iris_phase": "Immerse",
                        "title": f"Understanding {key_word} Fundamentals",
                        "content": f"Immerse yourself in the world of {topic}. Understand the landscape, key players, and fundamental concepts that drive innovation in this field. Build a solid foundation for practical application.",
                        "topics": [f"{key_word} Fundamentals", "Market & Landscape Analysis", "Core Concepts & Terminology", "Real-World Case Studies"],
                        "learning_goals": [f"Understand the core principles of {topic}", "Analyze the current market landscape", "Identify key stakeholders and use cases"],
                        "activities": [
                            {"type": "lecture", "title": f"Introduction to {topic}", "duration_minutes": 45},
                            {"type": "discussion", "title": "Industry trends and opportunities", "duration_minutes": 30},
                            {"type": "exercise", "title": "Stakeholder mapping exercise", "duration_minutes": 45}
                        ]
                    },
                    {
                        "phase": "Design",
                        "iris_phase": "Realize",
                        "title": f"Designing {key_word} Solutions",
                        "content": f"Transform your understanding into actionable plans. Learn to architect solutions, create prototypes, and develop implementation roadmaps for {topic} projects.",
                        "topics": ["Solution Architecture", "Prototype Development", "Technical Planning", "Resource Allocation"],
                        "learning_goals": ["Design scalable solution architectures", "Create functional prototypes", "Develop project implementation plans"],
                        "activities": [
                            {"type": "lecture", "title": "Architecture patterns and best practices", "duration_minutes": 40},
                            {"type": "exercise", "title": "Design your solution", "duration_minutes": 60},
                            {"type": "quiz", "title": "Design principles check", "duration_minutes": 15}
                        ]
                    },
                    {
                        "phase": "Implement",
                        "iris_phase": "Iterate",
                        "title": f"Building & Testing {key_word} Projects",
                        "content": f"Get hands-on with implementation. Build functional solutions, conduct rigorous testing, and iterate based on feedback. Apply industry best practices throughout development.",
                        "topics": ["Implementation Best Practices", "Quality Assurance & Testing", "Iteration & Refinement", "Performance Optimization"],
                        "learning_goals": ["Implement solutions using best practices", "Conduct comprehensive testing", "Iterate based on feedback"],
                        "activities": [
                            {"type": "lecture", "title": "Implementation strategies", "duration_minutes": 30},
                            {"type": "exercise", "title": "Build your solution", "duration_minutes": 90},
                            {"type": "demo", "title": "Peer review and feedback", "duration_minutes": 30}
                        ]
                    },
                    {
                        "phase": "Operate",
                        "iris_phase": "Scale",
                        "title": f"Deploying & Scaling {key_word} Solutions",
                        "content": f"Launch your solution to the real world. Learn deployment strategies, monitoring, and optimization techniques. Prepare for growth and sustainable operation.",
                        "topics": ["Deployment Strategies", "Monitoring & Analytics", "Scaling & Growth", "Sustainability Planning"],
                        "learning_goals": ["Deploy solutions effectively", "Set up monitoring and analytics", "Plan for scale and sustainability"],
                        "activities": [
                            {"type": "lecture", "title": "Deployment and DevOps", "duration_minutes": 35},
                            {"type": "exercise", "title": "Deploy your project", "duration_minutes": 50},
                            {"type": "reflect", "title": "Lessons learned & next steps", "duration_minutes": 25}
                        ]
                    }
                ],
                "capstone_project": f"Build and deploy a complete {topic} solution that addresses a real-world problem. Your project will demonstrate mastery of all four CDIO phases: from initial concept through to operational deployment. Final deliverables include a working prototype, technical documentation, and a presentation showcasing your solution's impact."
            }

        if provider not in clients:
            continue

        try:
            content = ""
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                content = response.text
            else:
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[
                        {"role": "system", "content": "You are Alexandria AI, an expert curriculum designer. Output only valid JSON."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    max_tokens=4000
                )
                content = response.choices[0].message.content

            # Parse JSON
            try:
                content = content.replace("```json", "").replace("```", "").strip()
                result = json.loads(content)
                # Ensure required fields
                if "modules" not in result:
                    result["modules"] = []
                return result
            except:
                import re
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    return json.loads(match.group())
                raise ValueError("Failed to parse JSON")

        except Exception as e:
            print(f"[{provider}] Course Gen Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {
        "title": f"Course on {topic}",
        "description": f"Content generation failed. Errors: {errors}",
        "modules": [],
        "capstone_project": "Error generating project."
    }

async def parse_course_material(
    file_base64: str,
    file_name: str,
    mime_type: str = "application/pdf"
) -> Dict:
    """
    Parse uploaded course material (PDF, PPTX, Image) and extract structured course content.
    Uses Gemini Vision for images/PDFs if available, falls back to text extraction + API.
    """
    
    # 1. Decode File
    try:
        file_bytes = base64.b64decode(file_base64)
    except Exception as e:
        return {"success": False, "error": f"Failed to decode file: {str(e)}"}
        
    extracted_text = ""
    is_vision_capable = False
    
    # 2. Extract Content based on type
    try:
        if mime_type == "application/pdf":
            # Attempt text extraction first (reliable for text-heavy PDFs)
            try:
                pdf_reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                text_parts = []
                # Limit to first 20 pages to avoid context limits
                for i in range(min(len(pdf_reader.pages), 20)):
                    text_parts.append(pdf_reader.pages[i].extract_text())
                extracted_text = "\n".join(text_parts)
            except Exception as e:
                print(f"PDF Text Extraction failed: {e}")
                
            # If text is sparse, mark for Vision fallback (if supported)
            if len(extracted_text) < 100: 
                is_vision_capable = True # For Gemini
                
        elif "presentation" in mime_type or file_name.endswith(".pptx"):
            try:
                prs = Presentation(io.BytesIO(file_bytes))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                extracted_text = "\n".join(text_parts)
            except Exception as e:
                print(f"PPTX Extraction failed: {e}")
                
        elif mime_type.startswith("image/"):
            is_vision_capable = True
            
        elif mime_type.startswith("text/"):
            extracted_text = file_bytes.decode("utf-8")
            
    except Exception as e:
        return {"success": False, "error": f"File processing failed: {str(e)}"}

    # 3. AI Analysis
    prompt = """
    Analyze this course material and extract a structured course outline.
    
    OUTPUT JSON format:
    {
        "title": "Proposed Course Title",
        "summary": "2-3 sentence summary of the material",
        "main_topics": [
            {
                "name": "Topic 1",
                "description": "Brief description",
                "subtopics": ["Sub 1", "Sub 2"]
            },
            ...
        ],
        "learning_objectives": ["Objective 1", "Objective 2"],
        "key_concepts": ["Concept 1", "Concept 2"],
        "difficulty_level": "Beginner/Intermediate/Advanced",
        "target_audience": "Who is this for?",
        "estimated_duration": "e.g. 4 weeks"
    }
    """
    
    errors = []
    
    # Priority: Gemini (Vision) -> OpenRouter (Text) -> OpenAI (Text)
    
    # A. Gemini Vision (Native PDF/Image support)
    if "gemini" in clients:
        try:
            # For Gemini 1.5/2.0, we can pass the image/pdf directly as Part
            # Convert to Part object
            import google.generativeai as genai
            
            content_parts = [prompt]
            
            if mime_type.startswith("image/"):
                img = PILImage.open(io.BytesIO(file_bytes))
                content_parts.append(img)
            elif mime_type == "application/pdf":
                 # Gemini PDF support usually requires upload, but we can try text fallback or 'pdf_part' if library supports
                 # For now, let's rely on the extracted text if we have it, OR just text.
                 # Actually, Gemini API supports inline data for images, but for PDF it usually wants 'File API'.
                 # Let's stick to text for PDFs unless we implement File API upload.
                 if not extracted_text:
                     return {"success": False, "error": "PDF Parsing failed and File API not implemented."}
                 content_parts.append(f"DOCUMENT CONTENT:\n{extracted_text[:30000]}") # 30k chars safety limit
            else:
                 content_parts.append(f"DOCUMENT CONTENT:\n{extracted_text[:30000]}")

            response = await clients["gemini"].generate_content_async(content_parts)
            content = response.text
            
            # Parse JSON
            content = content.replace("```json", "").replace("```", "").strip()
            # Find JSON brace
            import re
            match = re.search(r'\{.*\}', content, re.DOTALL)
            if match:
                data = json.loads(match.group())
                return {"success": True, "data": data}
                
        except Exception as e:
            errors.append(f"Gemini: {e}")
            
    # B. Text Fallback (OpenAI / OpenRouter)
    if not extracted_text:
         return {"success": False, "error": f"Could not extract text from file. Errors: {errors}"}
         
    for provider in ["openrouter", "openai"]:
        if provider not in clients: continue
        
        try:
            response = await clients[provider].chat.completions.create(
                model=MODELS[provider],
                messages=[
                    {"role": "system", "content": "You are an expert curriculum designer. Output valid JSON."},
                    {"role": "user", "content": f"{prompt}\n\nDOCUMENT CONTENT:\n{extracted_text[:15000]}"} # Limit tokens
                ],
                temperature=0.7
            )
            content = response.choices[0].message.content
            
            # Parse JSON
            content = content.replace("```json", "").replace("```", "").strip()
            import re
            match = re.search(r'\{.*\}', content, re.DOTALL)
            if match:
                data = json.loads(match.group())
                return {"success": True, "data": data}
                
        except Exception as e:
             errors.append(f"{provider}: {e}")
             
    return {"success": False, "error": f"All providers failed. Errors: {errors}"}


async def generate_t6_syllabus(
    course_title: str,
    course_description: str,
    duration_weeks: int,
    level: str,
    material_content: str = "",
    hexahelix_sectors: list = None
) -> Dict:
    """Generate T6-format syllabus (MIT OCW + TSEA-X CDIO Hybrid)"""
    
    # Calculate sections per CDIO phase based on duration
    sections_per_phase = max(1, duration_weeks // 4)
    total_sections = sections_per_phase * 4
    
    context_prompt = ""
    if material_content:
        context_prompt = f"""
        INSTRUCTOR MATERIALS (incorporate these into the syllabus):
        --- MATERIAL START ---
        {material_content[:15000]}
        --- MATERIAL END ---
        """
    
    sectors_str = ", ".join(hexahelix_sectors) if hexahelix_sectors else "Technology, Education"
    
    prompt = f"""You are an academic curriculum designer creating a T6 Syllabus (MIT OCW + TSEA-X CDIO Hybrid).

COURSE DETAILS:
- Title: {course_title}
- Description: {course_description}
- Duration: {duration_weeks} weeks
- Level: {level}
- Hexahelix Sectors: {sectors_str}

{context_prompt}

REQUIREMENTS:
1. Create EXACTLY {total_sections} sections ({sections_per_phase} per CDIO phase)
2. Each section = 1 week of content
3. Align sections to CDIO phases in order: Conceive -> Design -> Implement -> Operate
4. Include 3-5 specific topics per section
5. Include practical activities aligned to phase
6. Specify assessment types (Quiz, Assignment, Project, Presentation)

OUTPUT (valid JSON only):
{{
    "title": "Syllabus title",
    "overview": "Brief overview paragraph",
    "learning_outcomes": ["Outcome 1", "Outcome 2", ... (5-8 outcomes)],
    "assessment_strategy": {{
        "quizzes": 20,
        "assignments": 30,
        "project": 30,
        "capstone": 20
    }},
    "resources": ["Textbook 1", "Online resource", ...],
    "sections": [
        {{
            "order": 1,
            "title": "Section title",
            "cdio_phase": "conceive",
            "week_number": 1,
            "topics": ["Topic 1", "Topic 2", "Topic 3"],
            "activities": ["Activity 1", "Activity 2"],
            "assessment": "Quiz",
            "duration_hours": 3.0
        }},
        ... (repeat for all {total_sections} sections)
    ]
}}

Respond ONLY with valid JSON."""

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(2)
            # Generate mock sections
            mock_sections = []
            phases = ["conceive", "design", "implement", "operate"]
            phase_names = {
                "conceive": "Understanding & Analysis",
                "design": "Planning & Architecture", 
                "implement": "Building & Development",
                "operate": "Deployment & Iteration"
            }
            for i in range(total_sections):
                phase_idx = i // sections_per_phase
                phase = phases[min(phase_idx, 3)]
                mock_sections.append({
                    "order": i + 1,
                    "title": f"Week {i+1}: {phase_names[phase]}",
                    "cdio_phase": phase,
                    "week_number": i + 1,
                    "topics": [f"Topic {i+1}.1", f"Topic {i+1}.2", f"Topic {i+1}.3"],
                    "activities": [f"Activity for week {i+1}"],
                    "assessment": ["Quiz", "Assignment", "Project", "Presentation"][phase_idx % 4],
                    "duration_hours": 3.0
                })
            
            return {
                "title": f"T6 Syllabus: {course_title}",
                "overview": f"A {duration_weeks}-week comprehensive course on {course_title} using the CDIO framework.",
                "learning_outcomes": [
                    "Analyze and define complex problems",
                    "Design scalable solutions",
                    "Implement working prototypes",
                    "Deploy and iterate on solutions",
                    "Apply critical thinking skills"
                ],
                "assessment_strategy": {"quizzes": 20, "assignments": 30, "project": 30, "capstone": 20},
                "resources": ["Course textbook", "Online documentation", "Supplementary readings"],
                "sections": mock_sections
            }

        if provider not in clients:
            continue

        try:
            content = ""
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                content = response.text
            elif provider == "openrouter":
                # Try OpenRouter with fallbacks
                current_models = [MODELS[provider]] + [m for m in OPENROUTER_FALLBACKS if m != MODELS[provider]]
                last_error = None
                
                for model in current_models:
                    try:
                        print(f"[OpenRouter] Trying model: {model}")
                        response = await clients[provider].chat.completions.create(
                            model=model,
                            messages=[{"role": "system", "content": "You are a JSON curriculum designer."}, {"role": "user", "content": prompt}],
                            temperature=0.7
                        )
                        content = response.choices[0].message.content
                        break # Success
                    except Exception as e:
                        print(f"[OpenRouter] Model {model} failed: {e}")
                        last_error = e
                        continue
                
                if not content and last_error:
                    raise last_error
            else:
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[{"role": "system", "content": "You are a JSON curriculum designer."}, {"role": "user", "content": prompt}],
                    temperature=0.7
                )
                content = response.choices[0].message.content

            try:
                content = content.replace("```json", "").replace("```", "").strip()
                return json.loads(content)
            except:
                import re
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    return json.loads(match.group())
                raise ValueError("Failed to parse JSON")

        except Exception as e:
            print(f"[{provider}] Syllabus Gen Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {
        "title": f"Syllabus for {course_title}",
        "overview": f"Generation failed. Errors: {errors}",
        "learning_outcomes": [],
        "assessment_strategy": {},
        "resources": [],
        "sections": []
    }

async def test_ai_connection() -> Dict:
    """Test the AI connection and return a simple response"""
    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            return {"status": "success", "provider": "mock", "message": "This is a simulated response."}
        
        if provider not in clients:
            continue
            
        try:
            prompt = "Say 'Hello from AI!' and mention which model you are."
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                return {"status": "success", "provider": "gemini", "message": response.text}
            else:
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[{"role": "user", "content": prompt}]
                )
                return {"status": "success", "provider": "openai" if provider == "openai" else "openrouter", "message": response.choices[0].message.content}
        except Exception as e:
            errors.append(f"{provider}: {str(e)}")
            continue
            
    return {"status": "error", "provider": "none", "message": f"All providers failed: {errors}"}

async def generate_embedding(text: str) -> List[float]:
    """Generate vector embedding for text"""
    
    # Mock fallback for development without keys
    if "mock" in PRIORITY or not any(k in clients for k in ["openai", "openrouter"]):
        # Return random 1536-dim vector for testing
        return [random.uniform(-0.1, 0.1) for _ in range(1536)]

    for provider in PRIORITY:
        if provider not in clients or provider == "mock":
            continue
            
        try:
            if provider == "openai" or provider == "openrouter":
                response = await clients[provider].embeddings.create(
                    model="text-embedding-3-small", # or text-embedding-ada-002
                    input=text
                )
                return response.data[0].embedding
            elif provider == "gemini":
                # Gemini embedding implementation
                result = genai.embed_content(
                    model="models/embedding-001",
                    content=text,
                    task_type="retrieval_document",
                    title="Embedding of text"
                )
                return result['embedding']
                
        except Exception as e:
            print(f"[{provider}] Embedding Error: {e}")
            continue
            
    # Fallback if all fail
    print("All embedding providers failed, returning random vector")
    return [random.uniform(-0.1, 0.1) for _ in range(1536)]

async def generate_grading_feedback(
    submission_type: str,
    project_title: str,
    student_name: str,
    submission_content: Dict
) -> Dict:
    """Generate grading feedback and score for a student submission"""
    
    prompt = f"""You are an expert instructor grading a student project.
    
    PROJECT DETAILS:
    Title: {project_title}
    Student: {student_name}
    Type: {submission_type}
    
    SUBMISSION CONTENT:
    {json.dumps(submission_content, indent=2)}
    
    TASK:
    Evaluate the submission based on:
    1. Clarity and completeness
    2. Feasibility (for Charters/Designs) or Code Quality (for Implementation)
    3. Alignment with CDIO principles
    
    Provide a JSON object with:
    1. "grade": A letter grade (A, B, C, D, F)
    2. "score": A numeric score (0-100)
    3. "feedback": A paragraph of constructive feedback (max 100 words)
    4. "strengths": List of 2-3 strong points
    5. "weaknesses": List of 2-3 areas for improvement
    
    Respond ONLY with valid JSON."""

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(1.5)
            return {
                "grade": "B+",
                "score": 88,
                "feedback": f"Good effort on the {submission_type}. The core concepts are there, but some details could be refined. (Simulated AI Feedback)",
                "strengths": ["Clear problem definition", "Good structure"],
                "weaknesses": ["Needs more specific metrics", "Consider edge cases"]
            }

        if provider not in clients:
            continue

        try:
            content = ""
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                content = response.text
            else:
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[{"role": "system", "content": "You are a JSON grading assistant."}, {"role": "user", "content": prompt}],
                    temperature=0.3
                )
                content = response.choices[0].message.content

            # Parse JSON
            try:
                content = content.replace("```json", "").replace("```", "").strip()
                return json.loads(content)
            except:
                import re
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    return json.loads(match.group())
                raise ValueError("Failed to parse JSON")

        except Exception as e:
            print(f"[{provider}] Grading AI Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {
        "grade": "N/A",
        "score": 0,
        "feedback": f"AI Grading Unavailable. Errors: {errors}",
        "strengths": [],
        "weaknesses": []
    }


# ===== SMART COURSE CREATION (Alexandria AI Engine) =====

async def parse_materials_multimodal(
    file_base64: str,
    file_name: str,
    mime_type: str = "application/pdf"
) -> Dict:
    """
    Parse uploaded materials (slides, PDFs, images) using Gemini Vision.
    Extracts structured content for intelligent course creation.
    
    Powered by Alexandria AI Engine for intelligent content extraction.
    """
    
    prompt = """You are an expert curriculum designer analyzing educational materials.
    
Analyze this document/slide and extract the following in JSON format:

{
    "title": "Suggested course title based on content",
    "summary": "2-3 sentence overview of the material",
    "main_topics": [
        {"name": "Topic name", "description": "Brief description", "subtopics": ["subtopic1", "subtopic2"]}
    ],
    "learning_objectives": ["What students will learn - list 4-6 objectives"],
    "key_concepts": ["List of key terms/concepts that should be taught"],
    "visual_elements": ["Description of any diagrams, charts, or images and their educational purpose"],
    "suggested_structure": {
        "immerse": {"title": "Phase title", "focus": "What to cover in Immerse phase"},
        "realize": {"title": "Phase title", "focus": "What to cover in Realize phase"},
        "iterate": {"title": "Phase title", "focus": "What to cover in Iterate phase"},
        "scale": {"title": "Phase title", "focus": "What to cover in Scale phase"}
    },
    "difficulty_level": "Beginner/Intermediate/Advanced",
    "estimated_duration": "e.g., 4 weeks",
    "target_audience": "Who this content is best suited for",
    "prerequisites": ["Any required prior knowledge"]
}

Respond ONLY with valid JSON. Be thorough in extracting educational value."""

    errors = []
    
    # Try Gemini Vision first (best for multi-modal)
    if "gemini_vision" in clients:
        try:
            import google.generativeai as genai
            from google.generativeai.types import content_types
            
            print(f"[AI Service] Parsing with Gemini Vision: {file_name}")
            
            # Gemini expects inline data in specific format
            # Use Part with inline_data for multimodal content
            file_part = {
                "inline_data": {
                    "mime_type": mime_type,
                    "data": file_base64
                }
            }
            
            # Send both the prompt and file data
            response = await clients["gemini_vision"].generate_content_async([
                prompt,
                file_part
            ])
            
            content = response.text
            print(f"[AI Service] Gemini Vision response length: {len(content)} chars")
            
            # Parse JSON
            try:
                content = content.replace("```json", "").replace("```", "").strip()
                return {"success": True, "data": json.loads(content), "provider": "gemini_vision"}
            except:
                import re
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    return {"success": True, "data": json.loads(match.group()), "provider": "gemini_vision"}
                raise ValueError("Failed to parse JSON from Gemini Vision response")
                
        except Exception as e:
            print(f"[gemini_vision] Parse Error: {e}")
            import traceback
            traceback.print_exc()
            errors.append(f"gemini_vision: {str(e)}")
    
    # Fallback to OpenAI GPT-4 Vision if available
    if "openai" in clients and mime_type.startswith("image/"):
        try:
            print(f"[AI Service] Fallback to OpenAI Vision: {file_name}")
            
            response = await clients["openai"].chat.completions.create(
                model="gpt-4o",  # GPT-4o has vision capabilities
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{file_base64}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=4096
            )
            
            content = response.choices[0].message.content
            print(f"[AI Service] OpenAI Vision response length: {len(content)} chars")
            
            # Parse JSON
            try:
                content = content.replace("```json", "").replace("```", "").strip()
                return {"success": True, "data": json.loads(content), "provider": "openai_vision"}
            except:
                import re
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    return {"success": True, "data": json.loads(match.group()), "provider": "openai_vision"}
                raise ValueError("Failed to parse JSON from OpenAI Vision response")
                
        except Exception as e:
            print(f"[openai_vision] Parse Error: {e}")
            errors.append(f"openai_vision: {str(e)}")
    
    # Fallback to text-only parsing if vision fails
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(2)
            
            # Extract intelligent content from filename
            import re
            base_name = file_name.rsplit('.', 1)[0] if '.' in file_name else file_name
            # Clean up filename: replace hyphens/underscores with spaces, title case
            clean_title = re.sub(r'[-_]', ' ', base_name).strip()
            # Extract key themes from the filename
            words = clean_title.split()
            key_themes = [w.title() for w in words if len(w) > 3][:5]
            
            # Generate intelligent topics based on themes
            if len(key_themes) >= 2:
                main_topic = ' '.join(key_themes[:3])
                topics = [
                    {"name": f"Foundations of {key_themes[0]}", "description": f"Core principles and fundamentals of {key_themes[0].lower()}", "subtopics": ["Key Principles", "Historical Context", "Core Terminology"]},
                    {"name": f"{key_themes[1] if len(key_themes) > 1 else main_topic} Fundamentals", "description": f"Essential concepts and frameworks", "subtopics": ["Conceptual Framework", "Best Practices", "Industry Standards"]},
                    {"name": f"Applied {main_topic}", "description": f"Practical applications and real-world scenarios", "subtopics": ["Case Studies", "Implementation Strategies", "Tools & Techniques"]},
                    {"name": f"Advanced {key_themes[0]} Strategies", "description": f"Deep dive into advanced concepts and mastery", "subtopics": ["Advanced Methods", "Optimization", "Future Trends"]}
                ]
            else:
                main_topic = clean_title or "Course Content"
                topics = [
                    {"name": "Introduction & Overview", "description": "Getting started with foundational knowledge", "subtopics": ["Overview", "Key Terms", "Learning Goals"]},
                    {"name": "Core Concepts & Theories", "description": "Essential frameworks and theoretical foundations", "subtopics": ["Main Theory", "Supporting Concepts", "Practical Framework"]},
                    {"name": "Practical Applications", "description": "Hands-on application of concepts", "subtopics": ["Real-World Examples", "Implementation Guide", "Best Practices"]},
                    {"name": "Advanced Topics & Mastery", "description": "Deep expertise and leadership-level content", "subtopics": ["Expert Techniques", "Innovation", "Future Directions"]}
                ]
            
            return {
                "success": True,
                "provider": "mock",
                "data": {
                    "title": f"{main_topic}",
                    "summary": f"A comprehensive curriculum covering {main_topic.lower()}, designed to take participants from foundational understanding to practical mastery through structured learning experiences.",
                    "main_topics": topics,
                    "learning_objectives": [
                        f"Understand the fundamental principles of {main_topic.lower()}",
                        f"Apply theoretical knowledge to real-world {key_themes[0].lower() if key_themes else 'professional'} challenges",
                        "Develop critical analysis and problem-solving skills",
                        f"Create implementation strategies for {key_themes[1].lower() if len(key_themes) > 1 else 'organizational'} contexts",
                        "Evaluate and optimize solutions based on industry best practices"
                    ],
                    "key_concepts": key_themes if key_themes else ["Strategy", "Framework", "Implementation", "Analysis"],
                    "visual_elements": ["Framework diagrams", "Process flowcharts", "Case study illustrations"],
                    "suggested_structure": {
                        "immerse": {"title": "Understanding the Landscape", "focus": f"Context, stakeholder analysis, and {main_topic.lower()} foundations"},
                        "realize": {"title": "Designing Solutions", "focus": "Strategic planning, architecture, and roadmap development"},
                        "iterate": {"title": "Building & Testing", "focus": "Implementation, feedback loops, and continuous improvement"},
                        "scale": {"title": "Deployment & Impact", "focus": "Launch strategies, measurement, and sustainable growth"}
                    },
                    "difficulty_level": "Intermediate",
                    "estimated_duration": "4 weeks",
                    "target_audience": "Professionals and emerging leaders",
                    "prerequisites": [f"Basic understanding of {key_themes[0].lower() if key_themes else 'the domain'}"]
                }
            }
        
        if provider not in clients:
            continue
    
    return {
        "success": False,
        "error": f"All providers failed: {errors}",
        "data": None
    }


async def generate_teaching_agenda(
    parsed_content: Dict,
    target_audience: str = "Intermediate",
    duration_weeks: int = 4
) -> Dict:
    """
    Generate MIT OCW-Standard Enterprise Curriculum with Alexandria AI Engine.
    
    Produces international-standard output with:
    - Bloom's Taxonomy learning objectives
    - IRIS/CDIO framework alignment
    - SFIA competency mapping
    - Detailed session schedules
    - Assignment rubrics
    - Instructor support notes
    
    Teaching action types:
    - EXPLAIN: Lecture/presentation content
    - DISCUSS: Discussion prompts and case studies
    - PRACTICE: Hands-on labs and exercises
    - QUIZ: Formative knowledge checks
    - DEMO: Live demonstrations
    - REFLECT: Metacognitive reflection prompts
    - COLLABORATE: Group activities and peer learning
    """
    
    content_json = json.dumps(parsed_content, indent=2)
    
    # SFIA competency reference for AI/Tech courses
    sfia_reference = """
SFIA Competency Levels:
- Level 1 (Follow): Basic awareness, follows instructions
- Level 2 (Assist): Applies skills with guidance
- Level 3 (Apply): Applies skills independently
- Level 4 (Enable): Designs approaches, mentors others
- Level 5 (Advise): Strategic influence, organizational impact

Common SFIA Skills for Tech Courses:
- PROG (Programming/Software Development)
- DTAN (Data Analysis)
- MLAI (Machine Learning/AI)
- SCTY (Information Security)
- ARCH (Solution Architecture)
- TEST (Testing)
- DESN (Systems Design)
"""

    # Bloom's Taxonomy reference
    blooms_reference = """
Bloom's Taxonomy Verbs by Level:
1. REMEMBER: Define, List, Identify, Recall, Name, Recognize
2. UNDERSTAND: Explain, Describe, Summarize, Compare, Interpret, Classify
3. APPLY: Implement, Execute, Use, Demonstrate, Solve, Calculate
4. ANALYZE: Analyze, Differentiate, Examine, Distinguish, Compare, Contrast
5. EVALUATE: Assess, Critique, Justify, Recommend, Judge, Validate
6. CREATE: Design, Develop, Construct, Produce, Build, Formulate

Level Guidelines:
- Beginner courses: Focus on Remember, Understand, Apply (Levels 1-3)
- Intermediate courses: Focus on Apply, Analyze, Evaluate (Levels 3-5)
- Advanced courses: Focus on Analyze, Evaluate, Create (Levels 4-6)
"""

    prompt = f"""You are Alexandria AI, an expert MIT-trained curriculum designer creating enterprise-grade course syllabi.

═══════════════════════════════════════════════════════════════════════════════
MISSION: Create an MIT OpenCourseWare-standard curriculum that rivals the best
         university courses in the world. This is NOT a template - be SPECIFIC.
═══════════════════════════════════════════════════════════════════════════════

EXTRACTED COURSE CONTENT:
{content_json}

TARGET AUDIENCE: {target_audience} level students
DURATION: {duration_weeks} weeks

{blooms_reference}

{sfia_reference}

═══════════════════════════════════════════════════════════════════════════════
REQUIREMENTS (MIT OCW STANDARD):
═══════════════════════════════════════════════════════════════════════════════

1. LEARNING OBJECTIVES must use Bloom's Taxonomy verbs appropriately:
   - Each module needs 3-5 specific, measurable objectives
   - Match verb complexity to the target audience level
   - Include the assessment method for each objective

2. SESSION SCHEDULE must be detailed:
   - Include timing in minutes for each activity
   - Use all 7 teaching action types across the course
   - Balance theory (40%) with practice (60%)

3. READINGS must be specific:
   - Required readings with chapter/page references
   - Optional deep-dive materials
   - Online resources and tools

4. ASSIGNMENTS must include rubrics:
   - Clear deliverables
   - Weighted grading criteria
   - Due dates by week

5. SFIA COMPETENCY MAPPING:
   - Map course to 2-4 relevant SFIA skills
   - Specify target level (1-5) for each
   - Include evidence requirements

6. INSTRUCTOR NOTES:
   - Common student misconceptions
   - Engagement strategies
   - Differentiation for struggling/advanced students

═══════════════════════════════════════════════════════════════════════════════
OUTPUT JSON SCHEMA (MIT OCW ENTERPRISE STANDARD):
═══════════════════════════════════════════════════════════════════════════════

{{
    "course_code": "T6-XXX-101",
    "course_title": "Specific course title from content",
    "tagline": "Compelling one-line description that sells the course",
    "description": "2-3 paragraph detailed course description",
    "level": "{target_audience}",
    "credits": 3,
    "duration_weeks": {duration_weeks},
    "prerequisites": ["Prerequisite 1", "Prerequisite 2"],
    
    "program_learning_outcomes": [
        {{
            "id": "LO1",
            "objective": "By course completion, students will be able to [Bloom's verb] [specific skill]",
            "blooms_level": "apply",
            "blooms_verb": "implement",
            "assessment_method": "Capstone Project",
            "sfia_skills": ["PROG", "DTAN"]
        }},
        {{
            "id": "LO2",
            "objective": "Students will be able to [Bloom's verb] [specific skill]",
            "blooms_level": "analyze",
            "blooms_verb": "evaluate",
            "assessment_method": "Case Study Analysis",
            "sfia_skills": ["ARCH"]
        }}
    ],
    
    "sfia_competency_map": [
        {{
            "skill_code": "PROG",
            "skill_name": "Programming/Software Development",
            "target_level": 3,
            "level_name": "Apply",
            "evidence_required": ["Working code submissions", "Technical documentation"],
            "modules_addressed": [1, 2, 3]
        }}
    ],
    
    "modules": [
        {{
            "week": 1,
            "iris_phase": "immerse",
            "cdio_phase": "conceive",
            "title": "Week 1: [Specific Module Title]",
            "subtitle": "Understanding [specific topic] context",
            "duration_hours": 3,
            
            "learning_objectives": [
                {{
                    "id": "M1-LO1",
                    "text": "[Bloom's verb] the key concepts of [topic]",
                    "blooms_level": "understand",
                    "blooms_verb": "explain"
                }},
                {{
                    "id": "M1-LO2", 
                    "text": "[Bloom's verb] [specific skill] to [context]",
                    "blooms_level": "apply",
                    "blooms_verb": "demonstrate"
                }}
            ],
            
            "session_schedule": [
                {{
                    "order": 1,
                    "type": "EXPLAIN",
                    "title": "Lecture: [Specific Topic]",
                    "description": "Detailed content coverage description",
                    "duration_minutes": 45,
                    "key_points": ["Point 1", "Point 2", "Point 3"],
                    "materials": ["Slides", "Whiteboard"]
                }},
                {{
                    "order": 2,
                    "type": "DISCUSS",
                    "title": "Case Study: [Real-world example]",
                    "description": "Analysis of [specific case]",
                    "duration_minutes": 20,
                    "discussion_prompts": ["Question 1?", "Question 2?"],
                    "facilitation_notes": "Encourage diverse perspectives"
                }},
                {{
                    "order": 3,
                    "type": "DEMO",
                    "title": "Live Demo: [Specific tool/technique]",
                    "description": "Instructor demonstrates [specific skill]",
                    "duration_minutes": 25,
                    "tools_used": ["Tool 1", "Tool 2"],
                    "key_takeaways": ["Takeaway 1", "Takeaway 2"]
                }},
                {{
                    "order": 4,
                    "type": "PRACTICE",
                    "title": "Hands-on Lab: [Specific exercise]",
                    "description": "Students practice [specific skill]",
                    "duration_minutes": 50,
                    "instructions": ["Step 1", "Step 2", "Step 3"],
                    "expected_output": "What students should produce",
                    "common_pitfalls": ["Pitfall 1", "Pitfall 2"]
                }},
                {{
                    "order": 5,
                    "type": "QUIZ",
                    "title": "Knowledge Check",
                    "description": "Formative assessment of module concepts",
                    "duration_minutes": 15,
                    "question_count": 5,
                    "passing_threshold": 70
                }},
                {{
                    "order": 6,
                    "type": "REFLECT",
                    "title": "Weekly Reflection",
                    "description": "Metacognitive synthesis",
                    "duration_minutes": 10,
                    "reflection_prompts": [
                        "What concept was most challenging and why?",
                        "How does this connect to your prior knowledge?"
                    ]
                }}
            ],
            
            "readings": {{
                "required": [
                    {{
                        "title": "Book/Article Title",
                        "author": "Author Name",
                        "chapters": "Chapters 1-2",
                        "pages": "pp. 1-45",
                        "url": "https://example.com if online"
                    }}
                ],
                "optional": [
                    {{
                        "title": "Deep Dive Resource",
                        "description": "For advanced learners",
                        "url": "https://example.com"
                    }}
                ]
            }},
            
            "assignment": {{
                "title": "Assignment 1: [Specific Title]",
                "description": "Detailed description of what students will do",
                "learning_outcomes_addressed": ["LO1", "LO2"],
                "deliverables": ["Deliverable 1", "Deliverable 2"],
                "due_week": 2,
                "weight_percent": 15,
                "rubric": [
                    {{
                        "criterion": "Technical Accuracy",
                        "weight": 40,
                        "excellent": "All requirements met with exceptional quality",
                        "proficient": "All requirements met",
                        "developing": "Most requirements met",
                        "beginning": "Few requirements met"
                    }},
                    {{
                        "criterion": "Documentation",
                        "weight": 30,
                        "excellent": "Clear, comprehensive documentation",
                        "proficient": "Adequate documentation",
                        "developing": "Partial documentation",
                        "beginning": "Minimal documentation"
                    }},
                    {{
                        "criterion": "Problem-Solving",
                        "weight": 30,
                        "excellent": "Creative, efficient solutions",
                        "proficient": "Functional solutions",
                        "developing": "Basic solutions with issues",
                        "beginning": "Incomplete solutions"
                    }}
                ]
            }},
            
            "resources": {{
                "software_tools": ["Tool 1", "Tool 2"],
                "online_resources": ["Resource 1 URL", "Resource 2 URL"],
                "lab_materials": ["Lab guide", "Dataset/VM"]
            }}
        }}
    ],
    
    "capstone_project": {{
        "title": "Capstone: [Specific Project Title]",
        "description": "Comprehensive description of the final project",
        "learning_outcomes_addressed": ["LO1", "LO2", "LO3", "LO4"],
        "sfia_skills_demonstrated": ["PROG", "ARCH", "DTAN"],
        "milestones": [
            {{
                "name": "Project Proposal",
                "due_week": 1,
                "weight_percent": 10,
                "deliverables": ["Problem statement", "Proposed solution", "Timeline"]
            }},
            {{
                "name": "Design Document",
                "due_week": 2,
                "weight_percent": 20,
                "deliverables": ["Architecture diagram", "Technical specifications"]
            }},
            {{
                "name": "Working Prototype",
                "due_week": 3,
                "weight_percent": 40,
                "deliverables": ["Functional code", "Test results", "Documentation"]
            }},
            {{
                "name": "Final Presentation",
                "due_week": {duration_weeks},
                "weight_percent": 30,
                "deliverables": ["Live demo", "Presentation slides", "Reflection report"]
            }}
        ],
        "rubric": {{
            "technical_excellence": 35,
            "innovation_creativity": 20,
            "presentation_communication": 20,
            "documentation": 15,
            "teamwork_collaboration": 10
        }}
    }},
    
    "assessment_strategy": {{
        "formative": [
            "Weekly knowledge check quizzes (not graded)",
            "In-class polling and discussions",
            "Peer code reviews"
        ],
        "summative": {{
            "assignments": 35,
            "quizzes": 15,
            "capstone_project": 40,
            "participation": 10
        }},
        "grading_scale": {{
            "A": "90-100",
            "B": "80-89",
            "C": "70-79",
            "D": "60-69",
            "F": "0-59"
        }}
    }},
    
    "instructor_notes": {{
        "common_misconceptions": [
            {{
                "concept": "[Specific concept]",
                "misconception": "Students often think...",
                "correction_strategy": "Address by..."
            }}
        ],
        "engagement_hooks": [
            "Use [specific current event] to introduce [topic]",
            "Start with [surprising fact] to capture attention"
        ],
        "differentiation": {{
            "struggling_learners": [
                "Provide step-by-step lab guides",
                "Pair with peer mentors",
                "Offer office hours"
            ],
            "advanced_learners": [
                "Assign extension challenges",
                "Independent research project option",
                "Peer teaching opportunities"
            ]
        }},
        "technology_requirements": [
            "Laptop with [specs]",
            "Software: [list]",
            "Internet access for [resources]"
        ]
    }},
    
    "auto_generated_quizzes": [
        {{
            "module": 1,
            "questions": [
                {{
                    "question": "Specific question based on module content?",
                    "options": ["A) Option A", "B) Option B", "C) Option C", "D) Option D"],
                    "correct": 0,
                    "explanation": "Detailed explanation of why this is correct",
                    "blooms_level": "understand"
                }}
            ]
        }}
    ],
    
    "_metadata": {{
        "framework": "IRIS + CDIO",
        "standard": "MIT OCW Enterprise",
        "blooms_taxonomy": true,
        "sfia_aligned": true,
        "version": "2.0"
    }}
}}

═══════════════════════════════════════════════════════════════════════════════
CRITICAL REQUIREMENTS:
═══════════════════════════════════════════════════════════════════════════════

1. Generate EXACTLY {duration_weeks} modules (one per week)
2. Distribute IRIS phases evenly: immerse → realize → iterate → scale
3. Each module must have 5-7 teaching actions with specific content
4. Include at least 2 required readings per module with real references
5. Each assignment must have a complete rubric with 3-4 criteria
6. Map at least 2 SFIA competencies to the course
7. Include 3+ specific instructor notes per category
8. Generate 3-5 quiz questions per module

BE SPECIFIC. Reference actual tools, frameworks, and real-world examples.
Use the extracted content to create DOMAIN-SPECIFIC curriculum, not generic templates.

Respond ONLY with valid JSON."""

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(2)
            
            # Extract title from parsed content
            course_title = parsed_content.get("title", "Enterprise Skills Course")
            main_topics = parsed_content.get("main_topics", [])
            
            # Generate enterprise-grade mock modules
            modules = []
            phases = ["immerse", "realize", "iterate", "scale"]
            cdio_phases = ["conceive", "design", "implement", "operate"]
            phase_titles = {
                "immerse": "Understanding & Exploration",
                "realize": "Planning & Architecture",
                "iterate": "Building & Testing",
                "scale": "Deployment & Optimization"
            }
            phase_subtitles = {
                "immerse": "Deep dive into context and stakeholder needs",
                "realize": "Designing solutions and technical architecture",
                "iterate": "Hands-on development and iterative refinement",
                "scale": "Production deployment and continuous improvement"
            }
            
            blooms_by_phase = {
                "immerse": [("understand", "explain"), ("remember", "identify"), ("understand", "describe")],
                "realize": [("apply", "design"), ("analyze", "differentiate"), ("apply", "demonstrate")],
                "iterate": [("apply", "implement"), ("analyze", "debug"), ("evaluate", "test")],
                "scale": [("evaluate", "assess"), ("create", "optimize"), ("apply", "deploy")]
            }
            
            for week in range(1, duration_weeks + 1):
                phase_idx = (week - 1) % 4
                phase = phases[phase_idx]
                cdio = cdio_phases[phase_idx]
                
                # Get topic from parsed content if available
                topic_name = main_topics[week-1]["name"] if week <= len(main_topics) else f"Module {week}"
                topic_desc = main_topics[week-1].get("description", "") if week <= len(main_topics) else ""
                
                # Generate learning objectives with proper Bloom's verbs
                blooms_verbs = blooms_by_phase.get(phase, [("apply", "implement")])
                learning_objectives = []
                for i, (level, verb) in enumerate(blooms_verbs[:3]):
                    learning_objectives.append({
                        "id": f"M{week}-LO{i+1}",
                        "text": f"{verb.capitalize()} the key concepts of {topic_name}" if i == 0 else f"{verb.capitalize()} {topic_name} principles in practical contexts",
                        "blooms_level": level,
                        "blooms_verb": verb
                    })
                
                # Generate session schedule
                session_schedule = [
                    {
                        "order": 1,
                        "type": "EXPLAIN",
                        "title": f"Lecture: Introduction to {topic_name}",
                        "description": f"Comprehensive overview of {topic_name} fundamentals and industry applications",
                        "duration_minutes": 45,
                        "key_points": [
                            f"Core concepts of {topic_name}",
                            "Industry best practices and standards",
                            "Real-world applications and case studies"
                        ],
                        "materials": ["Slides", "Reference documentation", "Industry examples"]
                    },
                    {
                        "order": 2,
                        "type": "DISCUSS",
                        "title": f"Case Study: {topic_name} in Practice",
                        "description": f"Analysis of real-world {topic_name} implementation",
                        "duration_minutes": 20,
                        "discussion_prompts": [
                            f"What challenges did the organization face with {topic_name}?",
                            "How would you approach this problem differently?"
                        ],
                        "facilitation_notes": "Encourage diverse perspectives and connect to student experiences"
                    },
                    {
                        "order": 3,
                        "type": "DEMO",
                        "title": f"Live Demo: {topic_name} Implementation",
                        "description": f"Instructor demonstrates {topic_name} techniques",
                        "duration_minutes": 25,
                        "tools_used": ["Development environment", "Framework tools", "Testing utilities"],
                        "key_takeaways": ["Step-by-step approach", "Common patterns", "Best practices"]
                    },
                    {
                        "order": 4,
                        "type": "PRACTICE",
                        "title": f"Hands-on Lab: {topic_name} Exercise",
                        "description": f"Students apply {topic_name} concepts in guided practice",
                        "duration_minutes": 50,
                        "instructions": [
                            "Set up your development environment",
                            f"Implement the core {topic_name} functionality",
                            "Test and validate your solution"
                        ],
                        "expected_output": f"Working {topic_name} implementation with documentation",
                        "common_pitfalls": ["Configuration errors", "Missing dependencies", "Scope creep"]
                    },
                    {
                        "order": 5,
                        "type": "QUIZ",
                        "title": "Knowledge Check",
                        "description": f"Formative assessment of {topic_name} concepts",
                        "duration_minutes": 15,
                        "question_count": 5,
                        "passing_threshold": 70
                    },
                    {
                        "order": 6,
                        "type": "REFLECT",
                        "title": "Weekly Reflection",
                        "description": "Metacognitive synthesis of learning",
                        "duration_minutes": 10,
                        "reflection_prompts": [
                            f"What was most challenging about {topic_name} and why?",
                            "How does this connect to your prior knowledge?",
                            "What questions do you still have?"
                        ]
                    }
                ]
                
                # Generate assignment with rubric
                assignment = {
                    "title": f"Assignment {week}: {topic_name} Implementation",
                    "description": f"Apply {topic_name} concepts to create a practical solution for a real-world scenario",
                    "learning_outcomes_addressed": [f"LO{(week-1) % 4 + 1}", f"LO{(week) % 4 + 1}"],
                    "deliverables": [
                        f"{topic_name} design document",
                        "Implementation code/artifacts",
                        "Testing documentation"
                    ],
                    "due_week": week + 1 if week < duration_weeks else week,
                    "weight_percent": round(35 / duration_weeks),
                    "rubric": [
                        {
                            "criterion": "Technical Accuracy",
                            "weight": 40,
                            "excellent": "All requirements met with exceptional quality and innovation",
                            "proficient": "All requirements met with good quality",
                            "developing": "Most requirements met with acceptable quality",
                            "beginning": "Few requirements met, significant gaps"
                        },
                        {
                            "criterion": "Documentation Quality",
                            "weight": 30,
                            "excellent": "Clear, comprehensive, professional documentation",
                            "proficient": "Adequate documentation covering key points",
                            "developing": "Partial documentation with gaps",
                            "beginning": "Minimal or missing documentation"
                        },
                        {
                            "criterion": "Problem-Solving Approach",
                            "weight": 30,
                            "excellent": "Creative, efficient solutions with clear reasoning",
                            "proficient": "Functional solutions with good approach",
                            "developing": "Basic solutions with some issues",
                            "beginning": "Incomplete or non-functional solutions"
                        }
                    ]
                }
                
                # Generate readings
                readings = {
                    "required": [
                        {
                            "title": f"Enterprise {topic_name} - Official Documentation",
                            "author": "Industry Consortium",
                            "chapters": "Chapters 1-3",
                            "pages": "pp. 1-45",
                            "url": "https://docs.example.com"
                        },
                        {
                            "title": f"Best Practices in {topic_name}",
                            "author": "Professional Association",
                            "chapters": "Chapter 2",
                            "pages": "pp. 20-35",
                            "url": "https://resources.example.com"
                        }
                    ],
                    "optional": [
                        {
                            "title": f"Advanced {topic_name} Techniques",
                            "description": "For students who want to go deeper",
                            "url": "https://advanced.example.com"
                        }
                    ]
                }
                
                modules.append({
                    "week": week,
                    "iris_phase": phase,
                    "cdio_phase": cdio,
                    "title": f"Week {week}: {topic_name}",
                    "subtitle": phase_subtitles[phase],
                    "duration_hours": 3,
                    "learning_objectives": learning_objectives,
                    "session_schedule": session_schedule,
                    "readings": readings,
                    "assignment": assignment,
                    "resources": {
                        "software_tools": ["VS Code", "Git", "Docker"],
                        "online_resources": ["Official documentation", "Video tutorials"],
                        "lab_materials": ["Lab environment access", "Sample datasets"]
                    }
                })
            
            # Generate enterprise-grade response
            return {
                "success": True,
                "provider": "mock",
                "agenda": {
                    "course_code": f"T6-{course_title[:3].upper()}-101",
                    "course_title": course_title,
                    "tagline": f"Master {course_title.lower()} through hands-on, project-based learning",
                    "description": f"This comprehensive {duration_weeks}-week course provides in-depth training in {course_title}. Students will gain practical skills through hands-on labs, real-world case studies, and a capstone project that demonstrates mastery.",
                    "level": target_audience,
                    "credits": 3,
                    "duration_weeks": duration_weeks,
                    "prerequisites": ["Basic computer literacy", "Familiarity with foundational concepts"],
                    
                    "program_learning_outcomes": [
                        {
                            "id": "LO1",
                            "objective": f"Explain the fundamental concepts and principles of {course_title}",
                            "blooms_level": "understand",
                            "blooms_verb": "explain",
                            "assessment_method": "Quizzes and assignments",
                            "sfia_skills": ["KNOW", "LEDA"]
                        },
                        {
                            "id": "LO2",
                            "objective": f"Apply {course_title} methodologies to solve real-world problems",
                            "blooms_level": "apply",
                            "blooms_verb": "implement",
                            "assessment_method": "Hands-on labs and assignments",
                            "sfia_skills": ["PROG", "DTAN"]
                        },
                        {
                            "id": "LO3",
                            "objective": f"Analyze and evaluate {course_title} solutions for effectiveness",
                            "blooms_level": "evaluate",
                            "blooms_verb": "assess",
                            "assessment_method": "Case studies and peer reviews",
                            "sfia_skills": ["ARCH", "TEST"]
                        },
                        {
                            "id": "LO4",
                            "objective": f"Design and create innovative {course_title} solutions",
                            "blooms_level": "create",
                            "blooms_verb": "design",
                            "assessment_method": "Capstone project",
                            "sfia_skills": ["DESN", "INOV"]
                        }
                    ],
                    
                    "sfia_competency_map": [
                        {
                            "skill_code": "PROG",
                            "skill_name": "Programming/Software Development",
                            "target_level": 3,
                            "level_name": "Apply",
                            "evidence_required": ["Working code submissions", "Technical documentation"],
                            "modules_addressed": list(range(1, duration_weeks + 1))
                        },
                        {
                            "skill_code": "DTAN",
                            "skill_name": "Data Analytics",
                            "target_level": 2,
                            "level_name": "Assist",
                            "evidence_required": ["Data analysis reports", "Visualization outputs"],
                            "modules_addressed": [2, 3, duration_weeks]
                        }
                    ],
                    
                    "modules": modules,
                    
                    "capstone_project": {
                        "title": f"Capstone: {course_title} Real-World Application",
                        "description": f"Apply all course concepts to build a comprehensive {course_title.lower()} solution for a real-world scenario. Students will work through the complete IRIS lifecycle.",
                        "learning_outcomes_addressed": ["LO1", "LO2", "LO3", "LO4"],
                        "sfia_skills_demonstrated": ["PROG", "ARCH", "DTAN"],
                        "milestones": [
                            {
                                "name": "Project Proposal",
                                "due_week": 1,
                                "weight_percent": 10,
                                "deliverables": ["Problem statement", "Proposed solution", "Timeline"]
                            },
                            {
                                "name": "Design Document",
                                "due_week": 2,
                                "weight_percent": 20,
                                "deliverables": ["Architecture diagram", "Technical specifications"]
                            },
                            {
                                "name": "Working Prototype",
                                "due_week": max(3, duration_weeks - 1),
                                "weight_percent": 40,
                                "deliverables": ["Functional code", "Test results", "Documentation"]
                            },
                            {
                                "name": "Final Presentation",
                                "due_week": duration_weeks,
                                "weight_percent": 30,
                                "deliverables": ["Live demo", "Presentation slides", "Reflection report"]
                            }
                        ],
                        "rubric": {
                            "technical_excellence": 35,
                            "innovation_creativity": 20,
                            "presentation_communication": 20,
                            "documentation": 15,
                            "teamwork_collaboration": 10
                        }
                    },
                    
                    "assessment_strategy": {
                        "formative": [
                            "Weekly knowledge check quizzes",
                            "In-class polling and discussions",
                            "Peer code reviews"
                        ],
                        "summative": {
                            "assignments": 35,
                            "quizzes": 15,
                            "capstone_project": 40,
                            "participation": 10
                        },
                        "grading_scale": {
                            "A": "90-100",
                            "B": "80-89",
                            "C": "70-79",
                            "D": "60-69",
                            "F": "0-59"
                        }
                    },
                    
                    "instructor_notes": {
                        "common_misconceptions": [
                            {
                                "concept": "Foundational concepts",
                                "misconception": "Students often confuse terminology in early modules",
                                "correction_strategy": "Use consistent terminology glossary and concept maps"
                            },
                            {
                                "concept": "Implementation details",
                                "misconception": "Students may skip testing and validation steps",
                                "correction_strategy": "Emphasize test-driven development from week 1"
                            }
                        ],
                        "engagement_hooks": [
                            "Use recent industry news to introduce each week's topic",
                            "Start with surprising statistics or case studies",
                            "Connect concepts to student career aspirations"
                        ],
                        "differentiation": {
                            "struggling_learners": [
                                "Provide step-by-step lab guides",
                                "Pair with peer mentors",
                                "Offer additional office hours"
                            ],
                            "advanced_learners": [
                                "Assign extension challenges",
                                "Independent research project option",
                                "Peer teaching opportunities"
                            ]
                        },
                        "technology_requirements": [
                            "Laptop with 8GB+ RAM",
                            "Software: VS Code, Git, Docker",
                            "Stable internet connection"
                        ]
                    },
                    
                    "auto_generated_quizzes": [
                        {
                            "module": week,
                            "questions": [
                                {
                                    "question": f"What is a key principle of the {phases[(week-1) % 4].capitalize()} phase?",
                                    "options": [
                                        "Understanding context and stakeholder needs",
                                        "Random implementation without planning",
                                        "Skipping documentation",
                                        "Ignoring user feedback"
                                    ],
                                    "correct": 0,
                                    "explanation": f"The {phases[(week-1) % 4].capitalize()} phase focuses on systematic approaches aligned with IRIS methodology.",
                                    "blooms_level": "understand"
                                },
                                {
                                    "question": f"Which Bloom's taxonomy level is most associated with the {phases[(week-1) % 4].capitalize()} phase?",
                                    "options": [
                                        blooms_by_phase[phases[(week-1) % 4]][0][0].capitalize(),
                                        "Memorization only",
                                        "Random guessing",
                                        "None of the above"
                                    ],
                                    "correct": 0,
                                    "explanation": f"The {phases[(week-1) % 4].capitalize()} phase emphasizes {blooms_by_phase[phases[(week-1) % 4]][0][0]} level cognitive skills.",
                                    "blooms_level": "remember"
                                }
                            ]
                        }
                        for week in range(1, duration_weeks + 1)
                    ],
                    
                    "_metadata": {
                        "framework": "IRIS + CDIO",
                        "standard": "MIT OCW Enterprise",
                        "blooms_taxonomy": True,
                        "sfia_aligned": True,
                        "version": "2.0",
                        "generated_by": "Alexandria AI (Mock Mode)"
                    }
                }
            }

        if provider not in clients:
            continue

        try:
            content = ""
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                content = response.text
            else:
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[{"role": "system", "content": "You are a JSON instructional designer."}, {"role": "user", "content": prompt}],
                    temperature=0.7
                )
                content = response.choices[0].message.content

            try:
                content = content.replace("```json", "").replace("```", "").strip()
                return {"success": True, "agenda": json.loads(content), "provider": provider}
            except:
                import re
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    return {"success": True, "agenda": json.loads(match.group()), "provider": provider}
                raise ValueError("Failed to parse JSON")

        except Exception as e:
            print(f"[{provider}] Teaching Agenda Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {
        "success": False,
        "error": f"All providers failed: {errors}",
        "agenda": None
    }


async def extract_knowledge_points(content: str) -> Dict:
    """
    Extract knowledge points for knowledge graph visualization.
    Returns nodes and edges for graph display.
    """
    
    prompt = f"""Analyze this educational content and extract a knowledge graph.

CONTENT:
{content[:8000]}

OUTPUT JSON:
{{
    "nodes": [
        {{"id": "node1", "label": "Concept Name", "type": "concept|skill|topic", "importance": "high|medium|low"}}
    ],
    "edges": [
        {{"source": "node1", "target": "node2", "relationship": "requires|enables|relates_to"}}
    ],
    "clusters": [
        {{"name": "Cluster Name", "nodes": ["node1", "node2"], "description": "What this cluster represents"}}
    ]
}}

Extract 10-20 key concepts with meaningful relationships. Respond ONLY with valid JSON."""

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(1)
            return {
                "success": True,
                "provider": "mock",
                "graph": {
                    "nodes": [
                        {"id": "n1", "label": "Core Concept", "type": "concept", "importance": "high"},
                        {"id": "n2", "label": "Foundation", "type": "topic", "importance": "high"},
                        {"id": "n3", "label": "Application", "type": "skill", "importance": "medium"},
                        {"id": "n4", "label": "Advanced Topic", "type": "topic", "importance": "medium"},
                        {"id": "n5", "label": "Practice", "type": "skill", "importance": "low"}
                    ],
                    "edges": [
                        {"source": "n2", "target": "n1", "relationship": "enables"},
                        {"source": "n1", "target": "n3", "relationship": "enables"},
                        {"source": "n1", "target": "n4", "relationship": "relates_to"},
                        {"source": "n3", "target": "n5", "relationship": "requires"}
                    ],
                    "clusters": [
                        {"name": "Fundamentals", "nodes": ["n1", "n2"], "description": "Core building blocks"}
                    ]
                }
            }

        if provider not in clients:
            continue

        try:
            response_content = ""
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                response_content = response.text
            else:
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.5
                )
                response_content = response.choices[0].message.content

            try:
                response_content = response_content.replace("```json", "").replace("```", "").strip()
                return {"success": True, "graph": json.loads(response_content), "provider": provider}
            except:
                import re
                match = re.search(r'\{.*\}', response_content, re.DOTALL)
                if match:
                    return {"success": True, "graph": json.loads(match.group()), "provider": provider}
                raise ValueError("Failed to parse JSON")

        except Exception as e:
            print(f"[{provider}] Knowledge Graph Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {"success": False, "error": f"All providers failed: {errors}", "graph": None}


async def refine_teaching_agenda(
    parsed_content: Dict,
    current_agenda: Dict,
    refinement_prompt: str,
    target_audience: str = "Intermediate",
    duration_weeks: int = 4,
) -> Dict:
    """
    Refine an existing Alexandria teaching agenda using a natural-language instruction.
    Keeps the same overall JSON schema while updating the curriculum plan.
    """

    parsed_json = json.dumps(parsed_content, indent=2)
    agenda_json = json.dumps(current_agenda, indent=2)

    prompt = f"""You are Alexandria AI, an expert curriculum designer.

You are given:
1. The original parsed course intent and source-derived structure
2. The current generated teaching agenda
3. A refinement instruction from the instructor

Your task is to revise the agenda while preserving a valid, production-ready JSON structure.

ORIGINAL PARSED CONTENT:
{parsed_json}

CURRENT AGENDA:
{agenda_json}

TARGET AUDIENCE: {target_audience}
DURATION WEEKS: {duration_weeks}

REFINEMENT INSTRUCTION:
{refinement_prompt}

REQUIREMENTS:
- Return ONLY valid JSON.
- Preserve the top-level agenda structure.
- Keep the course aligned to IRIS/CDIO framing.
- Preserve strong learning progression across modules.
- Update titles, subtitles, objectives, activities, quizzes, assignments, and capstone details where needed.
- Make the refinement concrete and specific, not superficial.
- Keep the number of modules equal to {duration_weeks} unless the instruction explicitly requires a different structure.
- Ensure every module still has meaningful learning actions.

Respond ONLY with the full refined agenda JSON."""

    errors = []
    for provider in PRIORITY:
        if provider == "mock":
            await asyncio.sleep(1)
            refined = json.loads(json.dumps(current_agenda))
            refined["tagline"] = f"{refined.get('tagline', '')} Refined: {refinement_prompt[:80]}".strip()
            for idx, module in enumerate(refined.get("modules", [])):
                if idx == 0:
                    subtitle = module.get("subtitle") or module.get("title") or ""
                    module["subtitle"] = f"{subtitle} — refined for {target_audience.lower()} learners".strip(" —")
            return {"success": True, "provider": "mock", "agenda": refined}

        if provider not in clients:
            continue

        try:
            response_content = ""
            if provider == "gemini":
                response = await clients[provider].generate_content_async(prompt)
                response_content = response.text
            else:
                response = await clients[provider].chat.completions.create(
                    model=MODELS[provider],
                    messages=[
                        {"role": "system", "content": "You are Alexandria AI, an expert curriculum designer. Return only valid JSON."},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.6,
                    max_tokens=8000,
                )
                response_content = response.choices[0].message.content

            try:
                response_content = response_content.replace("```json", "").replace("```", "").strip()
                return {"success": True, "agenda": json.loads(response_content), "provider": provider}
            except Exception:
                import re
                match = re.search(r'\{.*\}', response_content, re.DOTALL)
                if match:
                    return {"success": True, "agenda": json.loads(match.group()), "provider": provider}
                raise ValueError("Failed to parse JSON")

        except Exception as e:
            print(f"[{provider}] Agenda Refine Error: {e}")
            errors.append(f"{provider}: {str(e)}")
            continue

    return {"success": False, "error": f"All providers failed: {errors}", "agenda": current_agenda}
